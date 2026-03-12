"""
Generate HydroClaw 80-page PPT using python-pptx
Based on Presenton-style layout with professional design
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme - professional blue/tech theme
PRIMARY = RGBColor(0x00, 0x6D, 0xAA)  # Deep blue
SECONDARY = RGBColor(0x00, 0x91, 0xD5)  # Bright blue
ACCENT = RGBColor(0x00, 0xB4, 0xD8)  # Light blue
DARK = RGBColor(0x1A, 0x1A, 0x2E)  # Dark navy
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)  # Light gray-blue
TEXT_DARK = RGBColor(0x2D, 0x3A, 0x4A)  # Dark text
TEXT_GRAY = RGBColor(0x6B, 0x7B, 0x8D)  # Gray text
HIGHLIGHT = RGBColor(0x00, 0xC9, 0xA7)  # Teal accent
ORANGE = RGBColor(0xFF, 0x6B, 0x35)  # Orange accent

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    """Set slide background color"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, opacity=None):
    """Add a colored rectangle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    """Add a text box with specified properties"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_DARK, spacing=Pt(6)):
    """Add a bulleted list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = spacing
        p.level = 0
    return txBox


def title_slide(title, subtitle=""):
    """Create a title/cover slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK)
    # Left accent bar
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)
    # Title
    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(2),
                 title, font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(1.5),
                     subtitle, font_size=22, color=ACCENT, alignment=PP_ALIGN.LEFT)
    # Bottom line
    add_shape(slide, Inches(1.5), Inches(6.5), Inches(3), Inches(0.04), ACCENT)
    return slide


def chapter_slide(chapter_num, chapter_title, subtitle=""):
    """Create a chapter divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
                 f"第{chapter_num}章", font_size=24, color=RGBColor(0x80, 0xD0, 0xFF), bold=False)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(2),
                 chapter_title, font_size=42, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(1),
                     subtitle, font_size=20, color=RGBColor(0xB0, 0xE0, 0xFF))
    add_shape(slide, Inches(1.5), Inches(6.2), Inches(2), Inches(0.04), WHITE)
    return slide


def content_slide(title, bullets, note=""):
    """Standard content slide with title and bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    # Top bar
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 title, font_size=30, color=DARK, bold=True)
    # Separator line
    add_shape(slide, Inches(0.8), Inches(1.25), Inches(2), Inches(0.03), ACCENT)
    # Bullets
    add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
                    bullets, font_size=18, color=TEXT_DARK)
    if note:
        add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
                     note, font_size=12, color=TEXT_GRAY)
    return slide


def two_column_slide(title, left_title, left_items, right_title, right_items):
    """Two-column layout slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 title, font_size=30, color=DARK, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.25), Inches(2), Inches(0.03), ACCENT)
    # Left column
    add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(5.2), LIGHT_BG)
    add_text_box(slide, Inches(0.9), Inches(1.8), Inches(5), Inches(0.6),
                 left_title, font_size=20, color=PRIMARY, bold=True)
    add_bullet_list(slide, Inches(0.9), Inches(2.5), Inches(5), Inches(4),
                    left_items, font_size=16, color=TEXT_DARK)
    # Right column
    add_shape(slide, Inches(6.6), Inches(1.6), Inches(5.5), Inches(5.2), LIGHT_BG)
    add_text_box(slide, Inches(6.9), Inches(1.8), Inches(5), Inches(0.6),
                 right_title, font_size=20, color=PRIMARY, bold=True)
    add_bullet_list(slide, Inches(6.9), Inches(2.5), Inches(5), Inches(4),
                    right_items, font_size=16, color=TEXT_DARK)
    return slide


def highlight_slide(title, main_text, sub_text=""):
    """Highlight/emphasis slide with dark background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
                 title, font_size=22, color=ACCENT, bold=False)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(3),
                 main_text, font_size=34, color=WHITE, bold=True)
    if sub_text:
        add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(1.5),
                     sub_text, font_size=18, color=RGBColor(0xA0, 0xB0, 0xC0))
    return slide


def stats_slide(title, stats):
    """Stats/metrics slide with key numbers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 title, font_size=30, color=DARK, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.25), Inches(2), Inches(0.03), ACCENT)

    col_width = 10.5 / len(stats)
    for i, (num, label) in enumerate(stats):
        x = Inches(1.0 + i * col_width)
        add_shape(slide, x, Inches(2.0), Inches(col_width - 0.3), Inches(4.0), LIGHT_BG)
        add_text_box(slide, x, Inches(2.5), Inches(col_width - 0.3), Inches(1.5),
                     num, font_size=44, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(4.2), Inches(col_width - 0.3), Inches(1.5),
                     label, font_size=16, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
    return slide


def numbered_list_slide(title, items):
    """Numbered list slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 title, font_size=30, color=DARK, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.25), Inches(2), Inches(0.03), ACCENT)

    y = Inches(1.6)
    for i, item in enumerate(items):
        # Number circle
        add_text_box(slide, Inches(0.8), y, Inches(0.6), Inches(0.4),
                     str(i+1), font_size=16, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        # Item text
        add_text_box(slide, Inches(1.5), y, Inches(10), Inches(0.4),
                     item, font_size=17, color=TEXT_DARK)
        y += Inches(0.55)
    return slide


# ========================================
# SLIDE CONTENT - 80 PAGES
# ========================================

print("Generating 80-page HydroClaw PPT...")

# Slide 1: Cover
title_slide(
    "HydroClaw",
    "水网认知智能体系发布会\n中国水利水电科学研究院 · 2026年4月"
)

# Slide 2: Agenda/TOC
content_slide("目 录", [
    "第一章  愿景与定位 — 水利行业的AI革命",
    "第二章  三层产品架构 — HydroOS + HydroMAS + HydroTouch",
    "第三章  五层技术纵深 — L0-L4全栈架构",
    "第四章  应用场景 — 十大核心场景",
    "第五章  技术优势 — 开源、易用、专业",
    "第六章  发展路线图 — 从MVP到行业引领",
    "第七章  团队与合作 — 开放共赢的生态"
])

# ===== CHAPTER 1: Vision & Positioning =====
# Slide 3
chapter_slide("一", "愿景与定位", "水利行业的AI革命")

# Slide 4
highlight_slide("行业背景",
    "中国水网管理面临前所未有的挑战",
    "15万+水利工程 · 跨流域调度复杂性 · 极端气候频发")

# Slide 5
content_slide("水利行业三大痛点", [
    "🔸 数据孤岛 — 各水利工程系统独立，数据无法互通",
    "🔸 决策滞后 — 传统模型计算耗时长，无法实时响应",
    "🔸 人才短缺 — 高水平水利工程师供不应求",
    "",
    "传统水利信息化建设已进入瓶颈期",
    "亟需AI技术实现范式转变"
])

# Slide 6
content_slide("AI大模型时代的历史性机遇", [
    "大语言模型（LLM）的突破性进展",
    "多智能体系统（MAS）的成熟应用",
    "AI Agent 技术的实际落地",
    "GPU 算力的指数级增长",
    "水利行业数字化转型的政策驱动",
    "「数字孪生流域」国家战略的推进"
])

# Slide 7
highlight_slide("HydroClaw 是什么",
    "全球首个面向水网运行管理的\n认知智能体系",
    "Cognitive Intelligence System for Water Network Operations")

# Slide 8
content_slide("HydroClaw 核心定位", [
    "不是简单的「AI+水利」叠加",
    "而是深度融合的认知决策平台",
    "",
    "核心理念：",
    "让每一个水利从业者都拥有一个「AI首席工程师」",
    "",
    "从经验驱动 → 智能认知驱动"
])

# Slide 9
content_slide("使命与愿景", [
    "使命：用AI重新定义水网运行管理方式",
    "愿景：成为全球水利行业的「大脑」——万人万用的水网认知操作系统",
    "目标：5年内覆盖中国50%以上重要水利工程的智能决策支持",
    "",
    "让水利从「靠经验」到「靠智能」",
    "让决策从「事后处理」到「事前预判」"
])

# Slide 10
numbered_list_slide("产品六大版本", [
    "HydroClaw Lite — 个人助理版（面向个人水利工程师）",
    "HydroClaw Pro — 专业版（面向水利设计/咨询企业）",
    "HydroClaw Enterprise — 企业版（面向水利管理局/大型水务集团）",
    "HydroClaw Cloud — 云服务版（SaaS模式，按需付费）",
    "HydroClaw Edge — 边缘版（面向物联网/现场设备）",
    "HydroClaw Open — 开源社区版（推动行业生态）"
])

# Slide 11
two_column_slide("Lite vs Pro 对比",
    "Lite — 个人助理版", [
        "面向个人水利工程师",
        "AI助手问答",
        "基础数据查询",
        "简单报告生成",
        "免费使用"
    ],
    "Pro — 专业版", [
        "面向水利设计/咨询企业",
        "多Agent协同",
        "专业模型运行",
        "高级数据分析",
        "团队协作功能"
    ])

# Slide 12
two_column_slide("Enterprise vs Cloud 对比",
    "Enterprise — 企业版", [
        "面向水利管理局/大型水务集团",
        "私有化部署",
        "定制化开发",
        "全功能覆盖",
        "专属技术支持"
    ],
    "Cloud — 云服务版", [
        "SaaS模式，按需付费",
        "即开即用",
        "弹性扩缩",
        "自动运维",
        "低成本入门"
    ])

# ===== CHAPTER 2: Architecture =====
# Slide 13
chapter_slide("二", "三层产品架构", "HydroOS + HydroMAS + HydroTouch")

# Slide 14
highlight_slide("架构总览",
    "如果HydroClaw是一个人\nHydroOS是躯体，HydroMAS是大脑\nHydroTouch是五官和四肢",
    "三层架构是HydroClaw的核心创新")

# Slide 15
content_slide("三层架构概览", [
    "🔹 HydroOS — 计算底座层",
    "    提供数据、算力、工具等基础设施",
    "",
    "🔹 HydroMAS — 多智能体中枢层",
    "    15个Agent + 17个Skill + 认知决策引擎",
    "",
    "🔹 HydroTouch — 多端接入层",
    "    Web、飞书、桌面、移动、AR、API全覆盖"
])

# Slide 16
chapter_slide("二·一", "HydroOS — 计算底座层", "21个MCP Server · Ray分布式计算 · 知识图谱")

# Slide 17
content_slide("HydroOS 核心能力", [
    "21个MCP Server — 提供水文数据、GIS地理、气象预报、水质监测等专业工具",
    "Ray分布式计算 — 支持千核级并行水文模拟",
    "知识图谱 — 百万级水利实体和关系，构建行业知识库",
    "向量检索 — 基于水利语料的专业RAG系统",
    "数据湖 — 统一管理多源异构水利数据"
])

# Slide 18
numbered_list_slide("8大水文MCP", [
    "hydro_data_mcp — 水文数据采集与管理",
    "hydro_model_mcp — 水文模型调度",
    "hydro_gis_mcp — GIS空间分析",
    "hydro_forecast_mcp — 水文预报",
    "hydro_quality_mcp — 水质监测分析",
    "hydro_iot_mcp — 物联网设备管理",
    "hydro_knowledge_mcp — 知识图谱查询",
    "hydro_optimize_mcp — 优化调度"
])

# Slide 19
content_slide("其他MCP Server", [
    "hydro_report_mcp — 报告生成服务",
    "hydro_alert_mcp — 预警通知服务",
    "hydro_auth_mcp — 认证授权服务",
    "hydro_storage_mcp — 对象存储服务",
    "hydro_scheduler_mcp — 任务调度服务",
    "hydro_monitor_mcp — 系统监控服务",
    "hydro_log_mcp — 日志管理服务",
    "更多MCP Server持续扩展中..."
])

# Slide 20
content_slide("MCP协议 — Model Context Protocol", [
    "MCP是Anthropic提出的AI工具标准化协议",
    "为AI Agent提供统一的工具调用接口",
    "",
    "HydroClaw的MCP Server特性：",
    "• 标准化接口，可被任何MCP兼容客户端调用",
    "• 支持流式响应和异步处理",
    "• 内置权限管理和访问控制",
    "• 自动生成API文档"
])

# Slide 21
two_column_slide("HydroOS 技术栈",
    "计算与存储", [
        "计算框架：Ray + CUDA",
        "关系数据库：PostgreSQL + PostGIS",
        "图数据库：Neo4j",
        "缓存：Redis",
        "对象存储：MinIO"
    ],
    "基础设施", [
        "消息队列：RabbitMQ + Kafka",
        "容器化：Docker + Kubernetes",
        "监控：Prometheus + Grafana",
        "日志：ELK Stack",
        "CI/CD：GitHub Actions"
    ])

# Slide 22
content_slide("Ray分布式计算引擎", [
    "支持千核级并行计算",
    "自动任务调度和负载均衡",
    "水文模型并行化加速10-100倍",
    "支持GPU加速（CUDA）",
    "弹性伸缩，按需分配计算资源",
    "支持分布式训练和推理"
])

# Slide 23
content_slide("知识图谱与向量检索", [
    "知识图谱：",
    "• 百万级水利实体（水库、河流、闸门、泵站等）",
    "• 千万级关系（上下游、供排水、调度逻辑等）",
    "• 支持推理和知识发现",
    "",
    "向量检索（RAG）：",
    "• 水利法规、标准、规范全文索引",
    "• 历史案例和经验知识库",
    "• 语义搜索和智能问答"
])

# Slide 24
chapter_slide("二·二", "HydroMAS — 多智能体中枢层", "15个Agent · 17个Skill · 认知决策引擎")

# Slide 25
content_slide("HydroMAS 核心能力", [
    "15个专业Agent — 覆盖水利全生命周期",
    "17个Skill — 可组合的专业技能",
    "IntentRouter — 智能意图路由，自动分配最优Agent",
    "AgentCoordinator — 多Agent协同框架",
    "认知决策引擎 — 四阶认知循环"
])

# Slide 26
content_slide("15大Agent（上）", [
    "1. FloodGuardAgent — 防洪预警智能体",
    "2. DroughtWatchAgent — 旱情监测智能体",
    "3. WaterAllocAgent — 水资源调配智能体",
    "4. QualityPatrolAgent — 水质巡检智能体",
    "5. SchedulerAgent — 工程调度智能体",
    "6. InspectionAgent — 巡检维护智能体",
    "7. ReportAgent — 智能报告生成智能体",
    "8. AnalyticsAgent — 数据分析智能体"
])

# Slide 27
content_slide("15大Agent（下）", [
    "9. KnowledgeAgent — 知识问答智能体",
    "10. SimulationAgent — 水文模拟智能体",
    "11. AlertAgent — 预警管理智能体",
    "12. PlanningAgent — 规划辅助智能体",
    "13. ComplianceAgent — 合规审核智能体",
    "14. TrainingAgent — 培训指导智能体",
    "15. CoordinatorAgent — 协调总控智能体"
])

# Slide 28
content_slide("FloodGuardAgent — 防洪预警", [
    "实时监测降雨、水位、流量数据",
    "基于AI模型预测洪水演进过程",
    "自动生成洪水预警和调度建议",
    "支持多方案比选和风险评估",
    "与SchedulerAgent协同完成联合调度",
    "预报提前时间达6小时以上"
])

# Slide 29
content_slide("WaterAllocAgent — 水资源调配", [
    "跨流域水资源优化调配",
    "多目标优化（供水、灌溉、生态、发电）",
    "自动生成调度方案和分析报告",
    "实时跟踪水量平衡和用水效率",
    "支持「What-if」情景模拟",
    "水资源利用率提升15%+"
])

# Slide 30
content_slide("QualityPatrolAgent — 水质巡检", [
    "实时监测水质参数（pH、溶解氧、COD等）",
    "基于异常检测引擎自动识别水质异常",
    "利用知识图谱进行污染溯源",
    "自动生成水质巡检报告",
    "与AlertAgent联动进行预警",
    "污染事件发现时间从24h缩短到30min"
])

# Slide 31
content_slide("IntentRouter — 智能意图路由", [
    "基于大语言模型的意图理解",
    "自动将用户需求分配给最合适的Agent",
    "支持多轮对话和上下文管理",
    "复杂任务自动拆解和多Agent协同",
    "",
    "用户只需用自然语言描述需求",
    "HydroClaw自动完成任务编排"
])

# Slide 32
content_slide("17大Skill", [
    "数据可视化 · 报告生成 · 模型运行 · 预报分析",
    "GIS制图 · 知识检索 · 文档摘要 · 代码生成",
    "数据清洗 · 异常检测 · 趋势预测 · 对比分析",
    "方案评估 · 合规检查 · 多语言翻译 · 语音交互 · 实时监控",
    "",
    "Skill是可组合的原子能力单元",
    "Agent通过组合不同Skill完成复杂任务"
])

# Slide 33
content_slide("认知决策引擎", [
    "四阶认知循环：",
    "  感知 → 理解 → 决策 → 行动",
    "",
    "核心技术：",
    "• 多模态融合：文本 + 图像 + 时序 + 空间数据",
    "• 上下文管理：长期记忆 + 工作记忆 + 情景记忆",
    "• 不确定性推理：贝叶斯网络 + 模糊逻辑",
    "• 可解释AI：决策过程透明可追溯"
])

# Slide 34
chapter_slide("二·三", "HydroTouch — 多端接入层", "Web · 飞书 · 桌面 · 移动 · AR · API")

# Slide 35
numbered_list_slide("六大接入方式", [
    "Web端 — React + Ant Design 的专业工作台",
    "飞书集成 — 6大场景深度对接企业协作",
    "Tauri桌面端 — 轻量级跨平台原生应用",
    "Flutter移动端 — 现场巡检与应急响应",
    "AR眼镜 — 增强现实辅助巡检维护",
    "API / MCP Server — 开放接口供第三方系统调用"
])

# Slide 36
content_slide("Web端 — 专业工作台", [
    "技术栈：React + TypeScript + Ant Design",
    "实时数据看板和GIS地图",
    "多Agent对话式交互界面",
    "报告编辑和审批流程",
    "团队协作和权限管理",
    "响应式设计，支持大屏和移动端"
])

# Slide 37
numbered_list_slide("飞书集成6大场景", [
    "日常值班 — AI自动生成值班报告",
    "应急响应 — 洪水预警自动推送与处置建议",
    "数据查询 — 自然语言查询水文数据",
    "报告审批 — AI生成 + 人工审核 + 一键签发",
    "会议纪要 — AI总结会议要点与待办事项",
    "知识共享 — 团队知识库智能问答"
])

# Slide 38
two_column_slide("桌面端与移动端",
    "Tauri桌面端", [
        "轻量级跨平台（Win/Mac/Linux）",
        "原生性能体验",
        "离线模式支持",
        "本地数据缓存",
        "安装包 < 30MB"
    ],
    "Flutter移动端", [
        "iOS/Android双平台",
        "现场巡检数据采集",
        "应急响应快速处理",
        "GPS定位和拍照上传",
        "离线巡检功能"
    ])

# Slide 39
content_slide("AR眼镜 — 增强现实巡检", [
    "支持主流AR眼镜设备",
    "实时叠加水利工程信息",
    "语音交互，解放双手操作",
    "缺陷识别和自动记录",
    "远程专家协助和指导",
    "巡检效率提升300%"
])

# ===== CHAPTER 3: Technical Depth =====
# Slide 40
chapter_slide("三", "五层技术纵深", "L0-L4 全栈架构")

# Slide 41
content_slide("五层内部架构 L0-L4", [
    "L0 Data Layer — 数据湖与存储引擎",
    "L1 Compute Layer — 分布式计算与模型服务",
    "L2 MCP Tools Layer — 21个标准化工具服务",
    "L3 Skills Layer — 17个可组合专业技能",
    "L4 Agents Layer — 15个自主决策智能体",
    "",
    "每一层都可独立扩展和升级"
])

# Slide 42
content_slide("L0 Data Layer — 数据湖", [
    "多源异构数据统一管理",
    "支持结构化、半结构化、非结构化数据",
    "时序数据高效存储和查询",
    "空间数据（GIS）原生支持",
    "数据血缘和质量管理",
    "PB级数据存储能力"
])

# Slide 43
content_slide("L1 Compute Layer — 分布式计算", [
    "Ray集群：弹性伸缩，自动负载均衡",
    "GPU加速：CUDA/cuDNN深度学习推理",
    "模型服务：统一的模型训练和推理平台",
    "流计算：Kafka Streams实时数据处理",
    "批计算：Spark大规模历史数据分析"
])

# Slide 44
content_slide("L2 MCP Tools Layer — 工具服务", [
    "21个标准化MCP Server",
    "每个Server提供独立的专业能力",
    "统一的认证和鉴权机制",
    "自动化API文档和SDK生成",
    "支持流式响应和长时间任务",
    "可热插拔，独立部署和升级"
])

# Slide 45
content_slide("L3 Skills Layer — 专业技能", [
    "17个可组合的原子技能",
    "技能编排引擎（Skill Orchestrator）",
    "声明式技能定义和注册",
    "技能版本管理和灰度发布",
    "性能监控和自动优化"
])

# Slide 46
content_slide("L4 Agents Layer — 智能体", [
    "15个自主决策智能体",
    "Agent间通信协议（A2A Protocol）",
    "任务分解和多Agent协同",
    "自主学习和经验积累",
    "人机协同决策机制",
    "Agent市场（可扩展第三方Agent）"
])

# Slide 47
content_slide("14个核心算法模块（上）", [
    "1. 水文时间序列预测 — LSTM / Transformer",
    "2. 洪水演进模拟 — Saint-Venant方程数值解",
    "3. 水资源优化调度 — 多目标遗传算法",
    "4. 水质扩散模型 — WASP / QUAL2K",
    "5. 降雨径流模型 — 新安江 / SCS-CN",
    "6. 管网水力计算 — EPANET集成",
    "7. 地下水流动模拟 — MODFLOW接口"
])

# Slide 48
content_slide("14个核心算法模块（下）", [
    "8. 水库群联合调度 — 动态规划 + 强化学习",
    "9. 河道冲淤演变 — 二维水沙模型",
    "10. 生态流量计算 — Tennant / IHA方法",
    "11. 旱情预警模型 — SPI / PDSI指数",
    "12. 异常检测引擎 — Isolation Forest + AutoEncoder",
    "13. 知识图谱推理 — TransE / RotatE嵌入",
    "14. 多模态融合模型 — 跨模态注意力机制"
])

# Slide 49
content_slide("水文时间序列预测", [
    "基于Transformer架构的预测模型",
    "支持多变量联合预测",
    "自适应时间步长（分钟/小时/日）",
    "不确定性量化（概率预测）",
    "在线学习，模型持续更新",
    "预测精度提升20%以上（vs传统方法）"
])

# Slide 50
content_slide("水库群联合调度算法", [
    "动态规划（DP）：求解最优调度策略",
    "强化学习（RL）：自适应复杂环境",
    "多目标优化：防洪 + 供水 + 发电 + 生态",
    "约束处理：库容、流量、水位安全约束",
    "实时滚动优化：每小时更新调度方案",
    "支持百座级水库群联合优化"
])

# Slide 51
stats_slide("测试与质量保障", [
    ("1851", "自动化测试"),
    (">80%", "代码覆盖率"),
    ("3层", "测试体系\n(单元/集成/E2E)"),
    ("24/7", "CI/CD流水线")
])

# ===== CHAPTER 4: Application Scenarios =====
# Slide 52
chapter_slide("四", "应用场景", "十大核心场景")

# Slide 53
content_slide("场景1：防洪调度", [
    "痛点：极端降雨频发，传统预报滞后",
    "",
    "方案：FloodGuardAgent + SimulationAgent + SchedulerAgent 协同",
    "",
    "效果：",
    "• 预报提前6小时",
    "• 调度决策时间从4小时缩短到30分钟",
    "• 减少洪灾损失40%以上"
])

# Slide 54
content_slide("防洪调度 — 技术细节", [
    "实时雨量站数据接入（5分钟间隔）",
    "雷达降雨估测和融合",
    "分布式水文模型实时运行",
    "洪水演进可视化展示",
    "多方案自动比选",
    "调度指令自动生成和下发"
])

# Slide 55
content_slide("场景2：水资源调配", [
    "痛点：跨流域调水的多目标冲突",
    "",
    "方案：WaterAllocAgent + 多目标优化算法",
    "",
    "效果：",
    "• 水资源利用率提升15%",
    "• 调度方案生成从3天缩短到2小时",
    "• 节水经济效益数千万元/年"
])

# Slide 56
content_slide("场景3：水质监测预警", [
    "痛点：污染事件发现滞后、溯源困难",
    "",
    "方案：QualityPatrolAgent + 异常检测引擎 + 知识图谱",
    "",
    "效果：",
    "• 污染事件发现时间从24小时缩短到30分钟",
    "• 溯源准确率达90%以上",
    "• 自动生成应急处置方案"
])

# Slide 57
content_slide("场景4：智能巡检", [
    "痛点：水利工程巡检人力成本高、覆盖不全",
    "",
    "方案：InspectionAgent + AR眼镜 + IoT传感器",
    "",
    "效果：",
    "• 巡检效率提升300%",
    "• 隐患发现率提升50%",
    "• 巡检报告自动生成"
])

# Slide 58
content_slide("场景5：应急响应", [
    "痛点：突发事件响应链条长、信息不畅",
    "",
    "方案：AlertAgent + CoordinatorAgent + 飞书集成",
    "",
    "效果：",
    "• 应急响应时间从2小时缩短到15分钟",
    "• 处置建议自动推送",
    "• 全流程可追溯"
])

# Slide 59
content_slide("场景6：知识管理", [
    "痛点：水利行业知识碎片化、传承困难",
    "",
    "方案：KnowledgeAgent + 知识图谱 + RAG",
    "",
    "效果：",
    "• 新人上手时间从6个月缩短到1个月",
    "• 知识检索效率提升10倍",
    "• 构建行业最大知识库"
])

# Slide 60
content_slide("场景7：智能报告", [
    "痛点：报告编写耗时、格式不统一",
    "",
    "方案：ReportAgent + 模板引擎 + 数据可视化",
    "",
    "效果：",
    "• 报告生成时间从1天缩短到10分钟",
    "• 格式统一规范",
    "• 数据自动填充和校验"
])

# Slide 61
content_slide("场景8：规划辅助", [
    "痛点：水利规划方案比选复杂、周期长",
    "",
    "方案：PlanningAgent + SimulationAgent + 多方案对比",
    "",
    "效果：",
    "• 规划方案生成效率提升500%",
    "• 多维度自动评估",
    "• 可视化对比展示"
])

# Slide 62
content_slide("场景9：合规审核", [
    "痛点：水利法规繁多、人工审核易遗漏",
    "",
    "方案：ComplianceAgent + 法规知识库",
    "",
    "效果：",
    "• 合规审核时间从1周缩短到1天",
    "• 遗漏率降低90%",
    "• 自动生成合规报告"
])

# Slide 63
content_slide("场景10：培训教育", [
    "痛点：水利人才培养周期长、实操机会少",
    "",
    "方案：TrainingAgent + 虚拟仿真 + AI导师",
    "",
    "效果：",
    "• 培训效率提升200%",
    "• 虚拟实操零风险",
    "• 个性化学习路径"
])

# Slide 64
stats_slide("十大场景效果总览", [
    ("6h+", "洪水预报\n提前时间"),
    ("15%", "水资源\n利用率提升"),
    ("300%", "巡检效率\n提升"),
    ("90%", "遗漏率\n降低")
])

# ===== CHAPTER 5: Advantages =====
# Slide 65
chapter_slide("五", "技术优势", "开源 · 易用 · 专业")

# Slide 66
highlight_slide("全栈开源",
    "代码 · 文档 · 模型权重\n全部开放",
    "许可证：MIT License  |  GitHub Stars 1000+")

# Slide 67
content_slide("开源优势", [
    "透明可信 — 代码完全公开，可审计可验证",
    "社区驱动 — 汇聚全球水利AI开发者",
    "快速迭代 — 社区贡献加速产品进化",
    "生态丰富 — 第三方插件和扩展",
    "无厂商锁定 — 自主可控，数据安全",
    "学术合作 — 支持科研和教育"
])

# Slide 68
content_slide("与竞品对比", [
    "vs 传统水利系统：",
    "  AI能力 ★★★★★ vs ★★，多Agent协同能力领先",
    "",
    "vs 通用AI平台：",
    "  水利专业性 ★★★★★ vs ★★，多端接入更全面",
    "",
    "vs 国际水利软件：",
    "  自主可控 ★★★★★，中国水利特色适配"
])

# Slide 69
content_slide("Docker一键部署", [
    "docker compose up -d",
    "",
    "5分钟完成全栈部署",
    "支持CPU / GPU双模式",
    "自动配置数据库、消息队列、存储",
    "包含监控和日志系统",
    "支持单机和集群模式"
])

# Slide 70
content_slide("知识飞轮效应", [
    "数据 → 模型 → 决策 → 反馈 → 更多数据",
    "",
    "越用越智能的自学习系统",
    "行业知识不断积累和优化",
    "决策质量持续提升",
    "用户越多，系统越强",
    "形成正向反馈闭环"
])

# Slide 71
two_column_slide("安全与合规",
    "数据安全", [
        "端到端加密传输",
        "静态数据加密存储",
        "细粒度权限控制",
        "数据脱敏处理",
        "审计日志全记录"
    ],
    "合规认证", [
        "等保三级认证",
        "ISO 27001信息安全",
        "水利行业标准合规",
        "隐私计算支持",
        "国产化适配"
    ])

# ===== CHAPTER 6: Roadmap =====
# Slide 72
chapter_slide("六", "发展路线图", "从MVP到行业引领")

# Slide 73
content_slide("Phase A（2026 Q1-Q2）：MVP", [
    "核心Agent上线：FloodGuard、WaterAlloc、QualityPatrol",
    "基础MCP Tools完成（8个核心MCP）",
    "Web端和飞书端开放",
    "基础知识图谱构建",
    "首批3家试点单位接入"
])

# Slide 74
content_slide("Phase B（2026 Q3-Q4）：增强", [
    "全部15个Agent上线",
    "认知决策引擎优化",
    "Tauri桌面端和Flutter移动端发布",
    "知识图谱扩展到百万级实体",
    "试点扩展到10家单位"
])

# Slide 75
content_slide("Phase C（2027）：生态", [
    "开发者平台上线",
    "第三方插件市场",
    "AR/IoT深度集成",
    "Agent市场开放",
    "社区贡献者达500+"
])

# Slide 76
content_slide("Phase D（2028）：规模", [
    "覆盖全国50+重要水利工程",
    "多语言国际化（英语、法语、西班牙语等）",
    "行业标准制定参与",
    "SaaS版本商业化运营",
    "年度收入突破亿元"
])

# Slide 77
content_slide("Phase E（2029+）：引领", [
    "全球水利AI标准贡献者",
    "数字孪生流域全覆盖",
    "自主进化认知系统",
    "覆盖全球主要水利工程",
    "成为水利行业的「操作系统」"
])

# ===== CHAPTER 7: Team =====
# Slide 78
chapter_slide("七", "团队与合作", "开放共赢的生态")

# Slide 79
content_slide("中国水利水电科学研究院", [
    "国家级水利科研机构",
    "70+年水利科技积累",
    "水利行业最权威的技术支撑单位",
    "",
    "技术团队：",
    "• AI研究组 — 大模型、多Agent系统、知识图谱",
    "• 水利专家组 — 水文、水资源、水环境、水工程",
    "• 工程团队 — 全栈开发、DevOps、质量保障"
])

# Slide 80
content_slide("开放合作", [
    "学术合作 — 高校联合研究",
    "产业合作 — 水务企业联合创新",
    "生态合作 — 开发者社区共建",
    "标准合作 — 行业标准共同制定",
    "",
    "联系方式：hydroclaw@iwhr.com",
    "GitHub: github.com/leixiaohui-1974/HydroClaw"
])

# Slide 81 (bonus closing)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
             "AI赋能水网，认知驱动未来", font_size=44, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(1.5),
             "HydroClaw不仅是一个产品\n更是一种全新的水利工作方式",
             font_size=24, color=ACCENT)
add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1),
             "当AI真正理解水，世界将变得不同。",
             font_size=20, color=RGBColor(0xA0, 0xB0, 0xC0))
add_shape(slide, Inches(1.5), Inches(6.5), Inches(3), Inches(0.04), ACCENT)

# Slide 82 (Thank you)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_text_box(slide, Inches(0), Inches(2.5), W, Inches(2),
             "谢 谢", font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(4.8), W, Inches(1),
             "hydroclaw@iwhr.com  |  github.com/leixiaohui-1974/HydroClaw",
             font_size=18, color=RGBColor(0xB0, 0xE0, 0xFF), alignment=PP_ALIGN.CENTER)

# Save
output_path = os.path.join("发布会26.4", "output", "HydroClaw_Presenton_80p.pptx")
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)

total_slides = len(prs.slides)
print(f"Done! Generated {total_slides} slides")
print(f"Saved to: {output_path}")
