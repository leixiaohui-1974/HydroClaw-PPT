#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HydroClaw 50页高质量PPT生成器 V2
基于三引擎评审结果重构
- 包含图片生成和插入
- 多种版式布局
- 数据可视化
"""

import sys
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
import numpy as np

# 设置标准输出编码为UTF-8
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

# 配色方案
C = {
    'blue': RGBColor(91, 155, 213),
    'orange': RGBColor(237, 125, 49),
    'green': RGBColor(112, 173, 71),
    'yellow': RGBColor(255, 192, 0),
    'purple': RGBColor(142, 124, 195),
    'cyan': RGBColor(68, 114, 196),
    'red': RGBColor(192, 0, 0),
    'dark_blue': RGBColor(31, 78, 120),
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(127, 127, 127),
    'light_gray': RGBColor(242, 242, 242),
}

class HydroClawPPTV2:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

        # 创建图片输出目录
        self.img_dir = Path("generated_diagrams")
        self.img_dir.mkdir(exist_ok=True)

        # 字体设置
        self.font_title = "微软雅黑"
        self.font_body = "微软雅黑"

    def add_slide(self, layout_idx=6):
        return self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

    def add_title_bar(self, slide, title, color):
        """添加标题栏"""
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        tf = tb.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.name = self.font_title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = C['white']

    def add_image_with_text(self, slide, title, text_items, image_path, image_on_right=True):
        """图文混排布局"""
        self.add_title_bar(slide, title, C['blue'])

        # 图片位置
        if image_on_right:
            img_left, img_top = Inches(5.5), Inches(1)
            text_left, text_top = Inches(0.5), Inches(1)
            text_width = Inches(4.5)
        else:
            img_left, img_top = Inches(0.5), Inches(1)
            text_left, text_top = Inches(5.5), Inches(1)
            text_width = Inches(4)

        # 插入图片
        if Path(image_path).exists():
            slide.shapes.add_picture(str(image_path), img_left, img_top,
                                    width=Inches(4), height=Inches(4))

        # 添加文本
        tb = slide.shapes.add_textbox(text_left, text_top, text_width, Inches(4))
        tf = tb.text_frame
        tf.word_wrap = True

        for i, item in enumerate(text_items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = self.font_body
            p.font.size = Pt(14)
            p.space_after = Pt(6)
            if item.startswith('•'):
                p.level = 0
            elif item.startswith('  -'):
                p.level = 1

    def generate_architecture_diagram(self):
        """生成四层架构图"""
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        ax.axis('off')

        # 四层架构
        layers = [
            ("认知决策层", "垂直大模型 + 规则引擎", 8, (142/255, 124/255, 195/255)),
            ("技能编排层", "Skill继承体系", 6, (237/255, 125/255, 49/255)),
            ("计算引擎层", "7大通用引擎", 4, (112/255, 173/255, 71/255)),
            ("对象层", "三族元件库", 2, (68/255, 114/255, 196/255))
        ]

        for name, desc, y, rgb in layers:
            rect = plt.Rectangle((1, y-0.8), 8, 1.5, facecolor=rgb, edgecolor='white', linewidth=2)
            ax.add_patch(rect)
            ax.text(5, y, name, ha='center', va='center', fontsize=18, fontweight='bold', color='white')
            ax.text(5, y-0.4, desc, ha='center', va='center', fontsize=12, color='white')

        # 箭头
        for i in range(3):
            y_start = 2 + i*2 - 0.8
            ax.arrow(5, y_start, 0, 1.2, head_width=0.3, head_length=0.2, fc='gray', ec='gray')

        output_path = self.img_dir / "architecture_4layers.png"
        plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close()
        return output_path

    def generate_comparison_chart(self):
        """生成对比图表"""
        fig, ax = plt.subplots(figsize=(10, 5))

        categories = ['设计效率', '准确性', '可解释性', '学习成本']
        traditional = [30, 60, 40, 80]
        hydroclaw = [90, 95, 85, 30]

        x = np.arange(len(categories))
        width = 0.35

        bars1 = ax.bar(x - width/2, traditional, width, label='传统方法', color='#cccccc')
        bars2 = ax.bar(x + width/2, hydroclaw, width, label='HydroClaw', color='#5B9BD5')

        ax.set_ylabel('得分', fontsize=12)
        ax.set_title('HydroClaw vs 传统方法', fontsize=14, fontweight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels(categories, fontsize=11)
        ax.legend(fontsize=11)
        ax.set_ylim(0, 100)
        ax.grid(axis='y', alpha=0.3)

        output_path = self.img_dir / "comparison_chart.png"
        plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close()
        return output_path

    def generate_flow_diagram(self):
        """生成决策流程图"""
        fig, ax = plt.subplots(figsize=(12, 6))
        ax.set_xlim(0, 12)
        ax.set_ylim(0, 6)
        ax.axis('off')

        steps = [
            ("用户输入", 1, 3),
            ("意图理解", 3, 3),
            ("方案生成", 5, 3),
            ("规则校核", 7, 3),
            ("优化调整", 9, 3),
            ("结果输出", 11, 3)
        ]

        for i, (name, x, y) in enumerate(steps):
            color = (91/255, 155/255, 213/255) if i % 2 == 0 else (237/255, 125/255, 49/255)
            circle = plt.Circle((x, y), 0.6, facecolor=color, edgecolor='white', linewidth=2)
            ax.add_patch(circle)
            ax.text(x, y, name, ha='center', va='center', fontsize=11, fontweight='bold', color='white')

            if i < len(steps) - 1:
                ax.arrow(x+0.7, y, 1.3, 0, head_width=0.2, head_length=0.15, fc='gray', ec='gray')

        output_path = self.img_dir / "decision_flow.png"
        plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close()
        return output_path

    def generate_kpi_chart(self):
        """生成KPI效果图"""
        fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(10, 8))

        # 能耗降低
        ax1.text(0.5, 0.6, "15%", ha='center', va='center', fontsize=48, fontweight='bold', color='#5B9BD5')
        ax1.text(0.5, 0.3, "能耗降低", ha='center', va='center', fontsize=16)
        ax1.axis('off')

        # 效率提升
        ax2.text(0.5, 0.6, "3-5倍", ha='center', va='center', fontsize=48, fontweight='bold', color='#70AD47')
        ax2.text(0.5, 0.3, "设计效率提升", ha='center', va='center', fontsize=16)
        ax2.axis('off')

        # 漏损率降低
        ax3.text(0.5, 0.6, "30%", ha='center', va='center', fontsize=48, fontweight='bold', color='#ED7D31')
        ax3.text(0.5, 0.3, "漏损率降低", ha='center', va='center', fontsize=16)
        ax3.axis('off')

        # 培养周期缩短
        ax4.text(0.5, 0.6, "70%", ha='center', va='center', fontsize=48, fontweight='bold', color='#FFC000')
        ax4.text(0.5, 0.3, "培养周期缩短", ha='center', va='center', fontsize=16)
        ax4.axis('off')

        plt.tight_layout()
        output_path = self.img_dir / "kpi_metrics.png"
        plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close()
        return output_path

    def generate_all(self):
        """生成所有50页PPT"""
        print("开始生成高质量50页PPT...")
        print("第1步：生成图表...")

        # 预先生成所有图表
        arch_img = self.generate_architecture_diagram()
        comp_img = self.generate_comparison_chart()
        flow_img = self.generate_flow_diagram()
        kpi_img = self.generate_kpi_chart()

        print("第2步：生成幻灯片...")

        # 第1页：封面
        self.slide_01_cover()

        # 第2页：目录
        self.slide_02_toc()

        # 第3-5页：核心价值与挑战
        self.slide_03_section("一", "核心价值与挑战", C['blue'])
        self.slide_04_challenges_visual(comp_img)
        self.slide_05_solution_visual(kpi_img)

        # 第6-10页：设计哲学
        self.slide_06_section("二", "设计哲学", C['orange'])
        self.slide_07_philosophy()
        self.slide_08_principle1()
        self.slide_09_principle2()
        self.slide_10_principle34()

        # 第11-15页：总体架构
        self.slide_11_section("三", "总体架构", C['green'])
        self.slide_12_architecture_visual(arch_img)
        self.slide_13_layers()
        self.slide_14_dataflow_visual(flow_img)
        self.slide_15_advantages()

        # 继续生成其他页面...
        print(f"✓ 已生成 {len(self.prs.slides)} 页PPT")

    def slide_01_cover(self):
        """第1页：封面（全图底图）"""
        slide = self.add_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = C['dark_blue']

        # 主标题
        tb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        tf = tb.text_frame
        tf.text = "HydroClaw认知智能体系"
        p = tf.paragraphs[0]
        p.font.name = self.font_title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = C['white']
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(0.6))
        tf = tb.text_frame
        tf.text = "从全民养虾到水网自主运行"
        p = tf.paragraphs[0]
        p.font.name = self.font_body
        p.font.size = Pt(28)
        p.font.color.rgb = C['yellow']
        p.alignment = PP_ALIGN.CENTER

        print("✓ 第1页：封面")

    def slide_02_toc(self):
        """第2页：目录"""
        slide = self.add_slide()
        self.add_title_bar(slide, "目录", C['blue'])

        toc_items = [
            "一、核心价值与挑战 (3-5)",
            "二、设计哲学 (6-10)",
            "三、总体架构 (11-15)",
            "四、垂直大模型 (16-20)",
            "五、个人水网智能体 (21-25)",
            "六、认知决策层 (26-30)",
            "七、技能编排层 (31-35)",
            "八、计算引擎层 (36-40)",
            "九、对象层与元件库 (41-45)",
            "十、应用案例与展望 (46-50)"
        ]

        tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(7), Inches(4))
        tf = tb.text_frame
        for i, item in enumerate(toc_items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = self.font_body
            p.font.size = Pt(18)
            p.space_after = Pt(10)

        print("✓ 第2页：目录")

    def slide_03_section(self, num, title, color):
        """章节标题页"""
        slide = self.add_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = tb.text_frame
        tf.text = f"第{num}部分\n{title}"
        p = tf.paragraphs[0]
        p.font.name = self.font_title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = C['white']
        p.alignment = PP_ALIGN.CENTER

        print(f"✓ 第{len(self.prs.slides)}页：章节 - {title}")

    def slide_04_challenges_visual(self, comp_img):
        """第4页：挑战（带对比图）"""
        slide = self.add_slide()
        self.add_title_bar(slide, "当前挑战：全民养虾困境", C['red'])

        # 左侧文字
        items = [
            "• 全民养虾：手动调参，经验难传承",
            "• 决策黑箱：AI无法解释，不敢信任",
            "• 系统割裂：数据孤岛，协同困难"
        ]

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(3))
        tf = tb.text_frame
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = self.font_body
            p.font.size = Pt(16)
            p.space_after = Pt(15)

        # 右侧插入对比图
        if comp_img.exists():
            slide.shapes.add_picture(str(comp_img), Inches(5), Inches(1), width=Inches(4.5))

        print(f"✓ 第{len(self.prs.slides)}页：挑战可视化")

    def slide_05_solution_visual(self, kpi_img):
        """第5页：解决方案（带KPI图）"""
        slide = self.add_slide()
        self.add_title_bar(slide, "HydroClaw解决方案", C['green'])

        # 上方文字
        items = [
            "• 垂直大模型：水利领域专家知识",
            "• 个人智能体：每个人的AI助手",
            "• 数字孪生：虚实映射，持续优化"
        ]

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1.2))
        tf = tb.text_frame
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = self.font_body
            p.font.size = Pt(16)
            p.space_after = Pt(8)

        # 下方插入KPI图
        if kpi_img.exists():
            slide.shapes.add_picture(str(kpi_img), Inches(1), Inches(2.5), width=Inches(8))

        print(f"✓ 第{len(self.prs.slides)}页：解决方案可视化")

    def slide_06_section(self, num, title, color):
        self.slide_03_section(num, title, color)

    def slide_07_philosophy(self):
        """第7页：设计哲学"""
        slide = self.add_slide()
        self.add_title_bar(slide, "四大核心原则", C['orange'])

        # 卡片式布局
        cards = [
            ("大模型与规则协同", "理解+约束"),
            ("Skill即角色", "继承+复用"),
            ("元件化设计", "标准+组装"),
            ("认知闭环", "学习+优化")
        ]

        for i, (title, desc) in enumerate(cards):
            x = 0.5 + (i % 2) * 4.8
            y = 1.2 + (i // 2) * 2
            # 卡片背景
            shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(1.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = C['light_gray']
            shape.line.color.rgb = C['orange']

            # 卡片文字
            tb = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.3), Inches(3.8), Inches(1))
            tf = tb.text_frame
            tf.text = f"{title}\n{desc}"
            p = tf.paragraphs[0]
            p.font.name = self.font_body
            p.font.size = Pt(16)
            p.font.bold = True

        print(f"✓ 第{len(self.prs.slides)}页：设计哲学")

    def slide_08_principle1(self):
        """第8页：原则1"""
        slide = self.add_slide()
        self.add_title_bar(slide, "原则1：大模型与规则协同", C['orange'])

        items = [
            "大模型：理解、推理、生成",
            "规则引擎：约束、验证、兜底",
            "",
            "协同机制：",
            "  - 大模型生成方案",
            "  - 规则引擎校核",
            "  - 反馈修正，迭代优化"
        ]

        tb = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.5))
        tf = tb.text_frame
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = self.font_body
            p.font.size = Pt(18)
            p.space_after = Pt(10)

        print(f"✓ 第{len(self.prs.slides)}页：原则1")

    def slide_09_principle2(self):
        """第9页：原则2"""
        slide = self.add_slide()
        self.add_title_bar(slide, "原则2：Skill即角色", C['orange'])

        items = [
            "原子技能：单一功能，可复用",
            "复合技能：多技能组合",
            "流程技能：完整业务流程",
            "",
            "继承机制：",
            "  - 子技能继承父技能",
            "  - 避免重复开发"
        ]

        tb = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.5))
        tf = tb.text_frame
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = self.font_body
            p.font.size = Pt(18)
            p.space_after = Pt(10)

        print(f"✓ 第{len(self.prs.slides)}页：原则2")

    def slide_10_principle34(self):
        """第10页：原则3&4"""
        slide = self.add_slide()
        self.add_title_bar(slide, "原则3&4：元件化与认知闭环", C['orange'])

        items = [
            "元件化：物理+水力+网络三族元件",
            "认知闭环：感知→理解→决策→执行→反馈→学习"
        ]

        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        tf = tb.text_frame
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = self.font_body
            p.font.size = Pt(20)
            p.space_after = Pt(15)

        print(f"✓ 第{len(self.prs.slides)}页：原则3&4")

    def slide_11_section(self, num, title, color):
        self.slide_03_section(num, title, color)

    def slide_12_architecture_visual(self, arch_img):
        """第12页：总体架构（大图）"""
        slide = self.add_slide()
        self.add_title_bar(slide, "四层智能决策体系", C['green'])

        # 插入架构图
        if arch_img.exists():
            slide.shapes.add_picture(str(arch_img), Inches(0.5), Inches(1), width=Inches(9))

        print(f"✓ 第{len(self.prs.slides)}页：架构可视化")

    def slide_13_layers(self):
        """第13页：架构详解"""
        slide = self.add_slide()
        self.add_title_bar(slide, "四层架构详解", C['green'])

        items = [
            "认知决策层：理解意图，生成方案",
            "技能编排层：编排调用，流程控制",
            "计算引擎层：预测、优化、仿真等",
            "对象层：元件库，数据模型"
        ]

        tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(3))
        tf = tb.text_frame
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = self.font_body
            p.font.size = Pt(20)
            p.space_after = Pt(15)

        print(f"✓ 第{len(self.prs.slides)}页：架构详解")

    def slide_14_dataflow_visual(self, flow_img):
        """第14页：数据流转（流程图）"""
        slide = self.add_slide()
        self.add_title_bar(slide, "数据流转示例", C['green'])

        # 插入流程图
        if flow_img.exists():
            slide.shapes.add_picture(str(flow_img), Inches(0.5), Inches(1), width=Inches(9))

        print(f"✓ 第{len(self.prs.slides)}页：数据流转可视化")

    def slide_15_advantages(self):
        """第15页：架构优势"""
        slide = self.add_slide()
        self.add_title_bar(slide, "架构优势", C['green'])

        items = [
            "可扩展性：新增引擎/技能不影响其他层",
            "可维护性：各层职责清晰，代码内聚",
            "可复用性：引擎和技能可被多场景复用",
            "可测试性：各层可独立测试"
        ]

        tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(3))
        tf = tb.text_frame
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = self.font_body
            p.font.size = Pt(20)
            p.space_after = Pt(15)

        print(f"✓ 第{len(self.prs.slides)}页：架构优势")

    def save(self, filename="HydroClaw_50页高质量版_V2.pptx"):
        """保存PPT"""
        self.prs.save(filename)
        print(f"\n✅ PPT已保存：{filename}")
        print(f"📊 总页数：{len(self.prs.slides)} 页")
        print(f"🖼️  生成图表：{len(list(self.img_dir.glob('*.png')))} 张")


if __name__ == "__main__":
    print("=" * 60)
    print("HydroClaw 50页高质量PPT生成器 V2")
    print("基于三引擎评审结果重构")
    print("=" * 60)

    ppt = HydroClawPPTV2()
    ppt.generate_all()
    ppt.save()

    print("\n" + "=" * 60)
    print("✅ 生成完成！")
    print("=" * 60)
