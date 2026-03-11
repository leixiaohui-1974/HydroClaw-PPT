#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HydroClaw 50页完整PPT生成器
基于发布会素材和文档内容
"""

import sys
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# 设置标准输出编码为UTF-8
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

# 配色方案
C = {
    'blue': RGBColor(91, 155, 213),
    'orange': RGBColor(237, 125, 49),
    'green': RGBColor(112, 173, 71),
    'yellow': RGBColor(255, 192, 0),
    'purple': RGBColor(142, 124, 195),
    'cyan': RGBColor(68, 114, 196),
    'red': RGBColor(192, 0, 0),
    'dark_blue': RGBColor(31, 78, 120),
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(127, 127, 127),
}

class HydroClawPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

    def add_slide(self, layout_idx=6):
        return self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

    def add_title_bar(self, slide, title, color):
        """添加标题栏"""
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
        tf = tb.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C['white']

    def add_bullets(self, slide, items, top=1.2):
        """添加项目符号列表"""
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(8.4), Inches(4))
        tf = tb.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(18 if not item.startswith(('1.', '2.', '3.', '4.', '5.')) else 20)
            p.font.bold = item.startswith(('1.', '2.', '3.', '4.', '5.'))
            if item.startswith(('1.', '2.', '3.', '4.', '5.')):
                p.font.color.rgb = C['dark_blue']
            p.space_after = Pt(8)

    def generate_all(self):
        """生成所有50页"""
        print("开始生成50页PPT...")

        # 第1页：封面
        self.slide_01_cover()
        # 第2页：目录
        self.slide_02_toc()

        # 第一部分：核心价值与挑战（3-5页）
        self.slide_03_section_title("一", "核心价值与挑战", "从全民养虾到智能决策", C['blue'])
        self.slide_04_challenges()
        self.slide_05_solution()

        # 第二部分：设计哲学（6-10页）
        self.slide_06_section_title("二", "设计哲学", "四大核心原则", C['orange'])
        self.slide_07_philosophy_overview()
        self.slide_08_principle1()
        self.slide_09_principle2()
        self.slide_10_principle34()

        # 第三部分：总体架构（11-15页）
        self.slide_11_section_title("三", "总体架构", "四层智能决策体系", C['green'])
        self.slide_12_architecture_overview()
        self.slide_13_architecture_layers()
        self.slide_14_data_flow()
        self.slide_15_architecture_advantages()

        # 第四部分：垂直大模型（16-20页）
        self.slide_16_section_title("四", "垂直大模型", "水利领域专家知识", C['purple'])
        self.slide_17_llm_comparison()
        self.slide_18_knowledge_base()
        self.slide_19_case_library()
        self.slide_20_rule_engine()

        # 第五部分：个人水网智能体（21-25页）
        self.slide_21_section_title("五", "个人水网智能体", "每个人的AI助手", C['cyan'])
        self.slide_22_agent_concept()
        self.slide_23_designer_scenario()
        self.slide_24_operator_scenario()
        self.slide_25_researcher_scenario()

        # 第六部分：认知决策层（26-30页）
        self.slide_26_section_title("六", "认知决策层", "大模型与规则协同", C['blue'])
        self.slide_27_llm_rule_division()
        self.slide_28_rule_inheritance()
        self.slide_29_reply_templates()
        self.slide_30_decision_process()

        # 第七部分：技能编排层（31-35页）
        self.slide_31_section_title("七", "技能编排层", "Skill即角色", C['orange'])
        self.slide_32_skill_inheritance()
        self.slide_33_atomic_skills()
        self.slide_34_composite_skills()
        self.slide_35_process_skills()

        # 第八部分：计算引擎层（36-40页）
        self.slide_36_section_title("八", "计算引擎层", "七大通用引擎", C['green'])
        self.slide_37_engines_overview()
        self.slide_38_predict_optimize()
        self.slide_39_simulate_learn()
        self.slide_40_verify_visualize()

        # 第九部分：对象层与元件库（41-45页）
        self.slide_41_section_title("九", "对象层与元件库", "三族元件体系", C['yellow'])
        self.slide_42_component_overview()
        self.slide_43_physical_components()
        self.slide_44_hydro_components()
        self.slide_45_network_assembly()

        # 第十部分：应用案例与展望（46-50页）
        self.slide_46_section_title("十", "应用案例与展望", "从理论到实践", C['red'])
        self.slide_47_case_study()
        self.slide_48_penetration_strategy()
        self.slide_49_success_cases()
        self.slide_50_conclusion()

        print(f"✓ 已生成 {len(self.prs.slides)} 页PPT")

    # ========== 具体页面生成方法 ==========

    def slide_01_cover(self):
        """第1页：封面"""
        slide = self.add_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = C['dark_blue']

        # 主标题
        tb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        tf = tb.text_frame
        tf.text = "HydroClaw认知智能体系"
        p = tf.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = C['white']
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(0.6))
        tf = tb.text_frame
        tf.text = "从全民养虾到水网自主运行"
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = C['yellow']
        p.alignment = PP_ALIGN.CENTER

        # 核心概念
        tb = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(6), Inches(0.5))
        tf = tb.text_frame
        tf.text = "垂直大模型 + 个人水网智能体 + 数字孪生"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = C['cyan']
        p.alignment = PP_ALIGN.CENTER

        print("✓ 第1页：封面")

    def slide_02_toc(self):
        """第2页：目录"""
        slide = self.add_slide()

        # 标题
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = tb.text_frame
        tf.text = "目录"
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = C['blue']

        toc = [
            ("一、核心价值与挑战", "3-5"),
            ("二、设计哲学", "6-10"),
            ("三、总体架构", "11-15"),
            ("四、垂直大模型", "16-20"),
            ("五、个人水网智能体", "21-25"),
            ("六、认知决策层", "26-30"),
            ("七、技能编排层", "31-35"),
            ("八、计算引擎层", "36-40"),
            ("九、对象层与元件库", "41-45"),
            ("十、应用案例与展望", "46-50"),
        ]

        # 左列
        tb = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(4), Inches(4))
        tf = tb.text_frame
        for i, (title, pages) in enumerate(toc[:5]):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"{title}  ......  {pages}"
            p.font.size = Pt(16)
            p.font.color.rgb = C['dark_blue']
            p.space_after = Pt(12)

        # 右列
        tb = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(4))
        tf = tb.text_frame
        for i, (title, pages) in enumerate(toc[5:]):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"{title}  ......  {pages}"
            p.font.size = Pt(16)
            p.font.color.rgb = C['dark_blue']
            p.space_after = Pt(12)

        print("✓ 第2页：目录")

    def slide_03_section_title(self, num, title, subtitle, color):
        """章节标题页模板"""
        slide = self.add_slide()

        # 背景
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        # 章节号
        tb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        tf = tb.text_frame
        tf.text = f"第{num}部分"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.color.rgb = C['white']
        p.alignment = PP_ALIGN.CENTER

        # 标题
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(1))
        tf = tb.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = C['white']
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        if subtitle:
            tb = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(0.5))
            tf = tb.text_frame
            tf.text = subtitle
            p = tf.paragraphs[0]
            p.font.size = Pt(20)
            p.font.color.rgb = C['yellow']
            p.alignment = PP_ALIGN.CENTER

        print(f"✓ 第{len(self.prs.slides)}页：章节标题 - {title}")

    def slide_04_challenges(self):
        """第4页：当前挑战"""
        slide = self.add_slide()
        self.add_title_bar(slide, "当前挑战：全民养虾困境", C['red'])

        items = [
            "当前水网运行的三大困境",
            "",
            "1. 全民养虾现象",
            "• 每个工程师都在手动调参、试错",
            "• 经验无法传承，知识无法复用",
            "• 新人培养周期长达3-5年",
            "",
            "2. 决策黑箱问题",
            "• AI给出方案，但无法解释原因",
            "• 工程师不敢信任，不敢采用",
            "• 责任归属不清晰",
            "",
            "3. 系统割裂困境",
            "• 预测、优化、仿真各自为政",
            "• 数据孤岛严重，协同困难",
            "• 重复开发，效率低下",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：当前挑战")

    def slide_05_solution(self):
        """第5页：HydroClaw解决方案"""
        slide = self.add_slide()
        self.add_title_bar(slide, "HydroClaw解决方案", C['green'])

        items = [
            "三位一体的智能决策体系",
            "",
            "1. 垂直大模型：水利领域专家知识",
            "• 整合规范、案例、物理约束",
            "• 规则引擎兜底，确保安全",
            "• 可解释性强，工程师可信任",
            "",
            "2. 个人水网智能体：每个人的AI助手",
            "• 设计工程师：自动生成方案",
            "• 运维人员：智能诊断故障",
            "• 科研人员：快速验证假设",
            "",
            "3. 数字孪生：虚实映射",
            "• 实时仿真，预测未来",
            "• 闭环优化，持续学习",
            "• 知识沉淀，经验复用",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：解决方案")

    # 第二部分方法
    def slide_06_section_title(self, num, title, subtitle, color):
        self.slide_03_section_title(num, title, subtitle, color)

    def slide_07_philosophy_overview(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "设计哲学：四大核心原则", C['orange'])
        items = [
            "HydroClaw的设计理念",
            "",
            "1. 大模型与规则协同",
            "• 大模型负责理解、推理、生成",
            "• 规则引擎负责约束、验证、兜底",
            "",
            "2. Skill即角色",
            "• 每个技能都是一个专业角色",
            "• 继承机制实现知识复用",
            "",
            "3. 元件化设计",
            "• 物理元件、水力元件、网络元件",
            "• 像搭积木一样组装水网",
            "",
            "4. 认知闭环",
            "• 感知→理解→决策→执行→反馈",
            "• 持续学习，不断优化",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：设计哲学概览")

    def slide_08_principle1(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "原则1：大模型与规则协同", C['orange'])
        items = [
            "为什么需要规则引擎？",
            "",
            "大模型的优势：",
            "• 理解自然语言，降低使用门槛",
            "• 推理能力强，处理复杂场景",
            "• 生成能力强，自动编写代码",
            "",
            "大模型的局限：",
            "• 可能产生幻觉，给出错误答案",
            "• 缺乏领域约束，违反物理规律",
            "• 无法保证100%正确",
            "",
            "规则引擎的作用：",
            "• 硬约束：流速、压力、管径等物理限制",
            "• 软约束：设计规范、经验法则",
            "• 兜底机制：大模型失效时接管",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：原则1")

    def slide_09_principle2(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "原则2：Skill即角色", C['orange'])
        items = [
            "技能继承体系",
            "",
            "基础技能（Atomic Skill）：",
            "• 单一功能，如计算流量、绘制图表",
            "• 可复用，可组合",
            "",
            "复合技能（Composite Skill）：",
            "• 多个基础技能组合",
            "• 如水力计算 = 流量计算 + 压力计算 + 校核",
            "",
            "流程技能（Process Skill）：",
            "• 完整业务流程",
            "• 如方案设计 = 需求分析 + 方案生成 + 优化 + 出图",
            "",
            "继承机制：",
            "• 子技能继承父技能的知识和能力",
            "• 避免重复开发，提高复用率",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：原则2")

    def slide_10_principle34(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "原则3&4：元件化与认知闭环", C['orange'])
        items = [
            "元件化设计：",
            "• 物理元件：管道、水泵、阀门、水池",
            "• 水力元件：节点、管段、边界条件",
            "• 网络元件：拓扑结构、连接关系",
            "• 优势：标准化、可复用、易维护",
            "",
            "认知闭环：",
            "• 感知：传感器数据、用户输入",
            "• 理解：大模型解析意图",
            "• 决策：规则引擎约束 + 优化算法",
            "• 执行：调用计算引擎",
            "• 反馈：结果评估 + 知识更新",
            "• 学习：案例库积累 + 模型微调",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：原则3&4")

    # 第三部分方法
    def slide_11_section_title(self, num, title, subtitle, color):
        self.slide_03_section_title(num, title, subtitle, color)

    def slide_12_architecture_overview(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "总体架构：四层智能决策体系", C['green'])
        items = [
            "从上到下的四层架构",
            "",
            "第一层：认知决策层",
            "• 垂直大模型 + 规则引擎",
            "• 理解用户意图，生成决策方案",
            "",
            "第二层：技能编排层",
            "• Skill继承体系",
            "• 编排调用下层引擎",
            "",
            "第三层：计算引擎层",
            "• 7大通用引擎：预测、优化、仿真等",
            "• 提供基础计算能力",
            "",
            "第四层：对象层与元件库",
            "• 物理元件、水力元件、网络元件",
            "• 数据模型与业务对象",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：架构概览")

    def slide_13_architecture_layers(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "四层架构详解", C['green'])
        items = [
            "各层职责与交互",
            "",
            "认知决策层 → 技能编排层：",
            "• 传递：用户意图、约束条件",
            "• 接收：执行结果、异常信息",
            "",
            "技能编排层 → 计算引擎层：",
            "• 传递：计算任务、参数配置",
            "• 接收：计算结果、性能指标",
            "",
            "计算引擎层 → 对象层：",
            "• 传递：CRUD操作、查询请求",
            "• 接收：对象数据、拓扑关系",
            "",
            "优势：",
            "• 分层解耦，各层独立演进",
            "• 职责清晰，易于维护扩展",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：架构详解")

    def slide_14_data_flow(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "数据流转示例", C['green'])
        items = [
            "用户请求：'优化XX水厂的调度方案'",
            "",
            "1. 认知决策层：",
            "   • 大模型理解意图：需要优化调度",
            "   • 规则引擎约束：满足供水需求、压力范围",
            "",
            "2. 技能编排层：",
            "   • 调用'调度优化'技能",
            "   • 分解为：数据准备 + 优化计算 + 结果评估",
            "",
            "3. 计算引擎层：",
            "   • 预测引擎：预测未来24h用水量",
            "   • 优化引擎：求解最优调度方案",
            "   • 仿真引擎：验证方案可行性",
            "",
            "4. 对象层：",
            "   • 读取水厂、管网、水池等对象数据",
            "   • 更新调度方案到数据库",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：数据流转")

    def slide_15_architecture_advantages(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "架构优势", C['green'])
        items = [
            "为什么选择四层架构？",
            "",
            "1. 可扩展性",
            "• 新增引擎：只需在计算引擎层添加",
            "• 新增技能：只需在技能编排层添加",
            "• 不影响其他层",
            "",
            "2. 可维护性",
            "• 各层职责清晰，代码内聚",
            "• 修改某层不影响其他层",
            "",
            "3. 可复用性",
            "• 引擎可被多个技能调用",
            "• 技能可被多个场景复用",
            "",
            "4. 可测试性",
            "• 各层可独立测试",
            "• Mock接口，隔离依赖",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：架构优势")

    # 第四部分方法
    def slide_16_section_title(self, num, title, subtitle, color):
        self.slide_03_section_title(num, title, subtitle, color)

    def slide_17_llm_comparison(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "通用大模型 vs 垂直大模型", C['purple'])
        items = [
            "为什么需要垂直大模型？",
            "",
            "通用大模型（GPT-4、Claude等）：",
            "• 优势：知识广博，理解能力强",
            "• 劣势：缺乏水利专业知识",
            "• 问题：可能给出不符合规范的方案",
            "",
            "垂直大模型（HydroClaw-LLM）：",
            "• 基于通用大模型微调",
            "• 注入水利领域知识：规范、案例、经验",
            "• 结合规则引擎，确保合规性",
            "",
            "优势：",
            "• 专业性强，符合行业规范",
            "• 可解释性强，工程师可信任",
            "• 安全性高，规则引擎兜底",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：大模型对比")

    def slide_18_knowledge_base(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "知识库：三大知识来源", C['purple'])
        items = [
            "垂直大模型的知识来源",
            "",
            "1. 规范标准（Regulations）",
            "• 国家标准：GB 50015、GB 50013等",
            "• 行业规范：CJJ 92、CJJ 140等",
            "• 地方标准：各省市技术规程",
            "",
            "2. 案例库（Cases）",
            "• 历史项目：设计方案、运行数据",
            "• 典型案例：成功经验、失败教训",
            "• 专家经验：老工程师的隐性知识",
            "",
            "3. 物理约束（Physics）",
            "• 水力学原理：连续性方程、能量方程",
            "• 设备特性：水泵曲线、阀门特性",
            "• 边界条件：流量范围、压力限制",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：知识库")

    def slide_19_case_library(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "案例库：经验沉淀与复用", C['purple'])
        items = [
            "案例库的价值",
            "",
            "案例结构：",
            "• 场景描述：用户需求、边界条件",
            "• 解决方案：设计参数、计算过程",
            "• 结果评估：性能指标、优缺点",
            "• 经验总结：注意事项、优化建议",
            "",
            "案例检索：",
            "• 向量化表示：将案例转为向量",
            "• 相似度匹配：找到最相似的历史案例",
            "• 迁移学习：借鉴历史经验",
            "",
            "案例积累：",
            "• 每次设计都沉淀为案例",
            "• 专家审核，确保质量",
            "• 持续积累，越用越智能",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：案例库")

    def slide_20_rule_engine(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "规则引擎：安全兜底机制", C['purple'])
        items = [
            "规则引擎的三层防护",
            "",
            "第一层：硬约束（Hard Constraints）",
            "• 物理规律：流速、压力、管径",
            "• 不可违反，否则方案无效",
            "",
            "第二层：软约束（Soft Constraints）",
            "• 设计规范：推荐值、经验范围",
            "• 可以违反，但需要说明理由",
            "",
            "第三层：优化目标（Objectives）",
            "• 经济性：投资最小、运行费用最低",
            "• 可靠性：供水保证率、冗余度",
            "• 环保性：能耗、碳排放",
            "",
            "工作流程：",
            "• 大模型生成方案 → 规则引擎校核",
            "• 违反硬约束 → 拒绝方案，重新生成",
            "• 违反软约束 → 警告提示，由用户决定",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：规则引擎")

    # 第五部分方法
    def slide_21_section_title(self, num, title, subtitle, color):
        self.slide_03_section_title(num, title, subtitle, color)

    def slide_22_agent_concept(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "个人水网智能体：每个人的AI助手", C['cyan'])
        items = [
            "什么是个人水网智能体？",
            "",
            "传统模式：",
            "• 工程师需要学习复杂的软件",
            "• 需要掌握大量专业知识",
            "• 新人培养周期长",
            "",
            "智能体模式：",
            "• 自然语言交互，降低使用门槛",
            "• AI助手辅助决策，提高效率",
            "• 知识自动沉淀，经验快速复用",
            "",
            "三类典型用户：",
            "• 设计工程师：自动生成设计方案",
            "• 运维人员：智能诊断故障",
            "• 科研人员：快速验证假设",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：智能体概念")

    def slide_23_designer_scenario(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "场景1：设计工程师", C['cyan'])
        items = [
            "用户：'设计一个供水量5000m³/d的水厂'",
            "",
            "智能体工作流程：",
            "",
            "1. 需求分析",
            "   • 解析用户意图：新建水厂设计",
            "   • 提取关键参数：供水量、服务人口",
            "",
            "2. 方案生成",
            "   • 检索相似案例：找到3个类似项目",
            "   • 生成初步方案：水源、水厂、管网",
            "   • 规则引擎校核：符合规范",
            "",
            "3. 优化调整",
            "   • 多目标优化：投资、能耗、可靠性",
            "   • 生成3个备选方案",
            "",
            "4. 成果输出",
            "   • 自动出图：平面图、高程图",
            "   • 生成报告：设计说明、计算书",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：设计工程师场景")

    def slide_24_operator_scenario(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "场景2：运维人员", C['cyan'])
        items = [
            "用户：'XX区域水压偏低，帮我诊断原因'",
            "",
            "智能体工作流程：",
            "",
            "1. 数据采集",
            "   • 读取实时监测数据：流量、压力",
            "   • 读取历史数据：对比正常值",
            "",
            "2. 故障诊断",
            "   • 水力仿真：模拟当前工况",
            "   • 原因分析：管道堵塞？阀门误操作？",
            "   • 给出诊断结果：3个可能原因",
            "",
            "3. 解决方案",
            "   • 生成应急方案：调整水泵、开启阀门",
            "   • 仿真验证：方案可行性",
            "",
            "4. 执行反馈",
            "   • 推送操作指令到SCADA",
            "   • 监测效果，持续优化",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：运维人员场景")

    def slide_25_researcher_scenario(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "场景3：科研人员", C['cyan'])
        items = [
            "用户：'验证新型调度算法的效果'",
            "",
            "智能体工作流程：",
            "",
            "1. 实验设计",
            "   • 理解研究目标：对比新旧算法",
            "   • 设计实验方案：测试场景、评价指标",
            "",
            "2. 数据准备",
            "   • 生成测试数据：典型工况、极端工况",
            "   • 构建仿真模型：数字孪生水网",
            "",
            "3. 算法对比",
            "   • 运行新算法：求解优化问题",
            "   • 运行基准算法：传统方法",
            "   • 仿真验证：评估性能",
            "",
            "4. 结果分析",
            "   • 生成对比图表：能耗、成本、可靠性",
            "   • 撰写分析报告：结论、建议",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：科研人员场景")

    # 第六部分方法
    def slide_26_section_title(self, num, title, subtitle, color):
        self.slide_03_section_title(num, title, subtitle, color)

    def slide_27_llm_rule_division(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "大模型与规则的分工", C['blue'])
        items = [
            "如何划分大模型和规则引擎的职责？",
            "",
            "大模型负责：",
            "• 理解：解析用户自然语言输入",
            "• 推理：分析问题，制定策略",
            "• 生成：编写代码、生成报告",
            "• 学习：从案例中总结经验",
            "",
            "规则引擎负责：",
            "• 约束：物理规律、设计规范",
            "• 校核：验证方案合规性",
            "• 兜底：大模型失效时接管",
            "• 解释：说明决策依据",
            "",
            "协同机制：",
            "• 大模型生成 → 规则校核 → 反馈修正",
            "• 迭代优化，直到满足约束",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：大模型与规则分工")

    def slide_28_rule_inheritance(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "规则继承体系", C['blue'])
        items = [
            "规则的层次化组织",
            "",
            "基础规则（Base Rules）：",
            "• 物理规律：连续性方程、能量方程",
            "• 适用于所有场景",
            "",
            "领域规则（Domain Rules）：",
            "• 给水规范：管径、流速、压力",
            "• 排水规范：坡度、充满度",
            "",
            "场景规则（Scenario Rules）：",
            "• 设计场景：管网设计、水厂设计",
            "• 运维场景：调度优化、故障诊断",
            "",
            "继承机制：",
            "• 子规则继承父规则",
            "• 可以覆盖（Override）父规则",
            "• 避免重复定义，提高维护性",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：规则继承")

    def slide_29_reply_templates(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "回复模板：结构化输出", C['blue'])
        items = [
            "为什么需要回复模板？",
            "",
            "问题：",
            "• 大模型输出格式不统一",
            "• 难以解析和处理",
            "",
            "解决方案：",
            "• 定义结构化回复模板",
            "• 大模型按模板生成输出",
            "",
            "模板示例（设计方案）：",
            "• 方案概述：设计思路、关键参数",
            "• 详细设计：管网布局、设备选型",
            "• 计算结果：流量、压力、能耗",
            "• 合规性说明：符合哪些规范",
            "• 优化建议：可改进的地方",
            "",
            "优势：",
            "• 输出规范，易于解析",
            "• 可解释性强，工程师易理解",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：回复模板")

    def slide_30_decision_process(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "决策流程示例", C['blue'])
        items = [
            "用户请求：'设计一个小区供水管网'",
            "",
            "Step 1：意图理解（大模型）",
            "• 识别任务类型：管网设计",
            "• 提取参数：小区规模、用水量",
            "",
            "Step 2：方案生成（大模型）",
            "• 检索相似案例",
            "• 生成初步方案：管网拓扑、管径",
            "",
            "Step 3：规则校核（规则引擎）",
            "• 检查流速：2-3 m/s ✓",
            "• 检查压力：0.28-0.35 MPa ✓",
            "• 检查管径：≥DN100 ✓",
            "",
            "Step 4：优化调整（大模型+优化引擎）",
            "• 多目标优化：投资最小、可靠性高",
            "",
            "Step 5：结果输出（大模型）",
            "• 生成设计说明、计算书、图纸",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：决策流程")

    # 第七部分方法
    def slide_31_section_title(self, num, title, subtitle, color):
        self.slide_03_section_title(num, title, subtitle, color)

    def slide_32_skill_inheritance(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "技能继承体系", C['orange'])
        items = [
            "Skill即角色的实现",
            "",
            "基类Skill：",
            "• 定义通用接口：execute()、validate()",
            "• 提供公共能力：日志、异常处理",
            "",
            "子类继承：",
            "• WaterSupplySkill：给水领域基础技能",
            "  - PipeDesignSkill：管网设计",
            "  - PumpSelectionSkill：水泵选型",
            "• DrainageSkill：排水领域基础技能",
            "  - SewerDesignSkill：排水管网设计",
            "",
            "优势：",
            "• 代码复用：避免重复开发",
            "• 知识传承：父类知识自动继承",
            "• 易于扩展：新增技能只需继承",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：技能继承")

    def slide_33_atomic_skills(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "原子技能：最小功能单元", C['orange'])
        items = [
            "原子技能示例",
            "",
            "1. 计算类技能",
            "• FlowCalculation：流量计算",
            "• PressureCalculation：压力计算",
            "• HeadLossCalculation：水头损失计算",
            "",
            "2. 查询类技能",
            "• ComponentQuery：元件查询",
            "• CaseRetrieval：案例检索",
            "• RuleQuery：规则查询",
            "",
            "3. 可视化技能",
            "• ChartPlotting：图表绘制",
            "• NetworkVisualization：管网可视化",
            "• ReportGeneration：报告生成",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：原子技能")

    def slide_34_composite_skills(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "复合技能：多技能组合", C['orange'])
        items = [
            "复合技能示例",
            "",
            "HydraulicAnalysis（水力分析）：",
            "• = FlowCalculation + PressureCalculation",
            "• + HeadLossCalculation + ResultValidation",
            "",
            "PumpStationDesign（泵站设计）：",
            "• = FlowPrediction + PumpSelection",
            "• + PipeDesign + CostEstimation",
            "",
            "NetworkOptimization（管网优化）：",
            "• = TopologyAnalysis + HydraulicAnalysis",
            "• + OptimizationSolver + ResultEvaluation",
            "",
            "优势：",
            "• 封装复杂流程，降低使用难度",
            "• 保证执行顺序，避免遗漏步骤",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：复合技能")

    def slide_35_process_skills(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "流程技能：完整业务流程", C['orange'])
        items = [
            "流程技能示例",
            "",
            "WaterSupplyDesign（供水系统设计）：",
            "1. 需求分析：用水量预测、服务范围",
            "2. 方案生成：水源选择、管网布局",
            "3. 水力计算：流量分配、压力校核",
            "4. 设备选型：水泵、阀门、水池",
            "5. 优化调整：多目标优化",
            "6. 成果输出：图纸、报告、清单",
            "",
            "FaultDiagnosis（故障诊断）：",
            "1. 数据采集：实时监测、历史数据",
            "2. 异常检测：识别异常工况",
            "3. 原因分析：水力仿真、规则推理",
            "4. 方案生成：应急措施、长期方案",
            "5. 效果评估：仿真验证、风险评估",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：流程技能")

    # 第八部分方法
    def slide_36_section_title(self, num, title, subtitle, color):
        self.slide_03_section_title(num, title, subtitle, color)

    def slide_37_engines_overview(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "七大通用引擎", C['green'])
        items = [
            "计算引擎层的核心能力",
            "",
            "1. 预测引擎（Prediction）",
            "• 用水量预测、水质预测",
            "",
            "2. 优化引擎（Optimization）",
            "• 调度优化、设计优化",
            "",
            "3. 仿真引擎（Simulation）",
            "• 水力仿真、水质仿真",
            "",
            "4. 学习引擎（Learning）",
            "• 模型训练、参数学习",
            "",
            "5. 验证引擎（Verification）",
            "• 方案校核、规则验证",
            "",
            "6. 可视化引擎（Visualization）",
            "• 图表生成、3D渲染",
            "",
            "7. 报告引擎（Reporting）",
            "• 自动生成设计说明、计算书",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：七大引擎")

    def slide_38_predict_optimize(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "预测引擎 + 优化引擎", C['green'])
        items = [
            "预测引擎（Prediction Engine）",
            "",
            "功能：",
            "• 时间序列预测：LSTM、Prophet",
            "• 空间预测：克里金插值、GNN",
            "• 多变量预测：用水量、水质、压力",
            "",
            "应用场景：",
            "• 短期预测：未来24h用水量",
            "• 中期预测：未来7天需水量",
            "• 长期预测：未来10年供水规模",
            "",
            "优化引擎（Optimization Engine）",
            "",
            "算法库：",
            "• 线性规划：LP、MILP",
            "• 非线性规划：SQP、Interior Point",
            "• 启发式算法：GA、PSO、SA",
            "",
            "应用场景：",
            "• 调度优化：最小化能耗、成本",
            "• 设计优化：最优管径、泵站位置",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：预测+优化引擎")

    def slide_39_simulate_learn(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "仿真引擎 + 学习引擎", C['green'])
        items = [
            "仿真引擎（Simulation Engine）",
            "",
            "核心能力：",
            "• 水力仿真：EPANET、SWMM",
            "• 水质仿真：氯衰减、污染扩散",
            "• 数字孪生：实时仿真、预测未来",
            "",
            "应用场景：",
            "• 方案验证：设计方案可行性",
            "• 故障诊断：模拟异常工况",
            "• 应急演练：爆管、污染事故",
            "",
            "学习引擎（Learning Engine）",
            "",
            "技术栈：",
            "• 监督学习：回归、分类",
            "• 无监督学习：聚类、降维",
            "• 强化学习：调度策略优化",
            "",
            "应用场景：",
            "• 模型校准：管网粗糙系数",
            "• 异常检测：漏损、异常用水",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：仿真+学习引擎")

    def slide_40_verify_visualize(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "验证引擎 + 可视化引擎 + 报告引擎", C['green'])
        items = [
            "验证引擎（Verification Engine）",
            "• 规则校核：物理约束、设计规范",
            "• 结果验证：合理性检查",
            "",
            "可视化引擎（Visualization Engine）",
            "• 2D图表：折线图、柱状图、热力图",
            "• 3D渲染：管网三维可视化",
            "• 动态演示：水流动画、压力变化",
            "",
            "报告引擎（Reporting Engine）",
            "• 模板管理：设计说明、计算书",
            "• 自动填充：参数、图表、结论",
            "• 格式输出：Word、PDF、HTML",
            "",
            "三引擎协同：",
            "• 验证 → 可视化 → 报告",
            "• 一键生成完整设计文档",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：验证+可视化+报告引擎")

    # 第九部分方法
    def slide_41_section_title(self, num, title, subtitle, color):
        self.slide_03_section_title(num, title, subtitle, color)

    def slide_42_component_overview(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "三族元件体系", C['yellow'])
        items = [
            "元件化设计的三个层次",
            "",
            "第一族：物理元件（Physical Components）",
            "• 真实设备：管道、水泵、阀门、水池",
            "• 属性：尺寸、材质、性能参数",
            "• 行为：开关、调节、监测",
            "",
            "第二族：水力元件（Hydraulic Components）",
            "• 计算抽象：节点、管段、边界条件",
            "• 属性：流量、压力、水头",
            "• 行为：水力计算、能量平衡",
            "",
            "第三族：网络元件（Network Components）",
            "• 拓扑结构：图、树、环",
            "• 属性：连接关系、路径、分区",
            "• 行为：拓扑分析、路径搜索",
            "",
            "三族映射：",
            "• 物理元件 ↔ 水力元件 ↔ 网络元件",
            "• 虚实映射，数字孪生",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：三族元件")

    def slide_43_physical_components(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "物理元件库", C['yellow'])
        items = [
            "常用物理元件",
            "",
            "管道（Pipe）：",
            "• 属性：管径、长度、材质、粗糙系数",
            "• 类型：铸铁管、钢管、PE管、PCCP",
            "",
            "水泵（Pump）：",
            "• 属性：流量、扬程、效率、功率",
            "• 类型：离心泵、轴流泵、潜水泵",
            "• 特性曲线：Q-H、Q-η、Q-P",
            "",
            "阀门（Valve）：",
            "• 属性：口径、开度、阻力系数",
            "• 类型：闸阀、蝶阀、止回阀、调节阀",
            "",
            "水池（Tank）：",
            "• 属性：容积、水位、进出水管",
            "• 类型：清水池、高位水池、调节池",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：物理元件")

    def slide_44_hydro_components(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "水力元件库", C['yellow'])
        items = [
            "水力计算抽象",
            "",
            "节点（Node）：",
            "• 属性：高程、需水量、压力",
            "• 类型：用户节点、连接节点",
            "• 方程：连续性方程 Σ Q_in = Σ Q_out",
            "",
            "管段（Link）：",
            "• 属性：流量、流速、水头损失",
            "• 类型：管道、水泵、阀门",
            "• 方程：能量方程 H_loss = f(Q, D, L)",
            "",
            "边界条件（Boundary）：",
            "• 固定水头：水源、水池",
            "• 固定流量：用户节点",
            "• 压力约束：最小服务水头",
            "",
            "求解器：",
            "• 稳态求解：Newton-Raphson",
            "• 动态求解：时间步进法",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：水力元件")

    def slide_45_network_assembly(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "网络元件与组装", C['yellow'])
        items = [
            "网络拓扑元件",
            "",
            "图（Graph）：",
            "• 节点集合 V = {v1, v2, ..., vn}",
            "• 边集合 E = {e1, e2, ..., em}",
            "• 邻接矩阵、关联矩阵",
            "",
            "拓扑分析：",
            "• 连通性检查：是否存在孤立节点",
            "• 环路识别：找出所有基本环",
            "• 路径搜索：最短路径、关键路径",
            "",
            "网络组装：",
            "• 像搭积木一样组装水网",
            "• 拖拽元件 → 连接 → 配置参数",
            "• 自动生成拓扑结构",
            "• 一键水力计算",
            "",
            "优势：",
            "• 降低建模门槛，提高效率",
            "• 标准化元件，保证质量",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：网络组装")

    # 第十部分方法
    def slide_46_section_title(self, num, title, subtitle, color):
        self.slide_03_section_title(num, title, subtitle, color)

    def slide_47_case_study(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "应用案例：某市供水系统优化", C['red'])
        items = [
            "项目背景",
            "• 供水规模：50万m³/d",
            "• 问题：能耗高、漏损率高、水压不稳",
            "",
            "HydroClaw解决方案",
            "",
            "1. 数字孪生建模",
            "   • 构建全市管网数字孪生模型",
            "   • 接入SCADA实时数据",
            "",
            "2. 智能诊断",
            "   • 识别高漏损区域：3个重点片区",
            "   • 分析水压不稳原因：泵站调度不合理",
            "",
            "3. 优化方案",
            "   • 调度优化：降低能耗15%",
            "   • 分区计量：漏损率从18%降至12%",
            "   • 压力优化：水压合格率从85%提升至95%",
            "",
            "效果：",
            "• 年节约电费：300万元",
            "• 减少漏损：年节水200万m³",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：应用案例")

    def slide_48_penetration_strategy(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "渗透策略：从点到面", C['red'])
        items = [
            "HydroClaw的推广路径",
            "",
            "第一阶段：单点突破（0-6个月）",
            "• 目标：设计院、咨询公司",
            "• 场景：管网设计、方案优化",
            "• 价值：提高设计效率3-5倍",
            "",
            "第二阶段：线性扩展（6-18个月）",
            "• 目标：水务公司、运维单位",
            "• 场景：调度优化、故障诊断",
            "• 价值：降低能耗15%、减少漏损30%",
            "",
            "第三阶段：面状覆盖（18-36个月）",
            "• 目标：政府部门、监管机构",
            "• 场景：智慧水务、城市大脑",
            "• 价值：全域优化、应急指挥",
            "",
            "关键成功因素：",
            "• 降低使用门槛：自然语言交互",
            "• 快速见效：1-2周出成果",
            "• 持续优化：越用越智能",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：渗透策略")

    def slide_49_success_cases(self):
        slide = self.add_slide()
        self.add_title_bar(slide, "成功案例集锦", C['red'])
        items = [
            "典型应用场景",
            "",
            "案例1：某新区供水系统设计",
            "• 用时：3天（传统方法需2周）",
            "• 效果：生成3套方案，投资节约12%",
            "",
            "案例2：某水厂调度优化",
            "• 用时：1周（传统方法需1个月）",
            "• 效果：能耗降低18%，年节约80万元",
            "",
            "案例3：某市管网漏损治理",
            "• 用时：2周（传统方法需3个月）",
            "• 效果：精准定位漏损点，漏损率降低8%",
            "",
            "案例4：某科研团队算法验证",
            "• 用时：3天（传统方法需2周）",
            "• 效果：快速验证新算法，发表SCI论文",
            "",
            "用户反馈：",
            "• '像有了一个24小时在线的专家顾问'",
            "• '新人也能快速上手，大大缩短培养周期'",
        ]
        self.add_bullets(slide, items)
        print(f"✓ 第{len(self.prs.slides)}页：成功案例")

    def slide_50_conclusion(self):
        slide = self.add_slide()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = C['dark_blue']

        # 标题
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.8))
        tf = tb.text_frame
        tf.text = "HydroClaw：让水网自主运行"
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = C['white']
        p.alignment = PP_ALIGN.CENTER

        # 核心价值
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(7), Inches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True

        content = [
            "核心价值：",
            "",
            "• 降低门槛：从'全民养虾'到自然语言交互",
            "• 提高效率：设计效率提升3-5倍",
            "• 保证质量：规则引擎兜底，符合规范",
            "• 知识传承：案例库积累，越用越智能",
            "",
            "愿景：",
            "让每个水利工程师都有自己的AI助手",
            "让水网像自动驾驶一样自主运行",
        ]

        for i, line in enumerate(content):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = line
            if line.startswith(('核心价值', '愿景')):
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = C['yellow']
            else:
                p.font.size = Pt(18)
                p.font.color.rgb = C['white']
            p.space_after = Pt(8)

        # 联系方式
        tb = slide.shapes.add_textbox(Inches(3), Inches(5), Inches(4), Inches(0.4))
        tf = tb.text_frame
        tf.text = "感谢聆听！"
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = C['cyan']
        p.alignment = PP_ALIGN.CENTER

        print(f"✓ 第{len(self.prs.slides)}页：结束页")

    def save(self, filename="HydroClaw_50页完整版.pptx"):
        """保存PPT"""
        self.prs.save(filename)
        print(f"\n✅ PPT已保存：{filename}")
        print(f"📊 总页数：{len(self.prs.slides)} 页")


if __name__ == "__main__":
    print("=" * 60)
    print("HydroClaw 50页完整PPT生成器")
    print("=" * 60)

    ppt = HydroClawPPT()
    ppt.generate_all()
    ppt.save()

    print("\n" + "=" * 60)
    print("✅ 生成完成！")
    print("=" * 60)
