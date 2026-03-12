#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPTX → JSON 内容提取器

从已有 .pptx 文件中提取内容结构，输出与 presenton_cli.py 兼容的 JSON 格式。

用法:
    python pptx_to_json.py input.pptx -o output.json --images-dir extracted_images/

输出 JSON 格式示例:
{
  "title": "演示文稿标题",
  "slides": [
    {"type": "cover", "title": "...", "subtitle": "..."},
    {"type": "content", "title": "...", "bullets": ["...", "..."]},
    ...
  ]
}
"""

import argparse
import hashlib
import json
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ============================================================
# 常量
# ============================================================
SW, SH = 1280, 720  # Presenton 坐标系 (pt)

# 中文数字映射
CN_NUMS = {
    "一": "一", "二": "二", "三": "三", "四": "四", "五": "五",
    "六": "六", "七": "七", "八": "八", "九": "九", "十": "十",
}


# ============================================================
# 工具函数
# ============================================================

def emu_to_pt(emu):
    """EMU → point"""
    if emu is None:
        return 0
    return emu / 12700


def rgb_to_hex(rgb_color):
    """RGBColor → hex string (without #)"""
    if rgb_color is None:
        return None
    try:
        return f"{rgb_color[0]:02X}{rgb_color[1]:02X}{rgb_color[2]:02X}"
    except (TypeError, IndexError):
        return None


def get_shape_bg_color(shape):
    """尝试获取 shape 的填充颜色"""
    try:
        fill = shape.fill
        if fill.type is not None:
            fg = fill.fore_color
            if fg and fg.type is not None:
                return rgb_to_hex(fg.rgb)
    except Exception:
        pass
    return None


def get_slide_bg_color(slide_obj):
    """获取幻灯片背景颜色"""
    try:
        bg = slide_obj.background
        fill = bg.fill
        if fill.type is not None:
            fg = fill.fore_color
            if fg and fg.type is not None:
                return rgb_to_hex(fg.rgb)
    except Exception:
        pass
    return None


def is_dark_color(hex_color):
    """判断颜色是否为深色"""
    if not hex_color:
        return False
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        luminance = 0.299 * r + 0.587 * g + 0.114 * b
        return luminance < 128
    except (ValueError, IndexError):
        return False


def is_blue_ish(hex_color):
    """判断颜色是否偏蓝（用于检测章节页）"""
    if not hex_color:
        return False
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return b > 150 and b > r * 1.3 and b > g * 1.1
    except (ValueError, IndexError):
        return False


def extract_text_from_shape(shape):
    """从 shape 中提取所有文本段落"""
    paragraphs = []
    if not shape.has_text_frame:
        return paragraphs
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    return paragraphs


def get_shape_font_size(shape):
    """获取 shape 中最大的字号（pt）"""
    max_size = 0
    if not shape.has_text_frame:
        return max_size
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                size_pt = run.font.size.pt
                if size_pt > max_size:
                    max_size = size_pt
    return max_size


def get_shape_center(shape):
    """获取 shape 中心点坐标 (x_pt, y_pt)"""
    left = emu_to_pt(shape.left) if shape.left else 0
    top = emu_to_pt(shape.top) if shape.top else 0
    width = emu_to_pt(shape.width) if shape.width else 0
    height = emu_to_pt(shape.height) if shape.height else 0
    return (left + width / 2, top + height / 2)


def get_shape_rect(shape):
    """获取 shape 的矩形区域 (left, top, width, height) in pt"""
    return (
        emu_to_pt(shape.left) if shape.left else 0,
        emu_to_pt(shape.top) if shape.top else 0,
        emu_to_pt(shape.width) if shape.width else 0,
        emu_to_pt(shape.height) if shape.height else 0,
    )


def shape_area_ratio(shape, slide_width, slide_height):
    """计算 shape 面积占幻灯片面积的比例"""
    _, _, w, h = get_shape_rect(shape)
    slide_area = slide_width * slide_height
    if slide_area == 0:
        return 0
    return (w * h) / slide_area


def image_hash(blob):
    """为图片 blob 生成短 hash"""
    return hashlib.md5(blob).hexdigest()[:8]


# ============================================================
# 图片提取
# ============================================================

def extract_image(shape, slide_idx, images_dir, image_counter):
    """
    从 picture shape 中提取图片，保存到 images_dir。
    返回保存的文件路径（相对路径）或 None。
    """
    try:
        image = shape.image
        blob = image.blob
        content_type = image.content_type
        ext_map = {
            "image/png": ".png",
            "image/jpeg": ".jpg",
            "image/gif": ".gif",
            "image/bmp": ".bmp",
            "image/tiff": ".tiff",
            "image/svg+xml": ".svg",
            "image/x-wmf": ".wmf",
            "image/x-emf": ".emf",
        }
        ext = ext_map.get(content_type, ".png")
        img_hash = image_hash(blob)
        filename = f"slide{slide_idx}_{image_counter}_{img_hash}{ext}"
        filepath = os.path.join(images_dir, filename)
        os.makedirs(images_dir, exist_ok=True)
        with open(filepath, "wb") as f:
            f.write(blob)
        return filepath
    except Exception as e:
        print(f"  Warning: Failed to extract image from slide {slide_idx}: {e}", file=sys.stderr)
        return None


# ============================================================
# 幻灯片分析
# ============================================================

class SlideAnalyzer:
    """分析单张幻灯片并提取结构化信息"""

    def __init__(self, slide_obj, slide_idx, total_slides, images_dir, slide_width_pt, slide_height_pt):
        self.slide = slide_obj
        self.idx = slide_idx
        self.total = total_slides
        self.images_dir = images_dir
        self.sw = slide_width_pt
        self.sh = slide_height_pt

        # 分析结果
        self.bg_color = get_slide_bg_color(slide_obj)
        self.bg_dark = is_dark_color(self.bg_color)
        self.bg_blue = is_blue_ish(self.bg_color)

        self.text_shapes = []       # (shape, texts, font_size, rect)
        self.image_shapes = []      # shape
        self.other_shapes = []      # shape

        self._classify_shapes()

    def _classify_shapes(self):
        """将 shapes 分为文本、图片、其他"""
        for shape in self.slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                self.image_shapes.append(shape)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # 递归检查 group 里的图片
                self._process_group(shape)
            elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                texts = extract_text_from_shape(shape)
                if texts:
                    font_size = get_shape_font_size(shape)
                    rect = get_shape_rect(shape)
                    self.text_shapes.append((shape, texts, font_size, rect))
            else:
                self.other_shapes.append(shape)

    def _process_group(self, group_shape):
        """处理组合 shape"""
        try:
            for shape in group_shape.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self.image_shapes.append(shape)
                elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    texts = extract_text_from_shape(shape)
                    if texts:
                        font_size = get_shape_font_size(shape)
                        rect = get_shape_rect(shape)
                        self.text_shapes.append((shape, texts, font_size, rect))
        except Exception:
            pass

    def get_title(self):
        """获取标题：最靠上区域中字号最大的文本"""
        if not self.text_shapes:
            return ""
        # 找最靠上的文本 top 位置
        min_top = min(rect[1] for _, _, _, rect in self.text_shapes)
        # 在顶部区域（top 值在最靠上 + 120pt 以内）中找最大字号
        top_region = [
            (shape, texts, fs, rect) for shape, texts, fs, rect in self.text_shapes
            if rect[1] < min_top + 120
        ]
        if top_region:
            best = max(top_region, key=lambda x: x[2])
            return "\n".join(best[1])
        # fallback: 最靠上的
        sorted_by_top = sorted(self.text_shapes, key=lambda x: x[3][1])
        return "\n".join(sorted_by_top[0][1])

    def get_body_texts(self, exclude_title=True):
        """获取正文文本（排除标题）"""
        if not self.text_shapes:
            return []
        title_text = self.get_title() if exclude_title else ""
        body = []
        # 按 top 位置排序
        sorted_shapes = sorted(self.text_shapes, key=lambda x: (x[3][1], x[3][0]))
        for shape, texts, font_size, rect in sorted_shapes:
            joined = "\n".join(texts)
            if exclude_title and joined == title_text:
                continue
            body.extend(texts)
        return body

    def get_texts_in_region(self, left_frac, right_frac):
        """获取指定水平区域（0-1 分数）内的文本"""
        left_bound = self.sw * left_frac
        right_bound = self.sw * right_frac
        result = []
        for shape, texts, font_size, rect in self.text_shapes:
            cx = rect[0] + rect[2] / 2
            if left_bound <= cx <= right_bound:
                result.extend(texts)
        return result

    def get_text_groups_lr(self):
        """将文本分为左半和右半两组"""
        left_texts = []
        right_texts = []
        mid = self.sw / 2
        # 排除标题（最靠上且最大的文本）
        title_text = self.get_title()
        for shape, texts, font_size, rect in self.text_shapes:
            joined = "\n".join(texts)
            if joined == title_text:
                continue
            cx = rect[0] + rect[2] / 2
            if cx < mid:
                left_texts.append((rect[1], texts))  # (top, texts)
            else:
                right_texts.append((rect[1], texts))
        # 按 top 排序
        left_texts.sort(key=lambda x: x[0])
        right_texts.sort(key=lambda x: x[0])
        return (
            [t for _, texts in left_texts for t in texts],
            [t for _, texts in right_texts for t in texts],
        )

    def has_large_image(self):
        """是否有占幻灯片面积 > 40% 的大图"""
        for shape in self.image_shapes:
            if shape_area_ratio(shape, self.sw, self.sh) > 0.40:
                return True
        return False

    def has_image_on_side(self):
        """是否有图片在左侧或右侧（占半边以上）"""
        for shape in self.image_shapes:
            l, t, w, h = get_shape_rect(shape)
            cx = l + w / 2
            area_ratio = shape_area_ratio(shape, self.sw, self.sh)
            if area_ratio > 0.15 and (cx > self.sw * 0.55 or cx < self.sw * 0.45):
                return True
        return False

    def get_image_side(self):
        """返回主图片在哪一侧: 'left', 'right', 或 None"""
        best = None
        best_area = 0
        for shape in self.image_shapes:
            l, t, w, h = get_shape_rect(shape)
            area = w * h
            if area > best_area:
                best_area = area
                cx = l + w / 2
                best = "right" if cx > self.sw / 2 else "left"
        return best

    def extract_images(self):
        """提取所有图片，返回路径列表"""
        paths = []
        for i, shape in enumerate(self.image_shapes):
            path = extract_image(shape, self.idx, self.images_dir, i)
            if path:
                paths.append(path)
        return paths

    def detect_chapter_pattern(self):
        """检测章节模式: 第X章 或 Chapter X"""
        all_text = " ".join(t for _, texts, _, _ in self.text_shapes for t in texts)
        # 第X章
        m = re.search(r"第([一二三四五六七八九十\d]+)章", all_text)
        if m:
            return m.group(1)
        # Chapter N
        m = re.search(r"Chapter\s+(\d+)", all_text, re.IGNORECASE)
        if m:
            return m.group(1)
        # PART N
        m = re.search(r"Part\s+(\d+)", all_text, re.IGNORECASE)
        if m:
            return m.group(1)
        return None

    def detect_stats_pattern(self):
        """检测统计数据模式: 大数字 + 标签"""
        stats = []
        # 找大字号的数字文本
        for shape, texts, font_size, rect in self.text_shapes:
            if font_size >= 30:
                for text in texts:
                    # 匹配数字模式: "95%", "1000+", ">80%", "24/7", "3层", "6h+" 等
                    if re.match(r'^[<>≥≤~约]?\d[\d,./]*[%+xX倍万亿层个项台套种h天月年秒分]*[+\-]?$', text.strip()):
                        stats.append(text.strip())
        return stats if len(stats) >= 2 else None

    def detect_numbered_list(self):
        """检测编号列表模式"""
        body = self.get_body_texts()
        if len(body) < 3:
            return False
        numbered = 0
        for text in body:
            if re.match(r'^\s*(\d+[\.\)、]|[①②③④⑤⑥⑦⑧⑨⑩])', text):
                numbered += 1
        return numbered >= 3

    def detect_three_card(self):
        """检测三卡片模式: 三个并排的区域"""
        # 简单判断: 有3个或更多文本区域大致在同一 y 位置但 x 不同
        if len(self.text_shapes) < 6:
            return False
        # 检查是否有3+个背景矩形
        bg_rects = []
        for shape in self.other_shapes:
            try:
                l, t, w, h = get_shape_rect(shape)
                color = get_shape_bg_color(shape)
                if w > 100 and h > 200 and color:
                    bg_rects.append((l, t, w, h))
            except Exception:
                pass
        # 如果有3个类似大小的矩形排列在同一行
        if len(bg_rects) >= 3:
            y_positions = [r[1] for r in bg_rects]
            # 检查是否大致在同一行
            if max(y_positions) - min(y_positions) < 50:
                return True
        return False


# ============================================================
# 幻灯片分类器
# ============================================================

def classify_slide(analyzer):
    """
    根据分析结果判断幻灯片类型。
    返回 (type_str, kwargs_dict)
    """
    idx = analyzer.idx
    total = analyzer.total
    bg_dark = analyzer.bg_dark
    bg_blue = analyzer.bg_blue
    text_shapes = analyzer.text_shapes
    image_shapes = analyzer.image_shapes
    title = analyzer.get_title()
    body = analyzer.get_body_texts()

    # --- 1. 封面页: 第一页，或深色背景+大标题+少量正文 ---
    if idx == 0:
        subtitle = "\n".join(body) if body else ""
        result = {"type": "cover", "title": title, "subtitle": subtitle}
        imgs = analyzer.extract_images()
        if imgs:
            result["image"] = imgs[0]
        return result

    # --- 2. 章节页: 检测 "第X章" 模式 ---
    chapter_num = analyzer.detect_chapter_pattern()
    if chapter_num and (bg_dark or bg_blue):
        # 从正文中移除章节标记
        filtered_body = []
        subtitle = ""
        for text in body:
            if re.search(r"第[一二三四五六七八九十\d]+章", text) or re.search(r"Chapter\s+\d+", text, re.IGNORECASE):
                continue
            filtered_body.append(text)
        if filtered_body:
            subtitle = filtered_body[0] if len(filtered_body) == 1 else "\n".join(filtered_body)
        result = {"type": "chapter", "chapter_num": chapter_num, "title": title}
        if subtitle:
            result["subtitle"] = subtitle
        imgs = analyzer.extract_images()
        if imgs:
            result["image"] = imgs[0]
        return result

    # --- 3. 统计数据页: 多个大数字 ---
    stats_nums = analyzer.detect_stats_pattern()
    if stats_nums:
        # 尝试配对: 大数字 + 紧邻的小文本作为标签
        stats_pairs = _extract_stats_pairs(analyzer)
        if len(stats_pairs) >= 2:
            return {"type": "stats", "title": title, "stats": stats_pairs}

    # --- 4. 全幅图片页: 大图占主导 ---
    if analyzer.has_large_image() and len(text_shapes) <= 3:
        imgs = analyzer.extract_images()
        result = {"type": "image", "title": title}
        if imgs:
            result["image"] = imgs[0]
        caption_texts = [t for t in body if len(t) < 80]
        if caption_texts:
            result["caption"] = caption_texts[0]
        return result

    # --- 5. 图文混排页: 一侧图片 + 一侧文字 ---
    if analyzer.has_image_on_side() and len(body) >= 2:
        imgs = analyzer.extract_images()
        side = analyzer.get_image_side() or "right"
        result = {
            "type": "image_content",
            "title": title,
            "bullets": body,
            "image_side": side,
        }
        if imgs:
            result["image"] = imgs[0]
        return result

    # --- 6. 三卡片页 ---
    if analyzer.detect_three_card():
        cards = _extract_three_cards(analyzer)
        if cards and len(cards) >= 2:
            return {"type": "three_card", "title": title, "cards": cards}

    # --- 7. 双栏页: 文本在左右两侧 ---
    left_texts, right_texts = analyzer.get_text_groups_lr()
    if len(left_texts) >= 2 and len(right_texts) >= 2:
        # 尝试提取左右标题
        left_title = left_texts[0] if left_texts else ""
        left_items = left_texts[1:] if len(left_texts) > 1 else left_texts
        right_title = right_texts[0] if right_texts else ""
        right_items = right_texts[1:] if len(right_texts) > 1 else right_texts
        return {
            "type": "two_column",
            "title": title,
            "left_title": left_title,
            "left_items": left_items,
            "right_title": right_title,
            "right_items": right_items,
        }

    # --- 8. 高亮强调页: 深色背景 + 少量大文本 ---
    if bg_dark and len(text_shapes) <= 4 and len(body) <= 5:
        # 找到最大字号文本作为 main_text
        if text_shapes:
            sorted_by_size = sorted(text_shapes, key=lambda x: x[2], reverse=True)
            main_text = "\n".join(sorted_by_size[0][1])
            label = ""
            sub_text = ""
            remaining = []
            for shape, texts, fs, rect in sorted_by_size[1:]:
                remaining.extend(texts)
            if remaining:
                label = remaining[0]
                sub_text = "\n".join(remaining[1:]) if len(remaining) > 1 else ""
            result = {"type": "highlight", "label": label, "main_text": main_text}
            if sub_text:
                result["sub_text"] = sub_text
            imgs = analyzer.extract_images()
            if imgs:
                result["image"] = imgs[0]
            return result

    # --- 9. 编号列表页 ---
    if analyzer.detect_numbered_list():
        return {"type": "numbered_list", "title": title, "items": body}

    # --- 10. 最后一页特殊处理（可能是致谢页）---
    if idx == total - 1 and len(body) <= 3:
        if bg_dark or bg_blue:
            main_text = title
            sub = "\n".join(body)
            return {"type": "highlight", "label": "", "main_text": main_text, "sub_text": sub}

    # --- 默认: content 页 ---
    result = {"type": "content", "title": title, "bullets": body if body else [""]}
    imgs = analyzer.extract_images()
    if imgs:
        result["image"] = imgs[0]
    return result


def _extract_stats_pairs(analyzer):
    """提取统计数据对: [(数字, 标签), ...]"""
    # 按 x 位置对文本分组
    number_shapes = []
    label_shapes = []

    for shape, texts, font_size, rect in analyzer.text_shapes:
        joined = "\n".join(texts)
        if joined == analyzer.get_title():
            continue
        if font_size >= 30:
            for text in texts:
                if re.match(r'^[<>≥≤~约]?\d[\d,./]*[%+xX倍万亿层个项台套种h天月年秒分]*[+\-]?$', text.strip()):
                    number_shapes.append((rect[0], text.strip()))
        else:
            for text in texts:
                if text.strip():
                    label_shapes.append((rect[0], text.strip()))

    # 按 x 位置配对
    pairs = []
    for nx, num in sorted(number_shapes, key=lambda x: x[0]):
        best_label = ""
        best_dist = float("inf")
        for lx, label in label_shapes:
            dist = abs(lx - nx)
            if dist < best_dist:
                best_dist = dist
                best_label = label
        pairs.append([num, best_label])
        # 移除已配对的 label
        label_shapes = [(lx, l) for lx, l in label_shapes if l != best_label]

    return pairs


def _extract_three_cards(analyzer):
    """提取三卡片内容"""
    # 按 x 位置将非标题文本分成3组
    title_text = analyzer.get_title()
    items = []
    for shape, texts, font_size, rect in analyzer.text_shapes:
        joined = "\n".join(texts)
        if joined == title_text:
            continue
        items.append((rect[0], rect[1], font_size, texts))

    if len(items) < 3:
        return None

    # 按 x 聚类成3组
    items.sort(key=lambda x: x[0])
    x_positions = [item[0] for item in items]
    x_min, x_max = min(x_positions), max(x_positions)
    if x_max - x_min < 100:
        return None

    third = (x_max - x_min) / 3
    groups = [[], [], []]
    for item in items:
        bucket = min(2, int((item[0] - x_min) / third)) if third > 0 else 0
        groups[bucket].append(item)

    cards = []
    for group in groups:
        if not group:
            continue
        # 最大字号或最靠上的作为卡片标题
        group.sort(key=lambda x: (-x[2], x[1]))
        card_title = group[0][3][0] if group[0][3] else ""
        card_body = []
        for item in group[1:]:
            card_body.extend(item[3])
        if not card_body:
            card_body = group[0][3][1:] if len(group[0][3]) > 1 else [""]
        cards.append([card_title, card_body])

    return cards if len(cards) >= 2 else None


# ============================================================
# 主提取逻辑
# ============================================================

def extract_pptx_to_json(pptx_path, images_dir="extracted_images"):
    """
    从 PPTX 文件提取内容结构为 JSON。

    参数:
        pptx_path: .pptx 文件路径
        images_dir: 图片输出目录

    返回:
        dict — 可直接序列化为 JSON
    """
    prs = Presentation(pptx_path)

    # 获取幻灯片尺寸 (pt)
    slide_width_pt = emu_to_pt(prs.slide_width)
    slide_height_pt = emu_to_pt(prs.slide_height)

    total_slides = len(prs.slides)
    slides_data = []

    print(f"Processing: {pptx_path}", file=sys.stderr)
    print(f"Slide dimensions: {slide_width_pt:.0f} x {slide_height_pt:.0f} pt", file=sys.stderr)
    print(f"Total slides: {total_slides}", file=sys.stderr)
    print(f"Images output: {images_dir}", file=sys.stderr)
    print("---", file=sys.stderr)

    for idx, slide_obj in enumerate(prs.slides):
        analyzer = SlideAnalyzer(
            slide_obj, idx, total_slides, images_dir,
            slide_width_pt, slide_height_pt
        )

        slide_data = classify_slide(analyzer)

        print(f"  Slide {idx}: type={slide_data['type']}, "
              f"title={slide_data.get('title', '')[:40]!r}, "
              f"bg={'dark' if analyzer.bg_dark else 'light'}"
              f"{'(blue)' if analyzer.bg_blue else ''}, "
              f"texts={len(analyzer.text_shapes)}, "
              f"images={len(analyzer.image_shapes)}",
              file=sys.stderr)

        slides_data.append(slide_data)

    # 尝试从第一张幻灯片提取演示文稿标题
    pptx_title = ""
    if slides_data:
        pptx_title = slides_data[0].get("title", "")

    result = {
        "title": pptx_title,
        "slides": slides_data,
    }

    return result


# ============================================================
# CLI
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="PPTX → JSON 内容提取器 (兼容 presenton_cli.py 格式)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "示例:\n"
            "  python pptx_to_json.py input.pptx\n"
            "  python pptx_to_json.py input.pptx -o output.json\n"
            "  python pptx_to_json.py input.pptx -o output.json --images-dir extracted_images/\n"
            "  python pptx_to_json.py input.pptx --verbose\n"
        ),
    )
    parser.add_argument("input", help="输入 .pptx 文件路径")
    parser.add_argument("-o", "--output", default=None, help="输出 JSON 文件路径 (默认输出到 stdout)")
    parser.add_argument("--images-dir", default="extracted_images", help="图片输出目录 (默认: extracted_images/)")
    parser.add_argument("--verbose", action="store_true", help="显示详细分析信息")
    parser.add_argument("--pretty", action="store_true", default=True, help="格式化JSON输出 (默认开启)")
    parser.add_argument("--no-pretty", action="store_true", help="关闭格式化JSON输出")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"Error: File not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    if not args.input.lower().endswith((".pptx",)):
        print(f"Warning: File may not be a .pptx file: {args.input}", file=sys.stderr)

    # 将 images_dir 解析为相对于输出文件的路径（如果指定了输出文件）
    images_dir = args.images_dir
    if args.output:
        output_dir = os.path.dirname(os.path.abspath(args.output))
        if not os.path.isabs(images_dir):
            images_dir = os.path.join(output_dir, images_dir)

    result = extract_pptx_to_json(args.input, images_dir)

    indent = None if args.no_pretty else 2
    json_str = json.dumps(result, ensure_ascii=False, indent=indent)

    if args.output:
        os.makedirs(os.path.dirname(os.path.abspath(args.output)), exist_ok=True)
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(json_str)
        print(f"\nDone! Extracted {len(result['slides'])} slides → {args.output}", file=sys.stderr)
    else:
        print(json_str)
        print(f"\nDone! Extracted {len(result['slides'])} slides", file=sys.stderr)


if __name__ == "__main__":
    main()
