#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HydroClaw 50页完整PPT生成器
综合使用所有素材、图片和文档内容
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
from pathlib import Path

# 配色方案
COLORS = {
    'blue': RGBColor(91, 155, 213),      # 主色
    'orange': RGBColor(237, 125, 49),    # 辅色
    'green': RGBColor(112, 173, 71),     # 绿色
    'yellow': RGBColor(255, 192, 0),     # 黄色
    'purple': RGBColor(142, 124, 195),   # 紫色
    'cyan': RGBColor(68, 114, 196),      # 青色
    'red': RGBColor(192, 0, 0),          # 红色
    'dark_blue': RGBColor(31, 78, 120),  # 深蓝
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(127, 127, 127),
}

# 图片资源路径
IMAGE_DIR = Path("./发布会26.4/ppt_assets/team/ppt/media")

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        # 设置宽屏16:9
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

    def add_cover_slide(self):
        """第1页：封面"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局

        # 背景渐变（深蓝到浅蓝）
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['dark_blue']

        # 主标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "HydroClaw认知智能体系"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "从全民养虾到水网自主运行"
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['yellow']
        p.alignment = PP_ALIGN.CENTER

        # 核心概念
        concept_box = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(6), Inches(0.5))
        concept_frame = concept_box.text_frame
        concept_frame.text = "垂直大模型 + 个人水网智能体 + 数字孪生"
        p = concept_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['cyan']
        p.alignment = PP_ALIGN.CENTER

    def add_toc_slide(self):
        """第2页：目录"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "目录"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS['blue']

        # 目录内容（两列布局）
        toc_items = [
            ("一、核心价值与挑战", "3-5"),
            ("二、设计哲学", "6-10"),
            ("三、总体架构", "11-15"),
            ("四、垂直大模型", "16-20"),
            ("五、个人水网智能体", "21-25"),
            ("六、认知决策层", "26-30"),
            ("七、技能编排层", "31-35"),
            ("八、计算引擎层", "36-40"),
            ("九、对象层与元件库", "41-45"),
            ("十、应用案例与展望", "46-50"),
        ]

        # 左列
        left_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(4), Inches(4))
        left_frame = left_box.text_frame
        for i, (title, pages) in enumerate(toc_items[:5]):
            p = left_frame.add_paragraph() if i > 0 else left_frame.paragraphs[0]
            p.text = f"{title}  ......  {pages}"
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['dark_blue']
            p.space_after = Pt(12)

        # 右列
        right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(4))
        right_frame = right_box.text_frame
        for i, (title, pages) in enumerate(toc_items[5:]):
            p = right_frame.add_paragraph() if i > 0 else right_frame.paragraphs[0]
            p.text = f"{title}  ......  {pages}"
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['dark_blue']
            p.space_after = Pt(12)
