#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HydroClaw 发布会PPT V3 - 50页完整版
支持10+种版式，图文混排，专业配色
"""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ─── 配置 ───
IMG_DIR = Path("D:/cowork/ppt/发布会26.4/images")
OUT_FILE = "D:/cowork/ppt/发布会26.4/output/HydroClaw_发布会_V3.pptx"

# 配色
C = {
    'blue': RGBColor(91, 155, 213),
    'dark_blue': RGBColor(31, 78, 120),
    'navy': RGBColor(13, 27, 42),
    'orange': RGBColor(237, 125, 49),
    'green': RGBColor(112, 173, 71),
    'yellow': RGBColor(255, 192, 0),
    'purple': RGBColor(142, 124, 195),
    'cyan': RGBColor(68, 114, 196),
    'red': RGBColor(192, 0, 0),
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(127, 127, 127),
    'light_gray': RGBColor(242, 242, 242),
    'bg_dark': RGBColor(13, 27, 42),
    'bg_card': RGBColor(27, 40, 56),
    'text': RGBColor(51, 51, 51),
    'text_light': RGBColor(127, 127, 127),
}

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class HydroClawPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.page = 0

    # ─── 基础工具 ───

    def _slide(self):
        self.page += 1
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank

    def _bg(self, slide, color=None):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = color or C['navy']

    def _rect(self, slide, left, top, w, h, color, alpha_pct=100):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _text(self, slide, left, top, w, h, text, size=14, bold=False, color=None,
              align=PP_ALIGN.LEFT, font=None, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font or FONT_BODY
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color or C['white']
        p.alignment = align
        try:
            tf.paragraphs[0].space_before = Pt(0)
            tf.paragraphs[0].space_after = Pt(0)
        except:
            pass
        return tb

    def _multitext(self, slide, left, top, w, h, items, size=13, color=None, spacing=6):
        """多行文本，items为列表"""
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.name = FONT_BODY
            p.font.size = Pt(size)
            p.font.color.rgb = color or C['white']
            p.space_after = Pt(spacing)
        return tb

    def _image(self, slide, name, left, top, w=None, h=None):
        path = IMG_DIR / name
        if path.exists():
            if w and h:
                slide.shapes.add_picture(str(path), left, top, width=w, height=h)
            elif w:
                slide.shapes.add_picture(str(path), left, top, width=w)
            else:
                slide.shapes.add_picture(str(path), left, top)
            return True
        return False

    def _page_num(self, slide):
        self._text(slide, SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.4),
                   Inches(0.6), Inches(0.3), str(self.page),
                   size=9, color=C['gray'], align=PP_ALIGN.RIGHT)

    def _section_indicator(self, slide, section_name, color):
        """页面左侧的章节指示条"""
        self._rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, color)
        self._text(slide, Inches(0.15), SLIDE_H - Inches(0.4),
                   Inches(2), Inches(0.3), section_name,
                   size=8, color=C['gray'])

    # ─── 版式类型 ───

    def layout_cover(self, title, subtitle, desc=""):
        """封面页"""
        slide = self._slide()
        self._bg(slide, C['navy'])

        # 装饰线
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), C['blue'])
        self._rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), C['blue'])

        # 左侧装饰
        self._rect(slide, Inches(0.8), Inches(2), Inches(0.06), Inches(3.5), C['blue'])

        self._text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.5),
                   title, size=44, bold=True, color=C['white'])
        self._text(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
                   subtitle, size=24, color=C['yellow'])
        if desc:
            self._text(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.6),
                       desc, size=14, color=C['gray'])

        # 底部信息
        self._text(slide, Inches(1.5), Inches(6.2), Inches(6), Inches(0.4),
                   '中国水利水电科学研究院  |  2026年4月', size=12, color=C['gray'])

        print(f"  [{self.page:02d}] 封面: {title}")

    def layout_toc(self, title, items):
        """目录页"""
        slide = self._slide()
        self._bg(slide, C['navy'])

        self._text(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.8),
                   title, size=32, bold=True, color=C['white'])
        self._rect(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), C['blue'])

        cols = 2
        per_col = (len(items) + 1) // cols
        for i, item in enumerate(items):
            col = i // per_col
            row = i % per_col
            x = Inches(1 + col * 6)
            y = Inches(1.8 + row * 0.7)

            num = f"{i+1:02d}"
            self._text(slide, x, y, Inches(0.6), Inches(0.5),
                       num, size=18, bold=True, color=C['blue'])
            self._text(slide, x + Inches(0.7), y + Inches(0.02), Inches(4.5), Inches(0.5),
                       item, size=14, color=C['white'])

        self._page_num(slide)
        print(f"  [{self.page:02d}] 目录")

    def layout_section(self, num, title, color, subtitle=""):
        """章节标题页"""
        slide = self._slide()
        self._bg(slide, color)

        self._text(slide, Inches(1.5), Inches(2), Inches(3), Inches(1),
                   f"第{num}部分", size=20, color=C['white'])
        self._rect(slide, Inches(1.5), Inches(3), Inches(2), Inches(0.04), C['white'])
        self._text(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(1.5),
                   title, size=40, bold=True, color=C['white'])
        if subtitle:
            self._text(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.6),
                       subtitle, size=16, color=C['white'])

        self._page_num(slide)
        print(f"  [{self.page:02d}] 章节: {title}")

    def layout_title_bullets(self, title, items, color=None, section=""):
        """标题 + 要点列表"""
        slide = self._slide()
        self._bg(slide, C['navy'])
        c = color or C['blue']
        if section:
            self._section_indicator(slide, section, c)

        # 标题栏
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1), c)
        self._text(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.7),
                   title, size=26, bold=True, color=C['white'])

        self._multitext(slide, Inches(1), Inches(1.4), Inches(11), Inches(5.5),
                        items, size=14, spacing=8)

        self._page_num(slide)
        print(f"  [{self.page:02d}] 要点: {title}")

    def layout_image_full(self, title, image_name, color=None, caption=""):
        """全幅图片页"""
        slide = self._slide()
        self._bg(slide, C['navy'])
        c = color or C['blue']

        # 标题
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), c)
        self._text(slide, Inches(0.8), Inches(0.12), Inches(11), Inches(0.65),
                   title, size=24, bold=True, color=C['white'])

        # 图片
        self._image(slide, image_name, Inches(0.3), Inches(1.1),
                    w=Inches(12.7), h=Inches(5.8))

        if caption:
            self._text(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.3),
                       caption, size=9, color=C['gray'])

        self._page_num(slide)
        print(f"  [{self.page:02d}] 全图: {title}")

    def layout_image_left(self, title, image_name, items, color=None, section=""):
        """左图右文"""
        slide = self._slide()
        self._bg(slide, C['navy'])
        c = color or C['blue']
        if section:
            self._section_indicator(slide, section, c)

        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), c)
        self._text(slide, Inches(0.8), Inches(0.12), Inches(11), Inches(0.65),
                   title, size=24, bold=True, color=C['white'])

        self._image(slide, image_name, Inches(0.3), Inches(1.1),
                    w=Inches(6.3), h=Inches(5.8))

        self._multitext(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(5.5),
                        items, size=13, spacing=7)

        self._page_num(slide)
        print(f"  [{self.page:02d}] 左图: {title}")

    def layout_cards(self, title, cards, color=None, section=""):
        """卡片式布局 (2-4张卡片)"""
        slide = self._slide()
        self._bg(slide, C['navy'])
        c = color or C['blue']
        if section:
            self._section_indicator(slide, section, c)

        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), c)
        self._text(slide, Inches(0.8), Inches(0.12), Inches(11), Inches(0.65),
                   title, size=24, bold=True, color=C['white'])

        n = len(cards)
        card_w = (12 - 0.3 * (n - 1)) / n
        for i, card in enumerate(cards):
            x = Inches(0.5 + i * (card_w + 0.3))
            y = Inches(1.2)
            w = Inches(card_w)
            h = Inches(5.8)

            card_color = card.get('color', c)
            self._rect(slide, x, y, w, h, C['bg_card'])

            # 顶部色条
            self._rect(slide, x, y, w, Inches(0.06), card_color)

            # 卡片标题
            self._text(slide, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.5),
                       card['title'], size=16, bold=True, color=card_color)

            # 卡片内容
            items = card.get('items', [])
            self._multitext(slide, x + Inches(0.3), y + Inches(0.9),
                           w - Inches(0.6), h - Inches(1.2),
                           items, size=11, spacing=5)

        self._page_num(slide)
        print(f"  [{self.page:02d}] 卡片: {title}")

    def layout_comparison(self, title, left_title, left_items, right_title, right_items,
                          left_color=None, right_color=None, section=""):
        """左右对比页"""
        slide = self._slide()
        self._bg(slide, C['navy'])
        lc = left_color or C['gray']
        rc = right_color or C['blue']
        if section:
            self._section_indicator(slide, section, rc)

        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), C['dark_blue'])
        self._text(slide, Inches(0.8), Inches(0.12), Inches(11), Inches(0.65),
                   title, size=24, bold=True, color=C['white'])

        # VS标记
        self._text(slide, Inches(6.1), Inches(3.5), Inches(1.1), Inches(0.8),
                   'VS', size=28, bold=True, color=C['yellow'], align=PP_ALIGN.CENTER)

        for side, (st, items, color) in enumerate([
            (left_title, left_items, lc), (right_title, right_items, rc)
        ]):
            x = Inches(0.5 + side * 6.7)
            w = Inches(5.8)

            self._rect(slide, x, Inches(1.2), w, Inches(5.8), C['bg_card'])
            self._rect(slide, x, Inches(1.2), w, Inches(0.06), color)
            self._text(slide, x + Inches(0.3), Inches(1.5), w - Inches(0.6), Inches(0.5),
                       st, size=18, bold=True, color=color)
            self._multitext(slide, x + Inches(0.3), Inches(2.2),
                           w - Inches(0.6), Inches(4.5),
                           items, size=12, spacing=6)

        self._page_num(slide)
        print(f"  [{self.page:02d}] 对比: {title}")

    def layout_big_number(self, title, numbers, color=None, section=""):
        """大数字展示页"""
        slide = self._slide()
        self._bg(slide, C['navy'])
        c = color or C['blue']
        if section:
            self._section_indicator(slide, section, c)

        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), c)
        self._text(slide, Inches(0.8), Inches(0.12), Inches(11), Inches(0.65),
                   title, size=24, bold=True, color=C['white'])

        n = len(numbers)
        for i, (num, label, desc) in enumerate(numbers):
            x = Inches(0.5 + i * (12 / n))
            w = Inches(12 / n - 0.3)

            self._text(slide, x, Inches(2), w, Inches(2),
                       num, size=64, bold=True, color=c, align=PP_ALIGN.CENTER)
            self._text(slide, x, Inches(4.2), w, Inches(0.5),
                       label, size=18, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
            self._text(slide, x, Inches(4.8), w, Inches(0.8),
                       desc, size=12, color=C['gray'], align=PP_ALIGN.CENTER)

        self._page_num(slide)
        print(f"  [{self.page:02d}] 大数: {title}")

    def layout_code_demo(self, title, code_lines, desc_items, color=None, section=""):
        """代码演示页（左代码右说明）"""
        slide = self._slide()
        self._bg(slide, C['navy'])
        c = color or C['blue']
        if section:
            self._section_indicator(slide, section, c)

        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), c)
        self._text(slide, Inches(0.8), Inches(0.12), Inches(11), Inches(0.65),
                   title, size=24, bold=True, color=C['white'])

        # 代码区域
        self._rect(slide, Inches(0.3), Inches(1.1), Inches(7.5), Inches(5.8),
                   RGBColor(20, 20, 30))

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(7), Inches(5.4))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(code_lines):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = line
            p.font.name = FONT_MONO
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(180, 220, 180)
            p.space_after = Pt(2)

        # 说明
        self._multitext(slide, Inches(8.2), Inches(1.3), Inches(4.8), Inches(5.5),
                        desc_items, size=12, spacing=6)

        self._page_num(slide)
        print(f"  [{self.page:02d}] 代码: {title}")

    def layout_summary(self, title, highlights, slogan=""):
        """总结页"""
        slide = self._slide()
        self._bg(slide, C['navy'])

        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), C['blue'])

        self._text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(1),
                   title, size=36, bold=True, color=C['white'])
        self._rect(slide, Inches(0.8), Inches(1.7), Inches(2), Inches(0.04), C['blue'])

        for i, (icon_text, highlight) in enumerate(highlights):
            y = Inches(2.2 + i * 0.8)
            self._text(slide, Inches(1), y, Inches(0.5), Inches(0.5),
                       icon_text, size=18, color=C['blue'])
            self._text(slide, Inches(1.7), y + Inches(0.05), Inches(10), Inches(0.5),
                       highlight, size=16, color=C['white'])

        if slogan:
            self._rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2), C['dark_blue'])
            self._text(slide, Inches(0.8), Inches(6), Inches(11.7), Inches(0.8),
                       slogan, size=24, bold=True, color=C['yellow'], align=PP_ALIGN.CENTER)

        self._page_num(slide)
        print(f"  [{self.page:02d}] 总结: {title}")

    # ─── 生成全部50页 ───

    def generate(self):
        print("=" * 50)
        print("HydroClaw 发布会PPT V3 - 50页完整版")
        print("=" * 50)

        # ═══ 第1页：封面 ═══
        self.layout_cover(
            "HydroClaw 认知智能体系",
            "从全民养虾到水网自主运行",
            "垂直大模型 + 个人水网智能体 + 数字孪生"
        )

        # ═══ 第2页：目录 ═══
        self.layout_toc("目录", [
            "一、核心价值与挑战",
            "二、设计哲学",
            "三、总体架构",
            "四、垂直大模型",
            "五、个人水网智能体",
            "六、认知决策层",
            "七、技能编排层",
            "八、计算引擎层",
            "九、对象层与元件库",
            "十、应用案例与展望",
        ])

        # ═══ 第一部分：核心价值与挑战（3-5页）═══
        self.layout_section("一", "核心价值与挑战", C['blue'],
                           "水利行业数字化转型的根本性挑战")

        self.layout_comparison(
            "水利行业的「全民养虾」困境",
            "当前困境", [
                ">> 每个水利项目都在「养虾」—— 养一套自己的软件系统",
                ">> 软件碎片化：同一个水网，10个系统各管一段",
                ">> 数据孤岛：监测数据、模型数据、调度数据互不相通",
                ">> 重复造轮子：每个设计院都在写相似的计算程序",
                ">> 依赖个人经验：资深工程师退休=知识流失",
                ">> 无法自主运行：所有决策都需要人工干预",
            ],
            "根本原因", [
                ">> 缺乏统一的认知智能框架",
                ">> 通用大模型不懂水利专业知识",
                ">> 没有把「规则」和「AI」有机结合",
                ">> 没有标准化的元件库和引擎体系",
                ">> 没有「选Skill即获角色」的设计理念",
                ">> 各环节无法形成认知闭环",
            ],
            left_color=C['red'], right_color=C['orange'],
            section="一、核心价值与挑战"
        )

        self.layout_big_number(
            "HydroClaw 解决方案",
            [
                ("10×", "效率提升", "设计全流程自动化"),
                ("1", "统一框架", "认知-技能-引擎-对象"),
                ("∞", "场景覆盖", "设计/运维/科研/教学"),
            ],
            color=C['blue'],
            section="一、核心价值与挑战"
        )

        # ═══ 第二部分：设计哲学（6-10页）═══
        self.layout_section("二", "设计哲学", C['orange'],
                           "四句话概括整个方案")

        self.layout_cards(
            "四大核心设计原则",
            [
                {'title': '原则1', 'color': C['blue'], 'items': [
                    '确定性的事用规则',
                    '不确定性的事用大模型',
                    '大模型的输出用规则兜底',
                    '',
                    '>> 规则处理60%交互',
                    '>> 大模型处理30%',
                    '>> 混合处理10%',
                ]},
                {'title': '原则2：Skill即角色', 'color': C['orange'], 'items': [
                    '用户选了什么Skill',
                    '就具备什么角色',
                    '行为即身份',
                    '',
                    '>> 无需声明角色',
                    '>> 无需切换模式',
                    '>> 自动加载规则',
                ]},
                {'title': '原则3：通用引擎', 'color': C['green'], 'items': [
                    '引擎通用不分领域',
                    '差异全在元件和参数',
                    '',
                    '>> 七大引擎与领域无关',
                    '>> 纯算法，不知道算什么水网',
                    '>> 新场景=新元件，无需改引擎',
                ]},
                {'title': '原则4：统一基类', 'color': C['purple'], 'items': [
                    '物理元件、水文元件',
                    '模型元件统一基类',
                    '混合组装任意水网',
                    '',
                    '>> 统一compute()接口',
                    '>> 像搭积木一样组装',
                ]},
            ],
            color=C['orange'],
            section="二、设计哲学"
        )

        self.layout_title_bullets(
            "原则1详解：大模型与规则协同",
            [
                "规则引擎处理（~60%交互）—— 确定性场景：",
                "  >> 查询闸门当前水位、执行标准巡检流程、ODD阈值判别、预警等级确定",
                "",
                "大模型处理（~30%交互）—— 需要理解和推理：",
                "  >> 分析工况并推荐方案、设计方案比选、跨域关联分析、可解释性生成",
                "",
                "大模型主导+规则兜底（~10%）—— 开放性问题：",
                "  >> 未预见的异常工况、科研探索性讨论、应急情况类比推理、创新方案构想",
                "",
                "核心机制：大模型生成方案 -> 规则引擎校核 -> 反馈修正",
                "安全底线：规则只紧不松、不编造数据、不确定性必须声明",
            ],
            color=C['orange'],
            section="二、设计哲学"
        )

        self.layout_title_bullets(
            "原则2详解：Skill即角色 —— 行为即身份",
            [
                "传统做法：用户先选角色 -> 再选功能 -> 规则静态绑定",
                "HydroClaw：用户直接选Skill -> 自动获得角色 -> 规则动态加载",
                "",
                "示例：同一个用户无缝切换：",
                "  上午用「仿真对比」Skill -> 自动进入科研模式 -> 深度分析输出",
                "  下午用「ODD预警」Skill -> 自动进入运维模式 -> 简洁步骤输出",
                "",
                "同一个B2预警模板，不同Skill上下文自动调整：",
                "  运维Skill触发 -> 简洁步骤式",
                "  科研Skill触发 -> 深度分析 + 论文引用",
                "",
                "技能继承体系：原子Skill -> 组合Skill -> 流程Skill",
                "每个Skill自带 (role, level) 二元标签",
            ],
            color=C['orange'],
            section="二、设计哲学"
        )

        self.layout_title_bullets(
            "原则3&4：通用引擎 + 统一基类",
            [
                "通用引擎 —— 与领域无关的纯算法：",
                "  >> 仿真引擎不知道自己在算城市供水还是河道防洪",
                "  >> 优化引擎不关心优化的是泵站启停还是闸门序列",
                "  >> 新增场景 = 新增元件 + 参数，引擎代码不变",
                "",
                "统一基类 —— 三族元件混合组装：",
                "  >> 物理元件(闸泵阀库河渠) + 水文元件(降雨产汇流)",
                "  >> + 模型元件(LSTM/GNN/RL策略)",
                "  >> 统一 compute(state, inputs, dt) -> dict 接口",
                "  >> 同一个水网里可以同时有物理元件和AI模型元件",
                "",
                "这意味着：物理仿真和AI预测可以在同一个管线中混合使用",
            ],
            color=C['orange'],
            section="二、设计哲学"
        )

        # ═══ 第三部分：总体架构（11-15页）═══
        self.layout_section("三", "总体架构", C['green'],
                           "四层设计，单向调用，各层只依赖下一层统一接口")

        self.layout_image_full(
            "四层智能决策体系",
            "chart_architecture_pyramid.png",
            color=C['green'],
            caption="四层之间的调用关系是单向的：①调②，②调③，③驱动④"
        )

        self.layout_cards(
            "四层架构详解",
            [
                {'title': '① 认知决策层', 'color': C['blue'], 'items': [
                    '大模型 + 规则引擎 + 模板渲染',
                    '理解意图、选择Skill',
                    '约束行为、格式化输出',
                    '五层规则继承',
                    '24个回复模板',
                ]},
                {'title': '② 技能编排层', 'color': C['orange'], 'items': [
                    'Skill体系：原子->组合->流程',
                    '知道做什么、按什么顺序',
                    '编排逻辑是领域知识',
                    '~25个原子 + ~10组合 + ~6流程',
                ]},
                {'title': '③ 计算引擎层', 'color': C['green'], 'items': [
                    '七大通用引擎 + 工具箱',
                    '怎么算 -- 纯算法，不分领域',
                    '仿真/辨识/调度/控制',
                    '优化/预测/学习 + 100+工具',
                ]},
                {'title': '④ 对象层', 'color': C['yellow'], 'items': [
                    '三族元件库',
                    '物理/水文/模型 统一基类',
                    '水网组装 + 外部工具接入',
                    '算什么 -- 具体元件实例',
                ]},
            ],
            color=C['green'],
            section="三、总体架构"
        )

        self.layout_title_bullets(
            "数据流向与调用关系",
            [
                "用户输入 -> 认知决策层：",
                "  >> 大模型解析意图，规则引擎匹配确定性场景",
                "  >> 确定调用哪个Skill，自动加载L0-L4规则",
                "",
                "认知层 -> 技能编排层：",
                "  >> Skill按预定义的编排逻辑，依次调用引擎",
                "  >> 编排逻辑是硬编码的领域知识，不由大模型即兴决定",
                "",
                "技能层 -> 计算引擎层：",
                "  >> 引擎提供统一接口：simulate / predict / optimize / learn",
                "  >> 引擎内部可切换后端：自研HydroOS / EPANET / HEC-RAS",
                "",
                "引擎层 <-> 对象层：",
                "  >> 引擎读取元件参数，驱动元件计算，更新元件状态",
            ],
            color=C['green'],
            section="三、总体架构"
        )

        self.layout_title_bullets(
            "架构优势分析",
            [
                "可扩展性：新增引擎或技能不影响其他层",
                "  >> 接入WNTR管网后端 = 1个MCP工具 + 1个S0 Skill包装，上层S1/S2不变",
                "",
                "可维护性：各层职责清晰，代码内聚",
                "  >> 认知层只管理解和约束，不做计算",
                "  >> 引擎层只管算法，不知道业务语境",
                "",
                "可复用性：引擎和技能可被多场景复用",
                "  >> 同一个优化引擎，城市供水和河道防洪都可调用",
                "",
                "可测试性：各层可独立测试",
                "  >> SIL平台天然就是RL训练场和验证环境",
                "",
                "安全性：L0安全底线不可覆盖",
            ],
            color=C['green'],
            section="三、总体架构"
        )

        # ═══ 第四部分：垂直大模型（16-20页）═══
        self.layout_section("四", "垂直大模型", C['purple'],
                           "通用大模型 vs 水利垂直大模型")

        self.layout_comparison(
            "通用大模型 vs 垂直大模型",
            "通用大模型", [
                ">> 广泛知识但缺乏水利专业深度",
                ">> 不了解水利规范和标准",
                ">> 可能给出物理上不可能的方案",
                ">> 没有工程安全意识",
                ">> 无法进行专业计算",
                ">> 幻觉风险高",
            ],
            "HydroClaw垂直大模型", [
                ">> 水利规范知识库（RAG增强）",
                ">> 工程案例库（历史最佳实践）",
                ">> 物理约束体系（流速/压力/管径等）",
                ">> 规则引擎兜底（大模型失效时接管）",
                ">> 五层规则继承（L0安全底线不可覆盖）",
                ">> 不确定性必须声明",
            ],
            left_color=C['gray'], right_color=C['purple'],
            section="四、垂直大模型"
        )

        self.layout_title_bullets(
            "水利规范知识库 + 工程案例库",
            [
                "水利规范知识库（RAG增强检索）：",
                "  >> 设计规范：管径选择、流速限制、安全系数",
                "  >> 运行规程：调度规则、应急预案、巡检标准",
                "  >> 施工标准：工艺要求、质量控制、验收标准",
                "",
                "工程案例库（历史最佳实践）：",
                "  >> 中线工程：冰期 Fr<=0.06，水位变化率<=0.3m/h",
                "  >> 胶东调水：泵站群启动先下游后上游",
                "  >> 双容水箱：溢流口高度、映射知识点",
                "",
                "物理约束体系：",
                "  >> 质量守恒、能量守恒不可违反",
                "  >> 管道流速、水泵工作点必须在合理范围",
                "  >> 安全约束只能收紧、不能放松",
            ],
            color=C['purple'],
            section="四、垂直大模型"
        )

        self.layout_title_bullets(
            "规则引擎兜底机制",
            [
                "为什么需要兜底？",
                "  >> 大模型可能产生幻觉 -- 规则引擎校验物理可行性",
                "  >> 大模型可能遗漏约束 -- 规则引擎补充必要校核",
                "  >> 大模型可能过于保守或激进 -- 规则引擎约束到合理范围",
                "",
                "兜底机制三层防线：",
                "  1. 实时校核：大模型每次输出，规则引擎同步验证",
                "  2. 物理预演：涉及执行的方案必须通过仿真验证",
                "  3. 人工确认：涉及实际操作的指令需二次确认",
                "",
                "L0安全底线（不可覆盖）：",
                "  >> 物理守恒定律不可违反",
                "  >> 安全约束只紧不松",
                "  >> 不编造数据",
                "  >> 不确定性必须声明",
            ],
            color=C['purple'],
            section="四、垂直大模型"
        )

        self.layout_image_full(
            "传统方法 vs HydroClaw 能力对比",
            "chart_comparison_radar.png",
            color=C['purple']
        )

        # ═══ 第五部分：个人水网智能体（21-25页）═══
        self.layout_section("五", "个人水网智能体", C['cyan'],
                           "选Skill即获角色，行为即身份")

        self.layout_image_full(
            "四大角色 -- 选Skill即获角色",
            "chart_persona_matrix.png",
            color=C['cyan']
        )

        self.layout_cards(
            "设计工程师场景 -- 个人超级工具",
            [
                {'title': '泵站选型', 'color': C['blue'], 'items': [
                    '输入：设计流量、扬程',
                    '过程：性能曲线拟合、变频/工频决策',
                    '输出：选型报告 + 全工况曲线',
                    '',
                    'L0: 出选型计算书',
                    'L1: 泵型比选决策',
                    'L2: 跨专业协调审批',
                ]},
                {'title': '断面优化', 'color': C['orange'], 'items': [
                    '输入：水力条件、地形约束',
                    '过程：多目标优化搜索',
                    '输出：最优断面 + 对比分析',
                    '',
                    '替代传统Excel试算',
                    '自动满足设计规范',
                ]},
                {'title': '管网计算', 'color': C['green'], 'items': [
                    '输入：管网拓扑、需水量',
                    '过程：WNTR/EPANET仿真',
                    '输出：压力分布 + 流量平衡',
                    '',
                    '可接入已有INP模型',
                    '泄漏检测与定位',
                ]},
            ],
            color=C['cyan'],
            section="五、个人水网智能体"
        )

        self.layout_cards(
            "运维 / 科研 / 教学场景",
            [
                {'title': '运维人员', 'color': C['orange'], 'items': [
                    '智能运维助手',
                    'ODD预警：水位超限自动上报',
                    '故障诊断：多因素根因分析',
                    '调度优化：全局最优分配',
                    '应急响应：保守策略优先',
                    '',
                    '三级权限：现地/管理段/调度中心',
                ]},
                {'title': '科研人员', 'color': C['green'], 'items': [
                    '科研实验平台',
                    '仿真对比：多算法并行验证',
                    '参数辨识：UKF/EnKF自动校准',
                    'RL策略：PPO/SAC自动训练',
                    '论文级深度输出',
                    '',
                    '三级权限：硕士/博士副研/PI博导',
                ]},
                {'title': '学生教师', 'color': C['purple'], 'items': [
                    '智能教学助手',
                    '分层引导：S0科普 -> S3研究生',
                    '苏格拉底式提问引导',
                    '大专生：步骤式操作',
                    '本科生：参数调整实验',
                    '研究生：自主设计实验',
                    '',
                    '教师端：场景配置 + 自动评估',
                ]},
            ],
            color=C['cyan'],
            section="五、个人水网智能体"
        )

        self.layout_title_bullets(
            "层级决定深度，而非能否使用",
            [
                "同一个Skill，不同层级的行为完全不同：",
                "",
                "以「泵站选型」为例：",
                "  designer_L0（设计员）：按参数查手册，出计算书，需L1校核",
                "  designer_L1（专业负责人）：泵型比选，变频决策，修改参数",
                "  designer_L2（总工）：跨专业协调，决定标准，批准非标方案",
                "",
                "以「水位控制实验」为例：",
                "  student_S0（低年级）：只观看动画演示，类比式解释",
                "  student_S1（大专生）：操作启停，不能改参数",
                "  student_S2（本科生）：调PID参数，苏格拉底引导",
                "  student_S3（研究生）：自行设计MPC，开放探索",
                "",
                "用户层级怎么确定？Skill菜单本身就是权限边界",
            ],
            color=C['cyan'],
            section="五、个人水网智能体"
        )

        # ═══ 第六部分：认知决策层（26-30页）═══
        self.layout_section("六", "认知决策层", C['blue'],
                           "大模型 + 规则引擎 + 模板渲染")

        self.layout_image_full(
            "认知决策五步流程",
            "chart_cognitive_flow.png",
            color=C['blue'],
            caption="每次用户交互，认知层执行固定五步：意图理解 -> Skill匹配 -> 规则加载 -> Skill执行 -> 模板渲染"
        )

        self.layout_image_full(
            "五层规则继承体系",
            "chart_rule_layers.png",
            color=C['blue'],
            caption="规则从Skill选择自动加载，用户不需要声明角色、不需要切换模式"
        )

        self.layout_title_bullets(
            "24个回复模板 -- 继承组合避免模板爆炸",
            [
                "基于CHS六要素（Plant-Actuator-Sensor-Disturbance-Constraint-Objective）：",
                "",
                "T0 六要素基类：T0-S状态查询 / T0-A分析诊断 / T0-D决策建议（3个原子模板）",
                "",
                "T1-T3 对象子模板（7个）：",
                "  T1 执行器类：status / diagnosis / control",
                "  T2 蓄水体类：status / prediction",
                "  T3 输水体类：status / simulation",
                "",
                "B1-B6 业务模板（6个）：预报/预警/预演/预案/SIL/HIL",
                "",
                "W1-W6 工作流模板（6个）：四预闭环/SIL验证/HIL验证/MBD全流程/应急/教学",
                "",
                "角色适配：不增加模板，用修饰器调整输出深度和术语",
                "案例适配：不增加模板，用L4参数填充具体数值",
            ],
            color=C['blue'],
            section="六、认知决策层"
        )

        self.layout_code_demo(
            "同一模板，不同Skill自动调整输出",
            [
                "# B2预警模板 - 运维Skill触发",
                ">>> skill = 'ODD预警'",
                ">>> context = 'operator_L1'",
                "",
                "输出：",
                "  ⚠ 渠池2水位42.35m超限",
                "  立即上报L1",
                "  建议：开启3号闸门泄流",
                "",
                "# B2预警模板 - 科研Skill触发",
                ">>> skill = '仿真对比'",
                ">>> context = 'researcher_L1'",
                "",
                "输出：",
                "  Saint-Venant分析显示上游闸门",
                "  调整后波前传播时间约45min",
                "  当前状态为积分饱和特征",
                "  参见Malaterre(2004)...",
            ],
            [
                "模板机制要点：",
                "",
                ">> 同一个B2预警模板",
                ">> 不同Skill上下文自动切换",
                "",
                "运维Skill触发：",
                "  - 简洁步骤式输出",
                "  - 直接给出行动建议",
                "  - 语言简练，重点突出",
                "",
                "科研Skill触发：",
                "  - 深度分析输出",
                "  - 包含数学推导",
                "  - 引用学术文献",
                "  - 详细物理解释",
                "",
                ">> 这就是「Skill即角色」的威力",
            ],
            color=C['blue'],
            section="六、认知决策层"
        )

        # ═══ 第七部分：技能编排层（31-35页）═══
        self.layout_section("七", "技能编排层", C['orange'],
                           "Skill体系：原子 -> 组合 -> 流程")

        self.layout_image_full(
            "Skill三层继承体系",
            "chart_skill_hierarchy.png",
            color=C['orange'],
            caption="S0原子Skill封装单个MCP工具调用 | S1组合Skill编排多个S0 | S2流程Skill编排多个S1"
        )

        self.layout_cards(
            "原子Skill清单（按引擎分类）",
            [
                {'title': '仿真/辨识/调度', 'color': C['blue'], 'items': [
                    '仿真类 (6个):',
                    '  明渠/管网/水库/水箱',
                    '  水锤/冰期仿真',
                    '',
                    '辨识类 (5个):',
                    '  糙率/泵曲线/闸系数',
                    '  UKF状态估计/模型校准',
                    '',
                    '调度类 (5个):',
                    '  全局优化/泵站经济运行',
                    '  闸门序列/水平衡/水库',
                ]},
                {'title': '控制/检测/预测', 'color': C['orange'], 'items': [
                    '控制类 (4个):',
                    '  PID/MPC/DMPC/控制器评估',
                    '',
                    '检测评估类 (5个):',
                    '  ODD检查/泄漏检测',
                    '  异常检测/WNAL评估/蒸发',
                    '',
                    '预测类 (5个):',
                    '  概念水文模型预报',
                    '  LSTM/Transformer时序',
                    '  GNN图网络/降阶模型/混合',
                ]},
                {'title': '学习/工具', 'color': C['green'], 'items': [
                    '学习类 (5个):',
                    '  RL策略训练/部署',
                    '  模仿学习/在线学习',
                    '  模型训练管理',
                    '',
                    '工具类 (18个):',
                    '  数据清洗/时序重采样',
                    '  曲线拟合/频率分析',
                    '  模型诊断/报告生成',
                    '  单位换算/公式计算器',
                ]},
            ],
            color=C['orange'],
            section="七、技能编排层"
        )

        self.layout_title_bullets(
            "组合Skill与流程Skill",
            [
                "组合Skill（~15个）—— 编排多个S0完成业务功能：",
                "  S1-FCST预报 = UKF状态估计 -> 仿真/数据预测 -> 不确定性量化",
                "  S1-WARN预警 = ODD检查 -> 趋势预测 -> 预警等级判定",
                "  S1-RHSL预演 = 方案解析 -> 闭环仿真 -> 约束校核 -> 方案排序",
                "  S1-SIL测试 = 工况加载 -> 闭环仿真xN -> 性能评估 -> 准入判定",
                "  S1-TRAIN训练 = 数据准备 -> 训练 -> 验证 -> 版本管理 -> 部署",
                "",
                "流程Skill（~6个）—— 编排多个S1完成完整流程：",
                "  S2-4P四预闭环 = 预报 -> 预警 -> 预演 -> 预案",
                "  S2-MBD全流程 = ODD定义 -> 控制设计 -> 预演 -> SIL -> HIL",
                "  S2-EMG应急响应 = 预警 -> 预案 -> 预演 -> 执行/人工确认",
                "  S2-EDU教学实验 = 场景配置 -> SIL(教学模式) -> 自动评估",
                "",
                "编排逻辑是硬编码的领域知识，不由大模型即兴决定",
            ],
            color=C['orange'],
            section="七、技能编排层"
        )

        self.layout_title_bullets(
            "Skill承载角色：四维度 x 多层级",
            [
                "运维维度（三级）：",
                "  L0 现地（执行指令/数据采集） | L1 管理段（段内协调） | L2 调度中心（全局决策）",
                "",
                "设计维度（三级）：",
                "  L0 设计员（参数计算/出图） | L1 专业负责人（设计/校核） | L2 总工（审定/跨专业）",
                "",
                "科研维度（三级）：",
                "  L0 硕士生（辅助计算/文献） | L1 博士生/副研（独立研究） | L2 PI/博导（方向把控）",
                "",
                "教学维度（教师三级 + 学生四级）：",
                "  教师：L0助教 / L1主讲 / L2课程负责人",
                "  学生：S0科普 / S1大专 / S2本科 / S3研究生",
                "",
                "层级决定的不是「能不能用」，而是「用到什么程度」",
            ],
            color=C['orange'],
            section="七、技能编排层"
        )

        # ═══ 第八部分：计算引擎层（36-40页）═══
        self.layout_section("八", "计算引擎层", C['green'],
                           "七大通用引擎 + 100+工具函数")

        self.layout_image_full(
            "七大通用计算引擎",
            "chart_seven_engines.png",
            color=C['green'],
            caption="所有引擎提供统一接口，Skill层不关心内部用哪个求解器"
        )

        self.layout_cards(
            "物理机理引擎 vs 数据学习引擎",
            [
                {'title': '物理机理引擎（5个）', 'color': C['blue'], 'items': [
                    '仿真引擎: Saint-Venant / MOC / ODE',
                    '辨识引擎: 最小二乘 / UKF / PINN',
                    '调度引擎: LP / MILP / DP / SDDP',
                    '控制引擎: PID / MPC / DMPC',
                    '优化引擎: SQP / GA / DE / PSO',
                    '',
                    '特点：基于物理方程',
                    '适用：有明确数学模型的场景',
                ]},
                {'title': '数据学习引擎（2个）', 'color': C['purple'], 'items': [
                    '预测引擎:',
                    '  概念水文 / LSTM / Transformer',
                    '  GNN / 降阶替代 / 物理-数据混合',
                    '',
                    '学习引擎:',
                    '  RL(PPO/SAC/DQN)',
                    '  模仿学习 / 在线学习',
                    '',
                    '特点：从数据中学习',
                    '适用：无解析方程或需要快速预测',
                ]},
            ],
            color=C['green'],
            section="八、计算引擎层"
        )

        self.layout_code_demo(
            "统一接口 -- Skill不关心后端实现",
            [
                "# 统一接口示例",
                "engine.simulate(network, bc, dt)",
                "  -> SimulationResult",
                "",
                "engine.predict(features, horizon)",
                "  -> PredictionResult",
                "",
                "engine.learn(env, reward_fn, episodes)",
                "  -> PolicyResult",
                "",
                "# 同一个预报Skill，可选不同后端：",
                "S1-FCST(预报):",
                "  方式A: SimulationEngine + SV方程",
                "         (精确但慢)",
                "  方式B: PredictionEngine + LSTM",
                "         (快但需训练数据)",
                "  方式C: PredictionEngine + IDZ降阶",
                "         (快且保留物理约束)",
                "  方式D: 混合模型(物理+数据融合)",
            ],
            [
                "引擎内部可切换后端：",
                "",
                "engine.simulate(city_network)",
                "  |-- 自研HydroOS求解器（默认）",
                "  |-- EPANET/WNTR后端",
                "  |-- HEC-RAS后端",
                "  |-- PSS/E后端（水电）",
                "",
                "自研体系是核心：",
                "  >> 什么水网都能从零建模",
                "  >> 不依赖任何外部工具",
                "",
                "外部工具是生态：",
                "  >> 各单位已有的成熟模型",
                "  >> 标准接口接入",
                "  >> 科研可同时对比验证",
            ],
            color=C['green'],
            section="八、计算引擎层"
        )

        self.layout_title_bullets(
            "外部工具生态 -- 第一梯队已就绪",
            [
                "第一梯队（pip install即可，已就绪）：",
                "  WNTR -- 城市供水管网仿真（EPA + Sandia国家实验室）",
                "  PySWMM -- 城市排水/雨洪管理（运行时控制注入）",
                "  ras-commander -- 河道防洪/溃坝（HEC-RAS Python自动化）",
                "  FloPy -- 地下水（USGS MODFLOW接口）",
                "  Pyomo -- 通用优化求解（已集成）",
                "",
                "第二梯队（需宿主软件）：HEC-HMS / MIKE SDK / OpenFOAM",
                "",
                "第四梯队（ML基础设施）：PyTorch / Stable-Baselines3 / scikit-learn",
                "  >> 这些是引擎的内部依赖，不直接暴露为MCP",
                "",
                "每个工具封装 = 1个MCP函数 + 1个S0 Skill",
                "上层S1/S2/规则/模板完全不变",
            ],
            color=C['green'],
            section="八、计算引擎层"
        )

        self.layout_title_bullets(
            "工具箱 -- 10大工具包 x 100+工具函数",
            [
                "TK1 数据处理：清洗/重采样/滤波/异常检测/物理校验",
                "TK2 曲线拟合：泵曲线/库容曲线/水位-流量关系/入渗曲线",
                "TK3 统计分析：频率分析/趋势检验/相关分析/极值统计",
                "TK4 诊断评估：传感器/水泵/闸门/管道/模型/控制器诊断",
                "TK5 单位坐标：单位换算/高程基准/坐标系转换",
                "TK6 格式I/O：SCADA/INP/HDF5/NetCDF/Shapefile",
                "TK7 水力计算：Manning/堰流/孔口/管损/水锤/临界水深/壅水",
                "TK8 水质计算：BOD衰减/余氯/混合/WQI/泥沙",
                "TK9 经济能耗：泵站能耗/发电量/生命周期成本/碳排放",
                "TK10 可视化：时序曲线/剖面图/管网拓扑/热力图/仿真动画",
                "",
                "工具箱不是第八个引擎，是跨引擎的公共基础设施",
                "Skill可直接调用，也被引擎内部调用",
            ],
            color=C['green'],
            section="八、计算引擎层"
        )

        # ═══ 第九部分：对象层（41-45页）═══
        self.layout_section("九", "对象层与元件库", C['yellow'],
                           "三族元件 -- 统一基类，混合组装任意水网")

        self.layout_image_full(
            "三族元件库 -- 统一基类架构",
            "chart_component_tree.png",
            color=C['yellow'],
            caption="所有元件继承同一个基类 Component，都有 compute() 方法、参数集、端口和约束"
        )

        self.layout_cards(
            "第一族：物理元件（被仿真引擎驱动）",
            [
                {'title': 'T1 执行器类', 'color': C['blue'], 'items': [
                    '闸门：平板/弧形/人字',
                    '水泵：离心/轴流/混流',
                    '阀门：蝶阀/球阀/PRV/FCV',
                    '水轮机：混流/轴流/冲击',
                    '船闸闸门、升船机',
                ]},
                {'title': 'T2 蓄水体类', 'color': C['green'], 'items': [
                    '水库、湖泊',
                    '水池/水箱',
                    '调蓄池',
                    '蓄滞洪区',
                    '',
                    'compute() -> 库容变化',
                ]},
                {'title': 'T3 输水体类', 'color': C['orange'], 'items': [
                    '明渠：梯形/矩形/U形',
                    '管道：球墨铸铁/PCCP/钢/PE',
                    '河道、隧洞',
                    '渡槽、倒虹吸',
                    '',
                    'compute() -> Saint-Venant',
                ]},
                {'title': '耦合器类', 'color': C['purple'], 'items': [
                    '水机电接口',
                    '闸-渠接口',
                    '泵-管接口',
                    '船闸充泄水',
                    '梯级耦合',
                ]},
            ],
            color=C['yellow'],
            section="九、对象层与元件库"
        )

        self.layout_cards(
            "第二族：水文元件 + 第三族：模型元件",
            [
                {'title': '水文元件', 'color': C['green'], 'items': [
                    'H1 产流类:',
                    '  流域/子流域/透水面/土壤层',
                    'H2 汇流类:',
                    '  单位线/马斯京根/运动波',
                    'H3 气象驱动:',
                    '  降雨场/蒸发场/温度场',
                    '',
                    '水文 = 外部扰动(Disturbance)',
                    '控制不了，只能预测和应对',
                ]},
                {'title': '模型元件', 'color': C['purple'], 'items': [
                    'M1 替代模型:',
                    '  IDZ/POD/DMD/PINN',
                    'M2 时序预测:',
                    '  LSTM/Transformer/GNN',
                    'M3 策略模型:',
                    '  RL(PPO/SAC)/模仿学习',
                    'M4 概念水文:',
                    '  新安江/HBV/GR4J',
                    '',
                    '同一元件可选多种计算方法',
                ]},
            ],
            color=C['yellow'],
            section="九、对象层与元件库"
        )

        self.layout_code_demo(
            "统一基类 -- compute()接口",
            [
                "class Component(ABC):",
                "  @abstractmethod",
                "  def compute(self, state, inputs, dt):",
                '    """物理解方程/水文算产汇流',
                '       模型做推断"""',
                "",
                "  def get_params(self) -> ParamSet",
                "  def get_ports(self) -> List[Port]",
                "  def get_constraints(self) -> ...",
                "",
                "# 物理元件：",
                "class OpenChannel(Component):",
                "  def compute(self, state, inputs, dt):",
                "    return saint_venant_step(...)",
                "",
                "# 模型元件：",
                "class LSTMPredictor(Component):",
                "  def compute(self, state, inputs, dt):",
                "    return self.model.forward(...)",
            ],
            [
                "统一基类的威力：",
                "",
                ">> 物理元件和AI模型元件",
                "   在同一个水网里混合使用",
                "",
                ">> 仿真引擎驱动物理元件",
                "   解Saint-Venant方程",
                "",
                ">> 预测引擎驱动模型元件",
                "   做LSTM/GNN推断",
                "",
                ">> 学习引擎驱动策略元件",
                "   做RL策略推断",
                "",
                ">> 像搭积木一样组装水网",
                "   新场景 = 新元件 + 参数",
            ],
            color=C['yellow'],
            section="九、对象层与元件库"
        )

        # ═══ 第十部分：应用案例与展望（46-50页）═══
        self.layout_section("十", "应用案例与展望", C['red'],
                           "从单点工具到全流程革命")

        self.layout_title_bullets(
            "运行示例：断面优化全流程",
            [
                "Step 1 -- 认知层：用户输入「优化河道断面」",
                "  >> 大模型解析意图 -> 匹配规则 -> 选择「断面优化」Skill",
                "  >> 自动加载设计规则（规范遵从、安全系数、变更追溯）",
                "",
                "Step 2 -- 技能层：Skill编排多个引擎",
                "  >> 调用优化引擎（多目标搜索）+ 仿真引擎（水力验证）",
                "  >> 调用工具箱（Manning公式、壅水计算）",
                "",
                "Step 3 -- 引擎层：执行算法",
                "  >> 优化引擎搜索最优断面参数",
                "  >> 仿真引擎验证水力性能（流速、水深、Fr数）",
                "",
                "Step 4 -- 对象层：读取/更新元件",
                "  >> 读取河道元件参数 -> 更新设计参数 -> 验证约束",
                "",
                "Step 5 -- 输出：优化参数 + 3D可视化 + 性能报告",
            ],
            color=C['red'],
            section="十、应用案例与展望"
        )

        self.layout_image_full(
            "设计院三阶段渗透策略",
            "chart_timeline.png",
            color=C['red'],
            caption="从单点工具替代Excel -> MBD验证降低返工 -> 全流程自动化效率提升10倍"
        )

        self.layout_image_full(
            "HydroClaw 核心效果指标",
            "chart_kpi_dashboard.png",
            color=C['red']
        )

        self.layout_summary(
            "总结与展望",
            [
                (">>", "从「全民养虾」到水网自主运行 -- 统一认知智能框架"),
                (">>", "垂直大模型 + 规则引擎兜底 -- 确定性用规则，不确定性用AI"),
                (">>", "个人水网智能体 -- 选Skill即获角色，行为即身份"),
                (">>", "四层架构 -- 认知-技能-引擎-对象，单向调用，职责清晰"),
                (">>", "七大通用引擎 + 100+工具 -- 与领域无关的纯算法层"),
                (">>", "三族元件统一基类 -- 物理+水文+模型，混合组装任意水网"),
            ],
            slogan="HydroClaw：让水网像大脑一样思考"
        )

    def save(self):
        Path(OUT_FILE).parent.mkdir(exist_ok=True)
        self.prs.save(OUT_FILE)
        print(f"\n{'='*50}")
        print(f"PPT已保存: {OUT_FILE}")
        print(f"总页数: {self.page} 页")
        size_kb = Path(OUT_FILE).stat().st_size // 1024
        print(f"文件大小: {size_kb}KB")
        print(f"{'='*50}")


if __name__ == "__main__":
    ppt = HydroClawPPT()
    ppt.generate()
    ppt.save()
