#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HydroClaw 发布会PPT V4 - 80页完整版
基于三层产品架构(HydroOS + HydroMAS + HydroTouch)重构
融入完整代码库：1851测试、21 MCP服务、15 Agent、17 Skill
支持AI生成图片(Nano Banana)占位符 + matplotlib图表 + 多种版式
"""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ─── 配置 ───
IMG_DIR = Path("D:/cowork/ppt/发布会26.4/images")
GEN_DIR = IMG_DIR / "generated"
OUT_FILE = "D:/cowork/ppt/发布会26.4/output/HydroClaw_发布会_V4.pptx"

# 配色方案 (深色科技风)
C = {
    'blue': RGBColor(91, 155, 213),
    'dark_blue': RGBColor(31, 78, 120),
    'navy': RGBColor(13, 27, 42),
    'orange': RGBColor(237, 125, 49),
    'green': RGBColor(112, 173, 71),
    'yellow': RGBColor(255, 192, 0),
    'purple': RGBColor(142, 124, 195),
    'cyan': RGBColor(68, 114, 196),
    'red': RGBColor(192, 0, 0),
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(127, 127, 127),
    'light_gray': RGBColor(242, 242, 242),
    'bg': RGBColor(13, 27, 42),
    'bg_card': RGBColor(27, 40, 56),
    'bg_card2': RGBColor(20, 33, 48),
    'accent': RGBColor(0, 180, 216),
    'text_dim': RGBColor(160, 170, 180),
}

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class HydroClawPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.page = 0

    # ═══════════════════════════════════════
    # 基础工具方法
    # ═══════════════════════════════════════

    def _slide(self):
        self.page += 1
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _bg(self, slide, color=None):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = color or C['bg']

    def _rect(self, slide, left, top, w, h, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _box(self, slide, left, top, w, h, color):
        """方角矩形"""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _text(self, slide, left, top, w, h, text, size=14, bold=False, color=None,
              align=PP_ALIGN.LEFT, font=None):
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font or FONT_BODY
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color or C['white']
        p.alignment = align
        return tb

    def _multitext(self, slide, left, top, w, h, items, size=13, color=None, spacing=6, bullet=""):
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"{bullet}{item}" if bullet else item
            p.font.name = FONT_BODY
            p.font.size = Pt(size)
            p.font.color.rgb = color or C['white']
            p.space_after = Pt(spacing)
        return tb

    def _image_or_placeholder(self, slide, img_name, left, top, w, h, label=""):
        """加载图片，不存在则显示占位符"""
        # 优先查找 generated/ 目录
        path = GEN_DIR / img_name
        if not path.exists():
            path = IMG_DIR / img_name
        if path.exists():
            slide.shapes.add_picture(str(path), left, top, width=w, height=h)
            return True
        else:
            # 占位符: 深色卡片 + 说明文字
            self._rect(slide, left, top, w, h, C['bg_card2'])
            self._text(slide, left + Inches(0.3), top + h // 2 - Inches(0.3),
                       w - Inches(0.6), Inches(0.6),
                       f"[AI生成图片占位] {label or img_name}",
                       size=11, color=C['text_dim'], align=PP_ALIGN.CENTER)
            return False

    def _chart_image(self, slide, chart_name, left, top, w, h):
        """加载 matplotlib 图表"""
        path = IMG_DIR / chart_name
        if path.exists():
            slide.shapes.add_picture(str(path), left, top, width=w, height=h)
            return True
        return False

    def _page_num(self, slide):
        self._text(slide, SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.4),
                   Inches(0.6), Inches(0.3), str(self.page),
                   size=9, color=C['gray'], align=PP_ALIGN.RIGHT)

    def _section_bar(self, slide, section, color):
        """左侧章节指示条"""
        self._box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, color)
        self._text(slide, Inches(0.12), SLIDE_H - Inches(0.4),
                   Inches(2), Inches(0.3), section, size=8, color=C['gray'])

    def _footer(self, slide, section="", color=None):
        """页脚：页码 + 章节"""
        self._page_num(slide)
        if section:
            self._section_bar(slide, section, color or C['blue'])

    # ═══════════════════════════════════════
    # 版式方法
    # ═══════════════════════════════════════

    def cover(self, title, subtitle, desc="", img="cover_hero.png"):
        slide = self._slide()
        self._bg(slide)
        # 顶部装饰线
        self._box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), C['accent'])
        self._box(slide, Inches(0), SLIDE_H - Inches(0.05), SLIDE_W, Inches(0.05), C['accent'])
        # 背景图 (如果有)
        self._image_or_placeholder(slide, img, Inches(6.5), Inches(0.5), Inches(6.5), Inches(6.5), "封面主视觉")
        # 左侧装饰
        self._box(slide, Inches(0.8), Inches(2), Inches(0.05), Inches(3.5), C['accent'])
        # 文字
        self._text(slide, Inches(1.3), Inches(2.2), Inches(5.5), Inches(1.5),
                   title, size=44, bold=True)
        self._text(slide, Inches(1.3), Inches(3.8), Inches(5.5), Inches(0.8),
                   subtitle, size=22, color=C['yellow'])
        if desc:
            self._text(slide, Inches(1.3), Inches(4.8), Inches(5.5), Inches(0.6),
                       desc, size=13, color=C['text_dim'])
        self._text(slide, Inches(1.3), Inches(6.2), Inches(6), Inches(0.4),
                   '中国水利水电科学研究院  |  2026年4月', size=12, color=C['gray'])
        print(f"  [{self.page:02d}] 封面: {title}")

    def toc(self, title, items):
        slide = self._slide()
        self._bg(slide)
        self._text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
                   title, size=32, bold=True)
        self._box(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), C['accent'])
        cols = 2
        per_col = (len(items) + 1) // cols
        for i, item in enumerate(items):
            col = i // per_col
            row = i % per_col
            x = Inches(1 + col * 6)
            y = Inches(1.8 + row * 0.6)
            self._text(slide, x, y, Inches(0.5), Inches(0.5),
                       f"{i+1:02d}", size=16, bold=True, color=C['accent'])
            self._text(slide, x + Inches(0.6), y + Inches(0.02), Inches(5), Inches(0.5),
                       item, size=14)
        self._page_num(slide)
        print(f"  [{self.page:02d}] 目录")

    def section_title(self, num, title, color, subtitle=""):
        slide = self._slide()
        self._bg(slide, color)
        self._text(slide, Inches(1.5), Inches(2), Inches(3), Inches(0.8),
                   f"Part {num}", size=18, color=RGBColor(255, 255, 255))
        self._box(slide, Inches(1.5), Inches(2.8), Inches(2), Inches(0.04), C['white'])
        self._text(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.5),
                   title, size=38, bold=True)
        if subtitle:
            self._text(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.6),
                       subtitle, size=14, color=RGBColor(220, 220, 220))
        self._page_num(slide)
        print(f"  [{self.page:02d}] === {title} ===")

    def title_bullets(self, title, bullets, section="", section_color=None):
        slide = self._slide()
        self._bg(slide)
        self._text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                   title, size=28, bold=True)
        self._box(slide, Inches(0.8), Inches(1.1), Inches(1.2), Inches(0.03), C['accent'])
        self._multitext(slide, Inches(1), Inches(1.5), Inches(11), Inches(5.5),
                        bullets, size=14, spacing=10, bullet="  ")
        self._footer(slide, section, section_color)
        print(f"  [{self.page:02d}] {title}")

    def title_bullets_image(self, title, bullets, img, img_label="", section="", section_color=None):
        """左文右图"""
        slide = self._slide()
        self._bg(slide)
        self._text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
                   title, size=28, bold=True)
        self._box(slide, Inches(0.8), Inches(1.1), Inches(1.2), Inches(0.03), C['accent'])
        self._multitext(slide, Inches(1), Inches(1.5), Inches(5.5), Inches(5.5),
                        bullets, size=13, spacing=8, bullet="  ")
        self._image_or_placeholder(slide, img, Inches(7), Inches(1), Inches(5.8), Inches(5.8), img_label)
        self._footer(slide, section, section_color)
        print(f"  [{self.page:02d}] {title}")

    def image_full(self, title, img, img_label="", caption="", section="", section_color=None):
        """全幅图片页"""
        slide = self._slide()
        self._bg(slide)
        self._text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                   title, size=26, bold=True)
        self._image_or_placeholder(slide, img, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.6), img_label)
        if caption:
            self._text(slide, Inches(0.8), SLIDE_H - Inches(0.6), Inches(11), Inches(0.4),
                       caption, size=10, color=C['text_dim'])
        self._footer(slide, section, section_color)
        print(f"  [{self.page:02d}] {title}")

    def cards(self, title, card_items, cols=3, section="", section_color=None):
        """卡片布局: card_items = [(标题, 内容, 颜色), ...]"""
        slide = self._slide()
        self._bg(slide)
        self._text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                   title, size=28, bold=True)
        self._box(slide, Inches(0.8), Inches(1.1), Inches(1.2), Inches(0.03), C['accent'])

        rows = (len(card_items) + cols - 1) // cols
        card_w = Inches((12 - 0.3 * (cols - 1)) / cols)
        card_h = Inches(min(5.5 / rows - 0.2, 2.5))

        for i, (ctitle, content, color) in enumerate(card_items):
            col = i % cols
            row = i // cols
            x = Inches(0.6) + col * (card_w + Inches(0.3))
            y = Inches(1.5) + row * (card_h + Inches(0.2))

            self._rect(slide, x, y, card_w, card_h, C['bg_card'])
            # 顶部色条
            self._box(slide, x, y, card_w, Inches(0.04), color)
            self._text(slide, x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.4), Inches(0.4),
                       ctitle, size=13, bold=True, color=color)
            if isinstance(content, list):
                self._multitext(slide, x + Inches(0.2), y + Inches(0.55),
                                card_w - Inches(0.4), card_h - Inches(0.7),
                                content, size=10, spacing=3)
            else:
                self._text(slide, x + Inches(0.2), y + Inches(0.55),
                           card_w - Inches(0.4), card_h - Inches(0.7),
                           content, size=10, color=C['text_dim'])

        self._footer(slide, section, section_color)
        print(f"  [{self.page:02d}] {title}")

    def comparison(self, title, left_title, left_items, right_title, right_items,
                   section="", section_color=None):
        """对比布局"""
        slide = self._slide()
        self._bg(slide)
        self._text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                   title, size=28, bold=True)
        self._box(slide, Inches(0.8), Inches(1.1), Inches(1.2), Inches(0.03), C['accent'])

        # 左边
        self._rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), C['bg_card'])
        self._text(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5),
                   left_title, size=16, bold=True, color=C['orange'])
        self._multitext(slide, Inches(0.8), Inches(2.2), Inches(5.2), Inches(4.3),
                        left_items, size=12, spacing=6)

        # VS
        self._text(slide, Inches(6.1), Inches(3.5), Inches(1.1), Inches(0.6),
                   "VS", size=20, bold=True, color=C['accent'], align=PP_ALIGN.CENTER)

        # 右边
        self._rect(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.3), C['bg_card'])
        self._text(slide, Inches(7.3), Inches(1.6), Inches(5), Inches(0.5),
                   right_title, size=16, bold=True, color=C['green'])
        self._multitext(slide, Inches(7.3), Inches(2.2), Inches(5.2), Inches(4.3),
                        right_items, size=12, spacing=6)

        self._footer(slide, section, section_color)
        print(f"  [{self.page:02d}] {title}")

    def big_number(self, title, numbers, section="", section_color=None):
        """大数字展示: numbers = [(数字, 单位, 说明, 颜色), ...]"""
        slide = self._slide()
        self._bg(slide)
        self._text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                   title, size=28, bold=True)

        cols = len(numbers)
        col_w = 12.0 / cols
        for i, (num, unit, desc, color) in enumerate(numbers):
            x = Inches(0.6 + i * col_w)
            self._text(slide, x, Inches(2.2), Inches(col_w), Inches(1.5),
                       str(num), size=60, bold=True, color=color, align=PP_ALIGN.CENTER)
            self._text(slide, x, Inches(3.8), Inches(col_w), Inches(0.5),
                       unit, size=16, color=color, align=PP_ALIGN.CENTER)
            self._text(slide, x, Inches(4.5), Inches(col_w), Inches(0.8),
                       desc, size=12, color=C['text_dim'], align=PP_ALIGN.CENTER)

        self._footer(slide, section, section_color)
        print(f"  [{self.page:02d}] {title}")

    def table_slide(self, title, headers, rows, section="", section_color=None):
        """表格页"""
        slide = self._slide()
        self._bg(slide)
        self._text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                   title, size=28, bold=True)

        cols = len(headers)
        col_w = Inches(11.5 / cols)
        # 表头
        for j, h in enumerate(headers):
            x = Inches(0.8 + j * 11.5 / cols)
            self._rect(slide, x, Inches(1.3), col_w - Inches(0.05), Inches(0.45), C['dark_blue'])
            self._text(slide, x + Inches(0.1), Inches(1.35), col_w - Inches(0.2), Inches(0.35),
                       h, size=11, bold=True, color=C['white'])
        # 数据行
        for i, row in enumerate(rows):
            bg = C['bg_card'] if i % 2 == 0 else C['bg_card2']
            for j, cell in enumerate(row):
                x = Inches(0.8 + j * 11.5 / cols)
                y = Inches(1.8 + i * 0.42)
                self._rect(slide, x, y, col_w - Inches(0.05), Inches(0.4), bg)
                self._text(slide, x + Inches(0.1), y + Inches(0.05),
                           col_w - Inches(0.2), Inches(0.3),
                           str(cell), size=10, color=C['white'])

        self._footer(slide, section, section_color)
        print(f"  [{self.page:02d}] {title}")

    def chart_page(self, title, chart_name, caption="", section="", section_color=None):
        """matplotlib图表页"""
        slide = self._slide()
        self._bg(slide)
        self._text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                   title, size=26, bold=True)
        ok = self._chart_image(slide, chart_name, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.6))
        if not ok:
            self._image_or_placeholder(slide, chart_name, Inches(0.5), Inches(1.1),
                                       Inches(12.3), Inches(5.6), chart_name)
        if caption:
            self._text(slide, Inches(0.8), SLIDE_H - Inches(0.6), Inches(11), Inches(0.4),
                       caption, size=10, color=C['text_dim'])
        self._footer(slide, section, section_color)
        print(f"  [{self.page:02d}] {title}")

    def quote_page(self, quote, author="", section="", section_color=None):
        """引言页"""
        slide = self._slide()
        self._bg(slide, C['dark_blue'])
        self._text(slide, Inches(2), Inches(2), Inches(9), Inches(3),
                   f'"{quote}"', size=28, color=C['white'], align=PP_ALIGN.CENTER)
        if author:
            self._text(slide, Inches(2), Inches(5), Inches(9), Inches(0.5),
                       f'— {author}', size=16, color=C['yellow'], align=PP_ALIGN.CENTER)
        self._page_num(slide)
        print(f"  [{self.page:02d}] 引言")

    def ending(self, title, subtitle, contact=""):
        slide = self._slide()
        self._bg(slide)
        self._box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), C['accent'])
        self._box(slide, Inches(0), SLIDE_H - Inches(0.05), SLIDE_W, Inches(0.05), C['accent'])
        self._image_or_placeholder(slide, "closing_vision.png",
                                   Inches(0), Inches(0), SLIDE_W, SLIDE_H, "结语愿景图")
        # 半透明遮罩
        self._rect(slide, Inches(2), Inches(2), Inches(9.3), Inches(3.5),
                   RGBColor(13, 27, 42))
        self._text(slide, Inches(2.5), Inches(2.5), Inches(8.3), Inches(1.5),
                   title, size=40, bold=True, align=PP_ALIGN.CENTER)
        self._text(slide, Inches(2.5), Inches(4), Inches(8.3), Inches(0.8),
                   subtitle, size=18, color=C['yellow'], align=PP_ALIGN.CENTER)
        if contact:
            self._text(slide, Inches(2.5), Inches(5), Inches(8.3), Inches(0.5),
                       contact, size=12, color=C['gray'], align=PP_ALIGN.CENTER)
        print(f"  [{self.page:02d}] 结语: {title}")

    # ═══════════════════════════════════════
    # 构建80页PPT内容
    # ═══════════════════════════════════════

    def build(self):
        print("=" * 60)
        print("HydroClaw 发布会PPT V4 — 80页完整版")
        print("=" * 60)

        SC = {  # section colors
            '总览': C['blue'],
            '架构': C['dark_blue'],
            '智能体': C['cyan'],
            '认知': C['purple'],
            '接入': C['green'],
            '应用': C['orange'],
            '技术': C['dark_blue'],
            '路线': C['accent'],
        }

        # ═══ 第0部分: 封面+目录 ═══
        self.cover(
            "HydroClaw",
            "水网认知智能体系",
            "从感知到决策，从工具到智慧  |  三层架构 · 多智能体 · 全端接入"
        )

        self.toc("目录 CONTENTS", [
            "产品愿景与定位",
            "三层产品架构总览",
            "Layer 1: HydroOS 计算底座",
            "Layer 2: HydroMAS 智能中枢",
            "认知决策与技能体系",
            "Layer 3: HydroTouch 多端接入",
            "核心应用场景",
            "技术纵深与安全合规",
            "发展路线图",
            "总结与展望",
        ])

        # ═══ 第1部分: 愿景与定位 (6页) ═══
        self.section_title("01", "产品愿景与定位", SC['总览'],
                           "水网全生命周期管理的多智能体智能决策平台")

        self.title_bullets_image(
            "行业痛点与机遇",
            [
                "水网运行高度依赖人工经验，缺乏智能决策支撑",
                "水文模型、水力学模型散落各处，无法统一调用",
                "科研成果到工程落地存在巨大鸿沟",
                "多端协同困难：办公室、调度中心、巡检现场断裂",
                "AI大模型快速发展，水利行业亟需拥抱认知智能",
            ],
            "vision_smart_water.png", "智慧水网愿景",
            "愿景", SC['总览']
        )

        self.quote_page(
            "让每一滴水都有AI守护\n让每一个水利人都有智能助手",
            "HydroClaw 设计理念"
        )

        self.title_bullets(
            "HydroClaw 是什么？",
            [
                "水网全生命周期管理的多智能体智能决策平台",
                "基于 HydroMAS 认知内核 + OpenClaw 代理架构",
                "覆盖科研、设计、运维、教学四大场景",
                "融合 CHS 理论体系：感知→认知→决策→控制",
                "已实现：15 Agent + 17 Skill + 21 MCP Server + 1851 Tests",
                "三层产品架构：HydroOS (算力) + HydroMAS (智能) + HydroTouch (接入)",
            ],
            "愿景", SC['总览']
        )

        self.big_number(
            "核心数据一览",
            [
                ("1851", "Tests", "全部通过的自动化测试", C['green']),
                ("15", "Agents", "领域+DevOps+内容智能体", C['blue']),
                ("21", "MCP", "标准化计算服务接口", C['cyan']),
                ("17", "Skills", "可编排的业务工作流", C['orange']),
            ],
            "愿景", SC['总览']
        )

        self.cards(
            "六大产品版本",
            [
                ("Personal 个人版", ["水利人的随身智库", "基础问答 + 7B 轻量模型", "个人知识管理"], C['blue']),
                ("Mobile 移动版", ["巡检+语音+告警", "1.5B 极轻量推理", "离线可用"], C['green']),
                ("Education 教育版", ["虚拟实验室+AI伴学", "三大教材沉浸渲染", "一键穿越沙箱"], C['purple']),
                ("Research 科研版", ["超级研究助理", "文献检索+大跨度仿真", "论文辅助写作"], C['orange']),
                ("Design 设计版", ["AI总工", "参数寻优+多方案设计", "四级校审流程"], C['cyan']),
                ("O&M 运维版", ["智能中枢神经", "私有化部署+四预闭环", "安全一票否决"], C['red']),
            ],
            cols=3, section="愿景", section_color=SC['总览']
        )

        # ═══ 第2部分: 三层架构总览 (5页) ═══
        self.section_title("02", "三层产品架构", SC['架构'],
                           "HydroOS + HydroMAS + HydroTouch")

        self.image_full(
            "三层产品架构全景",
            "three_layer_overview.png", "三层架构",
            "Layer 1: HydroOS (模型计算层) → Layer 2: HydroMAS (智能体层) → Layer 3: HydroTouch (多端接入层)",
            "架构", SC['架构']
        )

        self.cards(
            "三层职责划分",
            [
                ("HydroOS 计算底座", [
                    "21个MCP标准化计算服务",
                    "Ray分布式并行计算",
                    "知识图谱(190节点/354边)",
                    "TDengine时序 + Neo4j图数据库",
                    "GPU加速ML/水动力计算",
                ], C['cyan']),
                ("HydroMAS 智能中枢", [
                    "15个AI Agent (领域+DevOps+内容)",
                    "17个Skill工作流",
                    "IntentRouter意图路由",
                    "ReActEngine执行引擎",
                    "记忆/会话/个性化/RBAC",
                ], C['blue']),
                ("HydroTouch 多端接入", [
                    "Web全功能仪表盘",
                    "飞书微应用(已集成)",
                    "HydroDesktop桌面端(Tauri)",
                    "Flutter移动端",
                    "手表/AR眼镜(探索)",
                ], C['green']),
            ],
            cols=3, section="架构", section_color=SC['架构']
        )

        self.title_bullets(
            "跨层通信协议",
            [
                "Layer 1 ↔ Layer 2:  gRPC (高性能、强类型Protobuf) + Kafka/Pulsar消息队列",
                "Layer 2 ↔ Layer 3:  REST + WebSocket + SSE (Agent思考过程流式推送)",
                "Layer 1 内部:  Ray Object Store (零拷贝内存共享)",
                "统一认证:  OIDC/OAuth2.0 + Casdoor + 国密SM2/SM3/SM4全链路加密",
                "数据同步:  CRDTs (Yjs/Automerge) 实现多端无冲突状态合并",
                "语义缓存:  Redis + 向量匹配，高频查询绕过LLM直出",
            ],
            "架构", SC['架构']
        )

        self.chart_page(
            "五层架构详细展开 (L0-L4)",
            "chart_architecture_pyramid.png",
            "L0 数据感知 → L1 分布式计算 → L2 MCP工具 → L3 技能编排 → L4 认知智能体",
            "架构", SC['架构']
        )

        # ═══ 第3部分: HydroOS (8页) ═══
        self.section_title("03", "Layer 1: HydroOS", SC['架构'],
                           "云端/高性能服务器上的重型计算引擎")

        self.title_bullets_image(
            "HydroOS 核心定位",
            [
                "提供水文水力学建模、ML推理、知识检索等核心算力",
                "21个MCP Server: 水文(8) + 核心(5) + 外部(2) + 水网(4) + 知识(2)",
                "分布式计算: Ray + Dask 并行仿真/优化/数据清洗",
                "知识图谱: NetworkX → Neo4j, 190节点/354边/19领域",
                "时序数据库: TDengine 实时存储水位/流量/雨量",
                "扩容策略: Ray on K8s HPA, 防汛期自动扩容",
            ],
            "l0_core_algorithms.png", "14个核心算法模块",
            "HydroOS", SC['架构']
        )

        self.cards(
            "L0: 14个核心算法模块",
            [
                ("simulation 仿真", ["水箱ODE模型", "WNTR管网模型", "数字孪生引擎"], C['blue']),
                ("control 控制", ["PID控制器", "MPC模型预测控制", "阀门/泵站控制"], C['green']),
                ("prediction 预测", ["线性/多项式预测", "需求预测", "混合蒸发预测"], C['orange']),
                ("scheduling 调度", ["LP线性规划优化", "全局调度优化", "入流优化"], C['yellow']),
                ("identification 辨识", ["ARX模型辨识", "最小二乘参数ID", "系统辨识"], C['purple']),
                ("water_balance 水平衡", ["节点残差计算", "全厂水平衡", "异常检测"], C['cyan']),
                ("evaporation 蒸发", ["Merkel冷却塔模型", "焙烧蒸发", "赤泥堆场蒸发"], C['red']),
                ("detection 泄漏检测", ["GNN图神经网络", "声学传感器融合", "管段定位"], C['orange']),
                ("更多模块", ["设计优化", "评估评价", "数据清洗", "ODD安全域", "回用水调度"], C['gray']),
            ],
            cols=3, section="HydroOS", section_color=SC['架构']
        )

        self.image_full(
            "L1: Ray 分布式计算引擎",
            "l1_ray_distributed.png", "Ray分布式计算",
            "实时推理池(高优) + 长时演算池(低优) | GPU加速ML模型和2D水动力计算",
            "HydroOS", SC['架构']
        )

        self.image_full(
            "L2: 21个MCP标准化服务",
            "l2_mcp_servers.png", "MCP服务矩阵",
            "simulation | control | prediction | scheduling | evaluation | design | identification | water_balance | evaporation | leak_detection | reuse | hydrology(8) | CHS_RAG",
            "HydroOS", SC['架构']
        )

        self.table_slide(
            "8个水文学MCP服务详情",
            ["服务", "类名", "端口", "核心工具"],
            [
                ["Runoff 产流", "RunoffMCPServer", "8020", "SCS-CN, 新安江, 简化产流"],
                ["Routing 汇流", "RoutingMCPServer", "8021", "Muskingum, 单位线, M-C法"],
                ["Calibration 率定", "CalibrationMCPServer", "8022", "EnKF自动率定"],
                ["Precipitation 降水", "PrecipitationMCPServer", "8023", "IDW插值, 泰森多边形"],
                ["Hydraulic 1D", "Hydraulic1DMCPServer", "8024", "稳态回水, Preissmann非恒定流"],
                ["Hydraulic 2D", "Hydraulic2DMCPServer", "8025", "浅水方程, 溃坝模拟"],
                ["Assimilation 同化", "AssimilationMCPServer", "8026", "EnKF, 粒子滤波"],
                ["ML Hydro", "MLHydroMCPServer", "8027", "特征工程, sklearn训练/预测"],
            ],
            "HydroOS", SC['架构']
        )

        self.title_bullets_image(
            "知识图谱与RAG检索",
            [
                "CHS知识图谱: 190节点, 354条边, 19个领域",
                "覆盖: 水文学、水力学、水环境、控制论、方法论、应用",
                "TF-IDF字符n-gram向量检索, 指纹缓存自动重建",
                "关键词段落回退搜索 (无sklearn时)",
                "工具: search_chs_knowledge, list_books, read_chapter",
                "支持从飞书知识库同步文档, 扩展RAG数据源",
            ],
            "knowledge_graph_chs.png", "CHS知识图谱",
            "HydroOS", SC['架构']
        )

        self.chart_page(
            "核心KPI指标体系",
            "chart_kpi_dashboard.png",
            "节水率、回用率、能效比、安全合规率 — 量化评价水网运行效果",
            "HydroOS", SC['架构']
        )

        # ═══ 第4部分: HydroMAS 智能体 (10页) ═══
        self.section_title("04", "Layer 2: HydroMAS", SC['智能体'],
                           "部署在用户服务器/桌面的智能中枢")

        self.title_bullets_image(
            "HydroMAS 核心能力",
            [
                "IntentRouter: 规则(60%)+LLM(30%)+混合(10%) 意图路由",
                "ReActEngine: 工具/技能执行引擎, 支持多步推理",
                "15 Agents: 领域(7) + DevOps(4) + 内容(4)",
                "17 Skills: 四预+诊断+调度+报告+开发+内容",
                "RBAC 5角色: operator/designer/researcher/admin/teacher",
                "Memory: 长期记忆+日记+时间衰减搜索",
                "Session: per-user/per-group/main 多租户",
                "个性化: SOUL/USER/IDENTITY 三层人格体系",
            ],
            "l4_multi_agent.png", "15个AI智能体",
            "HydroMAS", SC['智能体']
        )

        self.cards(
            "15个AI智能体全览",
            [
                ("Orchestrator 编排器", ["主入口, 4级路由", "协作工作流调度"], C['blue']),
                ("Planning 规划", ["任务分解", "DAG依赖规划"], C['blue']),
                ("Analysis 分析", ["灵活数据分析", "水平衡/蒸发/回用"], C['blue']),
                ("Report 报告", ["Markdown报告", "日报/泄漏报告"], C['blue']),
                ("Safety 安全", ["ODD守护者", "12维安全边界", "一票否决"], C['red']),
                ("Handuo 瀚铎LLM", ["RAG领域问答", "水网大模型"], C['cyan']),
                ("RL Dispatch 强化学习", ["PPO/规则混合调度", "水资源优化分配"], C['green']),
                ("DevPlanner 开发规划", ["需求→DAG计划", "代码任务分解"], C['purple']),
                ("DevReviewer 代码审查", ["多维度Review", "风格/逻辑/安全/性能"], C['purple']),
                ("DevTester 测试生成", ["自动测试生成", "质量门控"], C['purple']),
                ("DevOrchestrator 开发编排", ["Analyse→Plan→", "Implement→Review→Test"], C['purple']),
                ("ContentPlanner 内容规划", ["需求→内容计划", "多渠道排期"], C['orange']),
                ("ContentReviewer 内容审查", ["结构/风格/技术", "合规性审查"], C['orange']),
                ("ContentPublisher 发布", ["飞书/微信/视频/PPT", "多渠道发布"], C['orange']),
                ("ContentOrchestrator", ["全流程编排", "写作→审查→发布"], C['orange']),
            ],
            cols=5, section="HydroMAS", section_color=SC['智能体']
        )

        self.chart_page(
            "17个Skill工作流",
            "chart_skill_hierarchy.png",
            "S0原子(35个) → S1组合(15个) → S2流程(6个) | 四预+泄漏+蒸发+回用+调度+日报+开发+内容",
            "HydroMAS", SC['智能体']
        )

        self.cards(
            "核心Skill详解",
            [
                ("四预闭环 FourPrediction", ["预报→预警→预演→预案", "完整决策支持链路"], C['blue']),
                ("泄漏诊断 LeakDiagnosis", ["水平衡→异常→GNN", "→定位→声学融合"], C['red']),
                ("蒸发优化 EvapOptimization", ["冷却塔→焙烧→赤泥", "→总量→优化建议"], C['orange']),
                ("回用调度 ReuseScheduling", ["水质匹配→LP优化", "→效益评估"], C['green']),
                ("全局调度 GlobalDispatch", ["需求预测→蒸发估算", "→全局调度→ODD检查"], C['cyan']),
                ("日报生成 DailyReport", ["水平衡→异常→KPI", "→蒸发→自动报告"], C['purple']),
            ],
            cols=3, section="HydroMAS", section_color=SC['智能体']
        )

        self.title_bullets(
            "多智能体基础设施",
            [
                "BaseAgent: 统一生命周期, AgentCard能力描述, AgentStatus状态管理",
                "AgentMessage + MessageBus: 异步消息传递, 发布/订阅模式",
                "AgentRegistry: 按能力/类型发现Agent, 动态注册与注销",
                "AgentContext: 共享黑板 + 调用追踪(Span), 跨Agent状态传递",
                "MultiAgentExecutor: DAG执行引擎, 并行/串行混合调度",
                "AgentHealthMonitor: 健康监控, 心跳检测, 自动故障恢复",
                "CapabilityNegotiator: 能力协商, 最佳Agent匹配",
                "CircuitBreaker + RateLimiter: 熔断与限流保护",
            ],
            "HydroMAS", SC['智能体']
        )

        self.title_bullets_image(
            "HydroClaw 工作台层",
            [
                "Personality: SOUL/USER/IDENTITY 三层人格体系",
                "  - SOUL.md: 小瀚 — 共享灵魂身份",
                "  - IDENTITY.md: 名称/表情/平台适配",
                "  - USER.md: 每组/每人个性化配置",
                "Memory: 长期记忆 + 日记 + 关键词搜索 + 时间衰减评分",
                "Session: per-user/per-group/main 三种作用域",
                "RBAC: 5角色 × 3权限粒度(EXECUTE/READ/DENIED)",
                "Heartbeat: 5项定期检查(系统健康/ODD/水平衡/记忆整理/资源)",
                "Evolution: 交互日志→分析→开发→测试→部署 自进化流水线",
            ],
            "intent_router_brain.png", "意图路由与认知决策",
            "HydroMAS", SC['智能体']
        )

        self.chart_page(
            "七大通用计算引擎",
            "chart_seven_engines.png",
            "仿真 | 辨识 | 调度 | 控制 | 优化 | 预测 | 学习 — 覆盖水网全场景计算需求",
            "HydroMAS", SC['智能体']
        )

        self.title_bullets(
            "边缘计算能力 (关键差异化)",
            [
                "本地LLM: Qwen-2.5-7B / Llama-3-8B via Ollama",
                "  → 断网下可执行基础意图路由和文档问答",
                "本地向量库: Chroma/pgvector 缓存知识图谱子集",
                "  → 无需联网即可进行领域知识检索",
                "离线预案: 预下载防汛预案SOP, 断网自动降级",
                "  → 关键时刻(如通信中断)仍可执行应急响应",
                "边缘驻留: Tauri桌面端 Rust+React, 内存极小",
                "  → 适合水库/泵站等边缘节点7×24长期运行",
            ],
            "HydroMAS", SC['智能体']
        )

        self.comparison(
            "Agent协作 vs 传统单体系统",
            "传统单体系统",
            [
                "单一程序处理所有逻辑",
                "模型耦合度高，修改牵一发动全身",
                "无法动态扩展能力",
                "故障扩散, 无隔离机制",
                "人工配置所有工作流",
                "无学习和进化能力",
            ],
            "HydroClaw 多智能体",
            [
                "15个专项Agent各司其职",
                "MCP标准接口, 松耦合组合",
                "动态注册新Agent/Skill",
                "熔断+限流+健康监控",
                "IntentRouter自动编排",
                "自进化流水线: 交互→分析→开发→测试→部署",
            ],
            "HydroMAS", SC['智能体']
        )

        # ═══ 第5部分: 认知决策 (8页) ═══
        self.section_title("05", "认知决策与技能体系", SC['认知'],
                           "从意图理解到智能决策的完整链路")

        self.image_full(
            "认知决策五步流程",
            "cognitive_decision_flow.png", "认知决策流程",
            "意图理解 → 技能匹配 → 规则加载 → 引擎执行 → 模板渲染",
            "认知决策", SC['认知']
        )

        self.title_bullets_image(
            "意图理解与路由",
            [
                "三种分类策略:",
                "  规则匹配(60%): 关键词、正则模式快速路由",
                "  LLM理解(30%): 复杂语义意图识别",
                "  混合模式(10%): LLM + 规则兜底保障",
                "",
                "输入: '帮我分析一下双容水箱的控制效果'",
                "输出: Intent(skill='S1-CTLD', role='researcher')",
                "",
                "认知API映射 (CHS理论):",
                "  感知: forecast, data_analysis",
                "  认知: warning, odd_assessment",
                "  决策: rehearsal, plan, leak_diagnosis",
                "  控制: optimization_design, global_dispatch",
            ],
            "intent_router_brain.png", "意图路由器",
            "认知决策", SC['认知']
        )

        self.image_full(
            "五层规则继承体系",
            "five_layer_rules.png", "五层规则",
            "L0安全底线(不可覆盖) → L1 Skill规则 → L2 产品规则 → L3 场景规则 → L4 案例参数",
            "认知决策", SC['认知']
        )

        self.cards(
            "五层规则详解",
            [
                ("L0 安全底线", ["不可被任何层覆盖", "水位上限/下限保护", "设备安全阈值", "紧急停机条件"], C['red']),
                ("L1 Skill规则", ["每个Skill携带角色规则", "科研/设计/运维差异化", "输入输出格式规范"], C['orange']),
                ("L2 产品规则", ["6大版本各自配置", "Personal/Edu/Design...", "功能开放度不同"], C['yellow']),
                ("L3 场景规则", ["上下文相关的覆盖", "防汛期 vs 日常运行", "白天 vs 夜间值班"], C['green']),
                ("L4 案例参数", ["用户可自定义", "水箱面积/阀门参数", "控制器增益/权重"], C['blue']),
            ],
            cols=5, section="认知决策", section_color=SC['认知']
        )

        self.image_full(
            "技能体系: S0原子→S1组合→S2流程",
            "skill_hierarchy_s0s1s2.png", "三层技能体系",
            "S0原子Skill(35个): 最小粒度计算操作 | S1组合Skill(15个): 固定工作流 | S2流程Skill(6个): 完整业务流程",
            "认知决策", SC['认知']
        )

        self.table_slide(
            "S2 流程Skill (6个完整业务流程)",
            ["编号", "名称", "步骤", "适用场景"],
            [
                ["S2-4P", "四预闭环", "预报→预警→预演→预案", "防汛调度、日常运行"],
                ["S2-SIL", "SIL策略验证", "建模→控制器设计→仿真验证→报告", "新控制策略上线前验证"],
                ["S2-HIL", "HIL硬件验证", "SIL+硬件接口→实物测试→对比", "控制器硬件在环测试"],
                ["S2-MBD", "MBD全流程", "需求→建模→设计→仿真→优化→验证", "Model-Based Design"],
                ["S2-EMG", "应急响应", "感知→评估→预案→执行→复盘", "突发洪水/设备故障"],
                ["S2-EDU", "教学实验", "理论→建模→实验→分析→报告", "水利工程教学"],
            ],
            "认知决策", SC['认知']
        )

        self.chart_page(
            "认知决策流程全景图",
            "chart_cognitive_flow.png",
            "从用户自然语言输入到结构化决策输出的完整数据流",
            "认知决策", SC['认知']
        )

        self.image_full(
            "24个响应模板 (覆盖CHS六要素)",
            "template_24_grid.png", "24模板矩阵",
            "状态查询 | 分析诊断 | 决策建议 | 控制指令 | 报告生成 | 知识问答",
            "认知决策", SC['认知']
        )

        # ═══ 第6部分: HydroTouch 多端接入 (7页) ═══
        self.section_title("06", "Layer 3: HydroTouch", SC['接入'],
                           "八大终端全覆盖, 统一接入网关")

        self.image_full(
            "多端设备全家福",
            "hydrotouch_multidevice.png", "多端接入",
            "Web | 飞书微应用 | Desktop | Android/iOS | 鸿蒙 | 手表 | AR眼镜",
            "HydroTouch", SC['接入']
        )

        self.table_slide(
            "八大终端详细规划",
            ["终端", "技术栈", "交互范式", "优先级"],
            [
                ["Web 全功能", "React + Vite + TailwindCSS", "仪表盘、计算面板、图谱可视化", "P0 已实现"],
                ["飞书微应用", "飞书开放平台 + H5", "嵌入办公流: 通知+审批+快查", "P0 已集成"],
                ["HydroDesktop", "Tauri (Rust+React)", "本地建模、参数率定、离线运行", "P1"],
                ["Android App", "Flutter", "移动巡检、实时告警、语音指令", "P1"],
                ["iOS App", "Flutter", "同Android", "P1"],
                ["鸿蒙 App", "ArkTS + Flutter适配", "华为生态、IoT联动", "P2"],
                ["智能手表", "WatchOS/WearOS Widget", "震动告警 + 语音指令", "P3"],
                ["AR 眼镜", "ARKit/ARCore + WebXR", "大坝巡检叠加实时数据", "P4"],
            ],
            "HydroTouch", SC['接入']
        )

        self.title_bullets(
            "HydroGateway 统一接入网关",
            [
                "所有终端 → HydroGateway → Layer 2 (HydroMAS)",
                "",
                "统一认证: OIDC/OAuth2.0 + 国密SM2",
                "  Token携带: 用户ID + 角色 + 行政区划代码 + 流域权限",
                "统一协议适配: REST / WebSocket / SSE / gRPC",
                "流量控制: 熔断 + 限流 + 优先级队列",
                "数据脱敏: 敏感坐标→相对坐标, 大模型输入脱敏",
                "统一日志: 全链路追踪, 审计合规",
            ],
            "HydroTouch", SC['接入']
        )

        self.title_bullets_image(
            "HydroDesktop 桌面端 (Tauri)",
            [
                "为什么选Tauri而不是Electron?",
                "  内存极小(~30MB vs 300MB+), 适合边缘长期驻留",
                "  Rust安全性, 无GC, 性能接近原生",
                "  Webview2前端, React组件复用",
                "",
                "核心功能:",
                "  本地LLM离线推理 (Ollama)",
                "  本地知识图谱缓存",
                "  离线防汛预案SOP",
                "  3D水网拓扑可视化",
                "  AI对话助手面板",
            ],
            "hydrodesktop_tauri.png", "桌面端界面",
            "HydroTouch", SC['接入']
        )

        self.title_bullets_image(
            "飞书深度集成 (6大场景)",
            [
                "场景1: 飞书机器人 → AI对话入口",
                "  @水网助手 当前水平衡状态如何？→ 实时数据返回",
                "场景2: 多维表格 → 项目管理 + 数据看板",
                "场景3: 知识库 → RAG增强 (规范/手册/案例)",
                "场景4: 审批流 → 设计文件四级校审",
                "场景5: 告警推送 → 卡片式实时监控",
                "  水平衡异常 → 飞书群告警卡片 → [一键诊断]",
                "场景6: Webhook → 自动化流水线",
            ],
            "feishu_integration.png", "飞书集成架构",
            "HydroTouch", SC['接入']
        )

        self.title_bullets_image(
            "AR眼镜: 大坝巡检数字孪生",
            [
                "通过AR眼镜实现空间计算:",
                "",
                "图像锚定: 自动识别物理水利设施",
                "数据叠加: 实时水位/压力/流量覆盖在真实设施上",
                "眼动追踪: 注视某设备 → 弹出数字孪生状态面板",
                "手势交互: 挥手切换数据维度",
                "语音指令: '显示上游水位趋势' → Agent执行",
                "",
                "应用场景: 大坝定期巡检、管道巡查、设备维护",
            ],
            "ar_dam_inspection.png", "AR巡检",
            "HydroTouch", SC['接入']
        )

        # ═══ 第7部分: 核心应用场景 (10页) ═══
        self.section_title("07", "核心应用场景", SC['应用'],
                           "科研 · 设计 · 运维 · 教学 四大场景")

        self.image_full(
            "四预闭环: 预报→预警→预演→预案",
            "four_prediction_loop.png", "四预闭环",
            "HydroClaw 核心业务流程 — 从数据感知到应急响应的完整决策链",
            "应用", SC['应用']
        )

        self.title_bullets_image(
            "数字孪生水网",
            [
                "物理水网的虚拟镜像, 实时同步运行状态",
                "",
                "核心能力:",
                "  DigitalTwinEngine: 状态估计+数据融合",
                "  实时水位/流量/压力映射",
                "  异常检测: 偏差超阈值自动报警",
                "  预演能力: 在孪生体上模拟各种调度方案",
                "",
                "应用场景:",
                "  防汛调度预演",
                "  管网泄漏快速定位",
                "  设备故障影响预估",
            ],
            "digital_twin_water_network.png", "数字孪生",
            "应用", SC['应用']
        )

        self.title_bullets_image(
            "GNN泄漏检测与定位",
            [
                "基于图自编码器(GAT)的智能泄漏检测:",
                "",
                "Step 1: 构建管网图结构 (节点=接合点, 边=管段)",
                "Step 2: GAT注意力机制学习正常流态特征",
                "Step 3: 重构误差异常 → 检测泄漏存在",
                "Step 4: 注意力权重定位 → 锁定泄漏管段",
                "Step 5: 声学传感器融合 → 精确定位泄漏点",
                "",
                "LeakDiagnosisSkill端到端:",
                "  水平衡→异常检测→GNN定位→声学融合→报告",
            ],
            "leak_detection_gnn.png", "GNN泄漏检测",
            "应用", SC['应用']
        )

        self.title_bullets_image(
            "氧化铝厂水网智能化",
            [
                "日取水10,400m³ (乌江7,800 + 洪渠2,600)",
                "12+车间节点, 蒸发损失~4,200m³/d",
                "",
                "目标:",
                "  节水15-20%",
                "  回用率 36%→50%+",
                "  泵站能耗 -8~12%",
                "",
                "已实现Skill:",
                "  EvapOptimization: Merkel+焙烧+赤泥多源耦合优化",
                "  GlobalDispatch: 需求→蒸发→调度→ODD检查",
                "  DailyReport: 自动化运行日报",
            ],
            "alumina_plant_overview.png", "氧化铝厂全景",
            "应用", SC['应用']
        )

        self.image_full(
            "调度中心大屏",
            "dispatch_center_dashboard.png", "调度中心",
            "水网拓扑实时图 + 水平衡仪表盘 + AI调度建议 + 告警列表",
            "应用", SC['应用']
        )

        self.title_bullets_image(
            "现场运维移动端",
            [
                "运维人员手机端 (飞书App / 小程序):",
                "",
                "今日巡检任务清单 (自动推送)",
                "AI助手实时对话:",
                "  '3号冷却塔出水温度偏高,42°C'",
                "  → Merkel模型分析: 正常范围35-38°C",
                "  → 可能原因: 填料结垢或风机减速",
                "  → [生成维修工单] [联系专家]",
                "",
                "一键工单 → 飞书审批 → 维修 → 自动更新台账",
            ],
            "mobile_field_inspection.png", "现场运维",
            "应用", SC['应用']
        )

        self.cards(
            "四大角色场景矩阵",
            [
                ("科研人员", [
                    "写作+建模+管理=科研",
                    "论文辅助写作",
                    "PID/MPC对比仿真",
                    "可复现实验框架",
                    "文献检索+RAG",
                ], C['blue']),
                ("设计人员", [
                    "写作+建模+管理=MBD设计",
                    "水箱/管网尺寸优化",
                    "灵敏度分析",
                    "四级校审流程",
                    "规范自动引用",
                ], C['cyan']),
                ("运维人员", [
                    "写作+建模=运维",
                    "实时监控+AI诊断",
                    "四预闭环日常执行",
                    "泄漏检测+定位",
                    "自动日报生成",
                ], C['orange']),
                ("教学人员", [
                    "虚拟实验室",
                    "一键穿越推演沙箱",
                    "AI伴学导师",
                    "三大教材沉浸渲染",
                    "学生实验自动评分",
                ], C['purple']),
            ],
            cols=4, section="应用", section_color=SC['应用']
        )

        self.chart_page(
            "传统方法 vs HydroClaw 对比",
            "chart_comparison_radar.png",
            "在响应速度、决策质量、自动化程度、多端协同等维度全面超越传统方式",
            "应用", SC['应用']
        )

        self.table_slide(
            "6大E2E测试场景",
            ["编号", "系统", "场景", "写作", "建模", "管理"],
            [
                ["R1", "双容水箱", "科研", "论文", "PID/MPC仿真", "实验管理"],
                ["D1", "双容水箱", "设计", "设计说明书", "尺寸优化", "校审流程"],
                ["O1", "双容水箱", "运维", "日报", "状态估计", "维修工单"],
                ["R2", "氧化铝厂", "科研", "蒸发论文", "多源耦合", "数据管理"],
                ["D2", "氧化铝厂", "设计", "设计文件", "水网设计", "质量管控"],
                ["O2", "氧化铝厂", "运维", "运行日报", "泄漏检测", "调度管理"],
            ],
            "应用", SC['应用']
        )

        # ═══ 第8部分: 技术纵深 (7页) ═══
        self.section_title("08", "技术纵深与安全合规", SC['技术'],
                           "产品化的最后一公里")

        self.chart_page(
            "规则层继承体系",
            "chart_rule_layers.png",
            "L0安全底线 → L1 Skill规则 → L2 产品规则 → L3 场景规则 → L4 案例参数",
            "技术", SC['技术']
        )

        self.title_bullets_image(
            "安全合规体系",
            [
                "等保三级/四级 认证准备",
                "信创适配: 国产CPU/OS/数据库全链路验证",
                "",
                "加密体系:",
                "  SM2: 非对称加密(公钥认证)",
                "  SM3: 哈希摘要(数据完整性)",
                "  SM4: 对称加密(数据传输)",
                "",
                "访问控制:",
                "  OIDC/OAuth2.0 统一认证",
                "  RBAC + ABAC 混合授权",
                "  大模型输入脱敏Agent",
                "  审计日志全链路追踪",
            ],
            "security_compliance.png", "安全合规",
            "技术", SC['技术']
        )

        self.title_bullets_image(
            "自进化流水线",
            [
                "HydroClaw 独特的AI自我进化能力:",
                "",
                "1. 交互日志记录 (JSONL格式)",
                "2. EvolutionAnalyzer 自动分析:",
                "   - 失败率统计",
                "   - 慢响应检测",
                "   - 技能缺口识别",
                "3. Claude Code 自动开发修复",
                "4. pytest 质量门控 (1851测试)",
                "5. 自动部署到生产环境",
                "",
                "定时任务: 每天凌晨3点自动执行全链路",
            ],
            "self_evolution_flywheel.png", "自进化飞轮",
            "技术", SC['技术']
        )

        self.title_bullets_image(
            "边缘部署方案",
            [
                "推荐组合: 懒猫微服 LC-03 + AI算力舱 X3",
                "",
                "LC-03 作为数据中心:",
                "  托管FastAPI Web平台 + TDengine + Neo4j",
                "  7盘位全固态, 96TB存储容量",
                "  内网穿透零配置远程访问",
                "",
                "AI算力舱 X3:",
                "  Jetson Orin 64GB, 275 TOPS, 60W",
                "  运行GNN泄漏检测实时推理",
                "  支持本地LLM (Qwen-2.5-7B)",
                "",
                "高校服务器做'重活'(训练), 边缘设备做'轻活'(推理)",
            ],
            "edge_computing_deploy.png", "边缘部署",
            "技术", SC['技术']
        )

        self.chart_page(
            "三族元件组件库",
            "chart_component_tree.png",
            "物理元件(闸/泵/阀/管道) + 水文元件(降雨/径流/地下水) + 模型元件(神经网络/GNN/RL)",
            "技术", SC['技术']
        )

        self.chart_page(
            "用户画像矩阵",
            "chart_persona_matrix.png",
            "科研人员/设计工程师/运维调度/教学场景 × 功能需求/工具偏好/痛点分析",
            "技术", SC['技术']
        )

        self.big_number(
            "技术栈统计",
            [
                ("1851", "Tests", "pytest全覆盖自动化测试", C['green']),
                ("14", "Core", "L0核心算法子模块", C['blue']),
                ("21", "MCP", "标准化FastMCP服务", C['cyan']),
                ("126", "New", "HydroClaw工作台新增测试", C['orange']),
            ],
            "技术", SC['技术']
        )

        # ═══ 第9部分: 路线图 (5页) ═══
        self.section_title("09", "发展路线图", SC['路线'],
                           "从MVP到规模化的18个月旅程")

        self.image_full(
            "五阶段发展路线",
            "roadmap_phases.png", "发展路线图",
            "Phase A: MVP核心(当前) → B: 桌面+移动 → C: 多端同构 → D: IoT前沿 → E: 规模化商业化",
            "路线", SC['路线']
        )

        self.cards(
            "Phase A: MVP核心闭环 (当前, 0-3个月)",
            [
                ("已完成", [
                    "21 MCP Servers + 知识图谱",
                    "15 Agents + 17 Skills",
                    "IntentRouter + ReActEngine",
                    "Web前端(React) + 飞书集成",
                    "1851 Tests 全部通过",
                ], C['green']),
                ("进行中", [
                    "FastAPI Hydrology路由",
                    "Docker部署打包",
                    "全链路E2E测试(Codex编写中)",
                    "认知决策层核心开发",
                ], C['yellow']),
                ("计划中", [
                    "Web前端功能增强",
                    "性能优化与压测",
                    "文档完善与API手册",
                    "首批试点用户接入",
                ], C['blue']),
            ],
            cols=3, section="路线", section_color=SC['路线']
        )

        self.cards(
            "Phase B-E: 中远期规划",
            [
                ("B: 桌面+移动 (3-6月)", [
                    "HydroDesktop (Tauri+React)",
                    "Flutter App (Android/iOS)",
                    "HydroGateway统一网关",
                    "OIDC认证中心",
                ], C['green']),
                ("C: 多端同构 (6-9月)", [
                    "CRDTs多端状态同步",
                    "鸿蒙App适配",
                    "边缘计算节点部署",
                    "飞书微应用增强",
                ], C['orange']),
                ("D: IoT前沿 (9-12月)", [
                    "智能手表适配",
                    "AR眼镜大坝巡检原型",
                    "语音终端(智能音箱)",
                    "主动播报水情",
                ], C['purple']),
                ("E: 规模化 (12-18月)", [
                    "多流域/多省SaaS部署",
                    "算力市场(按需付费)",
                    "插件市场(第三方扩展)",
                    "信创认证+等保测评",
                ], C['accent']),
            ],
            cols=4, section="路线", section_color=SC['路线']
        )

        self.chart_page(
            "发展时间线详图",
            "chart_timeline.png",
            "从MVP到商业化的完整里程碑规划",
            "路线", SC['路线']
        )

        self.table_slide(
            "关键技术决策",
            ["决策点", "选择", "理由"],
            [
                ["桌面端框架", "Tauri (非Electron)", "内存极小, Rust安全, 适合边缘"],
                ["移动端框架", "Flutter", "一码三端 (Android/iOS/鸿蒙)"],
                ["L1↔L2 通信", "gRPC", "强类型 + 高性能 + 流式支持"],
                ["L2↔L3 通信", "REST+WS+SSE", "Agent思考过程流式推送"],
                ["图数据库", "Neo4j → NebulaGraph", "分布式支持更好"],
                ["本地LLM", "Qwen-2.5-7B via Ollama", "中文能力强, 资源需求适中"],
                ["状态同步", "CRDTs (Yjs)", "无冲突合并, 支持弱网"],
                ["加密方案", "国密 SM2/SM3/SM4", "水利行业合规要求"],
            ],
            "路线", SC['路线']
        )

        # ═══ 补充: 开源生态与技术栈 (4页) ═══

        self.title_bullets(
            "开源技术栈全览",
            [
                "语言: Python 3.11+, Rust (Tauri), TypeScript (React)",
                "计算: NumPy, SciPy, PuLP (LP优化), Ray (分布式)",
                "AI/ML: PyTorch, torch_geometric (GNN), scikit-learn",
                "水利: WNTR (EPANET管网), 自研水文MCP集群",
                "Web: FastAPI + Jinja2 + React + Vite + TailwindCSS",
                "协议: FastMCP (工具协议), gRPC (层间通信)",
                "数据: TDengine (时序), Neo4j (图), SQLite (边缘)",
                "部署: Docker + K8s + Ollama (本地LLM)",
                "测试: pytest + pytest-asyncio, 1851 tests 全通过",
                "CI/CD: GitHub Actions + 自进化流水线",
            ],
            "技术", SC['技术']
        )

        self.title_bullets(
            "Docker多实例部署架构",
            [
                "docker-compose.yml 编排:",
                "  hydroclaw: 主服务 (FastAPI + 15 Agent + 17 Skill)",
                "  tdengine: 时序数据库 (水位/流量/雨量)",
                "  neo4j: 图数据库 (知识图谱)",
                "",
                "OpenClaw 多实例 (个性化部署):",
                "  openclaw-admin: 管理员角色",
                "  openclaw-student-a/b: 学生实例",
                "  openclaw-peer: 同行评审",
                "  openclaw-dev: 开发调试",
                "",
                "每个实例挂载独立人格文件, 共享核心引擎",
            ],
            "技术", SC['技术']
        )

        self.cards(
            "与主流平台对比",
            [
                ("传统水利软件", [
                    "单一功能(仿真/制图)",
                    "无AI决策能力",
                    "桌面单机运行",
                    "人工操作为主",
                ], C['gray']),
                ("通用AI平台", [
                    "无水利领域知识",
                    "幻觉问题严重",
                    "无物理模型校验",
                    "无法离线运行",
                ], C['orange']),
                ("HydroClaw", [
                    "CHS理论+AI+物理模型",
                    "MCP工具校验防幻觉",
                    "三层架构灵活部署",
                    "离线+在线双模式",
                    "1851测试质量保证",
                ], C['accent']),
            ],
            cols=3, section="技术", section_color=SC['技术']
        )

        self.title_bullets(
            "伴随式学习与知识飞轮",
            [
                "W-L-A 学习闭环 (Watch-Learn-Act):",
                "",
                "1. Watch: 用户在知识门户阅读教材, 划线提问",
                "   → AI结合教材上下文精准解答",
                "2. Learn: 点击[在沙箱推演], 浏览器穿越到HydroWeb",
                "   → 底层瞬间拉起隔离沙箱跑仿真",
                "3. Act: 成功解决真实调控问题",
                "   → 系统检测'测试通过'信号",
                "",
                "知识捕获: 成功经验自动入库",
                "交叉审查: 多引擎验证防幻觉",
                "专家入库: Admin审批 → YAML → RAG热更新",
            ],
            "技术", SC['技术']
        )

        # ═══ 第10部分: 总结与展望 (4页) ═══
        self.section_title("10", "总结与展望", C['accent'],
                           "让水网拥有认知智能")

        self.big_number(
            "项目成果总结",
            [
                ("3", "层", "产品架构\nHydroOS+MAS+Touch", C['blue']),
                ("15", "个", "AI智能体\n领域+DevOps+内容", C['green']),
                ("21", "个", "MCP服务\n标准化计算接口", C['cyan']),
                ("6", "大", "产品版本\n覆盖全行业形态", C['orange']),
            ],
            "展望", C['accent']
        )

        self.title_bullets(
            "核心竞争优势",
            [
                "唯一融合CHS水网理论的认知智能平台",
                "三层架构: 计算与智能分离, 灵活部署",
                "多Agent协作: 超越单一AI的集体智慧",
                "全端接入: 从桌面到AR眼镜, 无处不在",
                "自进化能力: 越用越聪明的AI系统",
                "信创合规: 满足水利行业安全要求",
                "开源生态: MCP标准接口, 社区可扩展",
                "学术+工程双轮驱动: 高校科研→产业落地",
            ],
            "展望", C['accent']
        )

        self.image_full(
            "自进化飞轮: 越用越智能",
            "self_evolution_flywheel.png", "自进化飞轮",
            "交互日志 → 进化分析 → 自动开发 → 测试门控 → 部署 → 用户交互 → ...",
            "展望", C['accent']
        )

        self.ending(
            "HydroClaw",
            "让每一滴水都有AI守护",
            "中国水利水电科学研究院  |  hydroclaw@iwhr.com  |  2026.04"
        )

        return self

    def save(self):
        Path(OUT_FILE).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(OUT_FILE)
        size_kb = Path(OUT_FILE).stat().st_size // 1024
        print(f"\n{'='*60}")
        print(f"保存成功: {OUT_FILE}")
        print(f"页数: {self.page}")
        print(f"大小: {size_kb} KB")
        print(f"{'='*60}")


if __name__ == "__main__":
    ppt = HydroClawPPT().build()
    ppt.save()
