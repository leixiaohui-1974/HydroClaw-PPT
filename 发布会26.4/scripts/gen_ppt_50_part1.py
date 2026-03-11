"""
HydroClaw 50页宽屏PPT生成器
基于现有内容扩展，包含所有图表和素材
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# 配色方案
C = {
    'primary': RGBColor(68, 84, 106),
    'dark': RGBColor(31, 73, 125),
    'blue': RGBColor(91, 155, 213),
    'orange': RGBColor(237, 125, 49),
    'green': RGBColor(112, 173, 71),
    'yellow': RGBColor(255, 192, 0),
    'purple': RGBColor(142, 124, 195),
    'red': RGBColor(192, 0, 0),
    'cyan': RGBColor(0, 176, 240),
    'text': RGBColor(63, 63, 63),
    'white': RGBColor(255, 255, 255),
    'light_bg': RGBColor(240, 240, 240),
    'light_blue': RGBColor(230, 240, 255),
    'light_green': RGBColor(230, 255, 230),
    'light_orange': RGBColor(255, 240, 230),
    'light_yellow': RGBColor(255, 255, 230)
}

# 创建宽屏PPT (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]
page_count = 0
TOTAL_PAGES = 50

def add_header(slide, title_text, color=C['primary']):
    """添加页面标题栏"""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    tf = header.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return header

def add_page_number(slide, page_num):
    """添加页码"""
    page_box = slide.shapes.add_textbox(Inches(12.3), Inches(7), Inches(0.9), Inches(0.35))
    tf = page_box.text_frame
    tf.text = f"{page_num}/{TOTAL_PAGES}"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = C['text']
    p.alignment = PP_ALIGN.RIGHT

def add_footer(slide, text="HydroClaw认知智能体系 | 中国水利水电科学研究院"):
    """添加页脚"""
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(11), Inches(0.35))
    tf = footer_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)

print("="*80)
print("HydroClaw 50页宽屏PPT生成器")
print("="*80)
print(f"格式: 16:9 宽屏 ({prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\")")
print(f"目标页数: {TOTAL_PAGES}")
print("="*80)

# ========== 第1页：封面 ==========
page_count += 1
slide = prs.slides.add_slide(blank)
bg = slide.background
fill = bg.fill
fill.gradient()
fill.gradient_angle = 90
fill.gradient_stops[0].color.rgb = C['dark']
fill.gradient_stops[1].color.rgb = C['primary']

title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(1.2))
tf = title_box.text_frame
tf.text = "HydroClaw认知智能体系"
p = tf.paragraphs[0]
p.font.size = Pt(68)
p.font.bold = True
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.333), Inches(0.7))
tf = subtitle_box.text_frame
tf.text = "从全民养虾到水网自主运行"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.color.rgb = C['blue']
p.alignment = PP_ALIGN.CENTER

subtitle2_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.4), Inches(10.333), Inches(0.6))
tf = subtitle2_box.text_frame
tf.text = "垂直大模型 + 个人水网智能体 + 数字孪生"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

author_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.3), Inches(10.333), Inches(0.5))
tf = author_box.text_frame
tf.text = "雷晓辉 | 中国水利水电科学研究院 | 2026年3月"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

print(f"✓ 第{page_count}页: 封面")

# ========== 第2页：目录 ==========
page_count += 1
slide = prs.slides.add_slide(blank)
add_header(slide, "目录")
add_page_number(slide, page_count)
add_footer(slide)

sections = [
    ("一、核心价值与挑战", "3-5", C['blue']),
    ("二、设计哲学", "6-10", C['orange']),
    ("三、总体架构", "11-15", C['green']),
    ("四、垂直大模型", "16-20", C['purple']),
    ("五、个人水网智能体", "21-25", C['cyan']),
    ("六、认知决策层", "26-30", C['blue']),
    ("七、技能编排层", "31-35", C['orange']),
    ("八、计算引擎层", "36-40", C['green']),
    ("九、对象层与元件库", "41-45", C['yellow']),
    ("十、应用案例与展望", "46-50", C['red'])
]

for i, (title, pages, color) in enumerate(sections):
    y = 1.1 + i * 0.58

    # 序号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.text = str(i+1)
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 标题
    text_box = slide.shapes.add_textbox(Inches(1.6), Inches(y), Inches(9.5), Inches(0.45))
    tf = text_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = C['text']
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 页码
    page_box = slide.shapes.add_textbox(Inches(11.3), Inches(y), Inches(1.5), Inches(0.45))
    tf = page_box.text_frame
    tf.text = f"P{pages}"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = C['text']
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

print(f"✓ 第{page_count}页: 目录")

print("\n" + "="*80)
print("第一部分完成，继续生成...")
print("="*80)
