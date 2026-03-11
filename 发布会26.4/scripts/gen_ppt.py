from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Colors
C = {
    'primary': RGBColor(68, 84, 106),
    'dark': RGBColor(31, 73, 125),
    'blue': RGBColor(91, 155, 213),
    'orange': RGBColor(237, 125, 49),
    'green': RGBColor(112, 173, 71),
    'yellow': RGBColor(255, 192, 0),
    'text': RGBColor(63, 63, 63),
    'white': RGBColor(255, 255, 255)
}

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Cover
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

bg = slide.background
fill = bg.fill
fill.gradient()
fill.gradient_angle = 90
fill.gradient_stops[0].color.rgb = C['dark']
fill.gradient_stops[1].color.rgb = C['primary']

title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
tf = title_box.text_frame
tf.text = "HydroClaw认知智能体系完整方案"
p = tf.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
tf = subtitle_box.text_frame
tf.text = "Cognitive AI Architecture for Autonomous Water Network Operations"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.color.rgb = C['blue']
p.alignment = PP_ALIGN.CENTER

subtitle2_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.6))
tf = subtitle2_box.text_frame
tf.text = "从全民养虾到水网自主运行"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

author_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
tf = author_box.text_frame
tf.text = "雷晓辉 | 中国水利水电科学研究院 | 2026.03"
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

print("Slide 1: Cover created")

# Slide 2: Philosophy
slide = prs.slides.add_slide(blank)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = C['primary']
header.line.fill.background()
tf = header.text_frame
tf.text = "设计哲学"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C['white']
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

principles = [
    "确定性的事用规则，不确定性的事用大模型",
    "用户选了什么Skill就具备什么角色",
    "引擎通用不分领域，差异全在元件和参数",
    "物理元件、水文元件、模型元件统一基类"
]

colors = [C['blue'], C['orange'], C['green'], C['yellow']]

for i, (text, color) in enumerate(zip(principles, colors)):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(1.5 + i*1.3), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.text = str(i+1)
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    text_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5 + i*1.3), Inches(7.5), Inches(0.5))
    tf = text_box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = C['text']
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

print("Slide 2: Philosophy created")

# Slide 3: Architecture
slide = prs.slides.add_slide(blank)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = C['primary']
header.line.fill.background()
tf = header.text_frame
tf.text = "总体架构：四层设计"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C['white']
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

layers = [
    ("认知决策层", "大模型 + 规则引擎", C['blue']),
    ("技能编排层", "Skill体系", C['orange']),
    ("计算引擎层", "七大通用引擎", C['green']),
    ("对象层", "三族元件库", C['yellow'])
]

for i, (title, desc, color) in enumerate(layers):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2 + i*1.4), Inches(7), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C['white']

    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = C['white']
    p.space_before = Pt(6)

print("Slide 3: Architecture created")

# Slide 4: Layer 1 - Cognitive Decision
slide = prs.slides.add_slide(blank)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = C['primary']
header.line.fill.background()
tf = header.text_frame
tf.text = "第①层：认知决策层"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C['white']
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Left box: LLM
left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(4.5), Inches(2.5))
left_box.fill.solid()
left_box.fill.fore_color.rgb = C['blue']
left_box.line.color.rgb = C['blue']
left_box.line.width = Pt(2)
tf = left_box.text_frame
tf.text = "大模型"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

items = ["处理不确定性", "自然语言理解", "意图识别", "上下文推理"]
for item in items:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(16)
    p.font.color.rgb = C['white']
    p.space_before = Pt(8)

# Right box: Rules
right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.2), Inches(4), Inches(2.5))
right_box.fill.solid()
right_box.fill.fore_color.rgb = C['orange']
right_box.line.color.rgb = C['orange']
right_box.line.width = Pt(2)
tf = right_box.text_frame
tf.text = "规则引擎"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

items = ["确定性逻辑", "五层继承", "24个回复模板"]
for item in items:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(16)
    p.font.color.rgb = C['white']
    p.space_before = Pt(8)

# Bottom: 5-step process
process_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(9), Inches(2.8))
process_box.fill.solid()
process_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
process_box.line.color.rgb = C['primary']
process_box.line.width = Pt(2)

tf = process_box.text_frame
tf.text = "五步决策流程"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = C['primary']
p.alignment = PP_ALIGN.CENTER

steps = [
    "1. 意图识别 → 2. 规则匹配 → 3. Skill选择",
    "4. 引擎调用 → 5. 结果封装"
]
for step in steps:
    p = tf.add_paragraph()
    p.text = step
    p.font.size = Pt(16)
    p.font.color.rgb = C['text']
    p.space_before = Pt(10)
    p.alignment = PP_ALIGN.CENTER

print("Slide 4: Cognitive Decision created")

# Slide 5: Layer 2 - Skills
slide = prs.slides.add_slide(blank)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = C['primary']
header.line.fill.background()
tf = header.text_frame
tf.text = "第②层：技能编排层"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C['white']
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

skill_types = [
    ("原子Skill", "单引擎调用", ["预测", "优化", "仿真", "学习"], C['blue']),
    ("组合Skill", "多引擎编排", ["断面优化", "泵站选型", "ODD预警"], C['orange']),
    ("流程Skill", "全流程自动化", ["MBD验证", "设计全流程"], C['green'])
]

for i, (title, desc, examples, color) in enumerate(skill_types):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3 + i*1.9), Inches(8.4), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C['white']

    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = C['white']
    p.space_before = Pt(6)

    p = tf.add_paragraph()
    p.text = "示例：" + "、".join(examples)
    p.font.size = Pt(14)
    p.font.color.rgb = C['white']
    p.space_before = Pt(8)

print("Slide 5: Skills created")

# Slide 6: Layer 3 - Engines
slide = prs.slides.add_slide(blank)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = C['primary']
header.line.fill.background()
tf = header.text_frame
tf.text = "第③层：七大通用引擎"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C['white']
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

engines = [
    ("预测引擎", "水文预报、负荷预测", C['blue']),
    ("优化引擎", "调度优化、断面优化", C['orange']),
    ("仿真引擎", "水力仿真、电力仿真", C['green']),
    ("学习引擎", "模式识别、参数校准", C['yellow']),
    ("验证引擎", "MBD验证、合规检查", C['blue']),
    ("可视引擎", "3D渲染、数据可视化", C['orange']),
    ("协同引擎", "多智能体协作", C['green'])
]

cols = 3
rows = 3
box_w = 2.8
box_h = 1.5
gap_x = 0.3
gap_y = 0.2
start_x = 0.8
start_y = 1.2

for i, (name, desc, color) in enumerate(engines):
    row = i // cols
    col = i % cols
    x = start_x + col * (box_w + gap_x)
    y = start_y + row * (box_h + gap_y)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(box_h))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.text = name
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C['white']
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = C['white']
    p.space_before = Pt(8)
    p.alignment = PP_ALIGN.CENTER

# Note box
note_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.8))
tf = note_box.text_frame
tf.text = "✓ 每个引擎内部可切换后端（开源/商业）  ✓ 通过MCP协议集成外部工具"
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.color.rgb = C['text']
p.alignment = PP_ALIGN.CENTER

print("Slide 6: Engines created")

# Slide 7: Layer 4 - Components
slide = prs.slides.add_slide(blank)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = C['primary']
header.line.fill.background()
tf = header.text_frame
tf.text = "第④层：三族元件库"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C['white']
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

components = [
    ("物理元件", ["水库", "河道", "闸门", "泵站", "管道", "阀门"], C['blue']),
    ("水文元件", ["降雨", "蒸发", "产流", "汇流", "地下水"], C['orange']),
    ("模型元件", ["水动力", "水质", "生态", "经济"], C['green'])
]

for i, (title, items, color) in enumerate(components):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3 + i*1.9), Inches(8.4), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C['white']

    p = tf.add_paragraph()
    p.text = "、".join(items)
    p.font.size = Pt(16)
    p.font.color.rgb = C['white']
    p.space_before = Pt(10)

# Bottom note
note_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.5))
tf = note_box.text_frame
tf.text = "统一基类 BaseComponent，四种参数属性：固有、设计、运行、学习"
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.color.rgb = C['text']
p.alignment = PP_ALIGN.CENTER

print("Slide 7: Components created")

# Slide 8: Runtime Example
slide = prs.slides.add_slide(blank)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = C['primary']
header.line.fill.background()
tf = header.text_frame
tf.text = "运行示例：断面优化Skill"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C['white']
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# User input
user_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.8))
user_box.fill.solid()
user_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
user_box.line.color.rgb = C['blue']
user_box.line.width = Pt(2)
tf = user_box.text_frame
tf.text = '用户："帮我优化这条河道的断面"'
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.color.rgb = C['text']
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Process steps
steps = [
    ("①认知层", "意图识别 → 匹配规则 → 选择'断面优化'Skill", C['blue']),
    ("②技能层", "编排：预测引擎 + 优化引擎 + 仿真引擎", C['orange']),
    ("③引擎层", "调用优化算法（NSGA-II）+ 水力仿真验证", C['green']),
    ("④对象层", "读取河道元件参数 → 计算 → 更新设计参数", C['yellow'])
]

for i, (layer, desc, color) in enumerate(steps):
    y = 2.3 + i * 1.1

    # Arrow
    if i > 0:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.5), Inches(y - 0.4), Inches(1), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(200, 200, 200)
        arrow.line.fill.background()

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(8.4), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.text = f"{layer}：{desc}"
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.color.rgb = C['white']
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Result
result_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.5))
result_box.fill.solid()
result_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
result_box.line.color.rgb = C['green']
result_box.line.width = Pt(2)
tf = result_box.text_frame
tf.text = "输出：优化后的断面参数 + 3D可视化 + 性能对比报告"
p = tf.paragraphs[0]
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = C['text']
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

print("Slide 8: Runtime Example created")

# Slide 9: Design Institute Strategy
slide = prs.slides.add_slide(blank)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = C['primary']
header.line.fill.background()
tf = header.text_frame
tf.text = "设计院切入策略：三阶段渗透"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C['white']
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

phases = [
    ("阶段1：单点工具", "泵站选型、管网计算、断面优化", "替代Excel和经验公式", C['blue']),
    ("阶段2：MBD验证", "全系统数字孪生验证", "发现设计缺陷，降低返工", C['orange']),
    ("阶段3：全流程", "从方案到施工图自动生成", "设计效率提升10倍", C['green'])
]

for i, (title, content, value, color) in enumerate(phases):
    y = 1.3 + i * 2

    # Phase number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.text = str(i+1)
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Content box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.6), Inches(y), Inches(7.6), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C['white']

    p = tf.add_paragraph()
    p.text = content
    p.font.size = Pt(15)
    p.font.color.rgb = C['white']
    p.space_before = Pt(6)

    p = tf.add_paragraph()
    p.text = f"价值：{value}"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = C['white']
    p.space_before = Pt(6)

print("Slide 9: Design Institute Strategy created")

# Slide 10: Summary
slide = prs.slides.add_slide(blank)

bg = slide.background
fill = bg.fill
fill.gradient()
fill.gradient_angle = 90
fill.gradient_stops[0].color.rgb = C['primary']
fill.gradient_stops[1].color.rgb = C['dark']

title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
tf = title_box.text_frame
tf.text = "核心价值"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

values = [
    "✓ 从全民养虾到水网自主运行",
    "✓ 确定性用规则，不确定性用大模型",
    "✓ 通用引擎 + 领域元件 = 无限可能",
    "✓ 设计、运行、教学全场景覆盖"
]

for i, value in enumerate(values):
    text_box = slide.shapes.add_textbox(Inches(1.5), Inches(3 + i*0.8), Inches(7), Inches(0.6))
    tf = text_box.text_frame
    tf.text = value
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = C['white']
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
tf = footer_box.text_frame
tf.text = "HydroClaw：让水网像大脑一样思考"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.italic = True
p.font.color.rgb = C['blue']
p.alignment = PP_ALIGN.CENTER

print("Slide 10: Summary created")

# Save
output = "HydroClaw_认知智能方案_重构版.pptx"
prs.save(output)
print(f"\nPPT saved: {output}")
print(f"Total slides: {len(prs.slides)}")
