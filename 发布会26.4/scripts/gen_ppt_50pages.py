from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Colors
C = {
    'primary': RGBColor(68, 84, 106),
    'dark': RGBColor(31, 73, 125),
    'blue': RGBColor(91, 155, 213),
    'orange': RGBColor(237, 125, 49),
    'green': RGBColor(112, 173, 71),
    'yellow': RGBColor(255, 192, 0),
    'purple': RGBColor(142, 124, 195),
    'red': RGBColor(192, 0, 0),
    'cyan': RGBColor(0, 176, 240),
    'text': RGBColor(63, 63, 63),
    'white': RGBColor(255, 255, 255),
    'light_bg': RGBColor(240, 240, 240),
    'light_blue': RGBColor(230, 240, 255),
    'light_green': RGBColor(230, 255, 230),
    'light_orange': RGBColor(255, 240, 230)
}

# 创建宽屏PPT (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]

def add_header(slide, title_text, color=C['primary']):
    """添加页面标题栏"""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    tf = header.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return header

def add_page_number(slide, page_num, total_pages):
    """添加页码"""
    page_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.7), Inches(0.3))
    tf = page_box.text_frame
    tf.text = f"{page_num}/{total_pages}"
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = C['text']
    p.alignment = PP_ALIGN.RIGHT

page_count = 0

print("开始生成50页宽屏PPT...")
print("="*60)

# ========== 第1页：封面 ==========
page_count += 1
slide = prs.slides.add_slide(blank)
bg = slide.background
fill = bg.fill
fill.gradient()
fill.gradient_angle = 90
fill.gradient_stops[0].color.rgb = C['dark']
fill.gradient_stops[1].color.rgb = C['primary']

title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
tf = title_box.text_frame
tf.text = "HydroClaw认知智能体系"
p = tf.paragraphs[0]
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(0.8))
tf = subtitle_box.text_frame
tf.text = "从全民养虾到水网自主运行"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = C['blue']
p.alignment = PP_ALIGN.CENTER

subtitle2_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.333), Inches(0.6))
tf = subtitle2_box.text_frame
tf.text = "垂直大模型 + 个人水网智能体 + 数字孪生"
p = tf.paragraphs[0]
p.font.size = Pt(22)
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

author_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.5))
tf = author_box.text_frame
tf.text = "雷晓辉 | 中国水利水电科学研究院 | 2026.03"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

print(f"第{page_count}页: 封面")

# ========== 第2页：目录 ==========
page_count += 1
slide = prs.slides.add_slide(blank)
add_header(slide, "目录")
add_page_number(slide, page_count, 50)

sections = [
    ("一、设计哲学", "3-6", C['blue']),
    ("二、总体架构", "7-10", C['orange']),
    ("三、垂直大模型", "11-15", C['green']),
    ("四、个人水网智能体", "16-20", C['purple']),
    ("五、认知决策层", "21-25", C['blue']),
    ("六、技能编排层", "26-30", C['orange']),
    ("七、计算引擎层", "31-35", C['green']),
    ("八、对象层", "36-40", C['yellow']),
    ("九、应用案例", "41-45", C['cyan']),
    ("十、总结与展望", "46-50", C['red'])
]

for i, (title, pages, color) in enumerate(sections):
    y = 1.2 + i * 0.55

    # 序号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.text = str(i+1)
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 标题
    text_box = slide.shapes.add_textbox(Inches(1.6), Inches(y), Inches(9), Inches(0.4))
    tf = text_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = C['text']
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 页码
    page_box = slide.shapes.add_textbox(Inches(11), Inches(y), Inches(1.5), Inches(0.4))
    tf = page_box.text_frame
    tf.text = f"P{pages}"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = C['text']
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

print(f"第{page_count}页: 目录")

# ========== 第3页：设计哲学 - 概述 ==========
page_count += 1
slide = prs.slides.add_slide(blank)
add_header(slide, "一、设计哲学：四大核心原则", C['blue'])
add_page_number(slide, page_count, 50)

principles = [
    "确定性用规则，不确定性用大模型",
    "用户选了什么Skill就具备什么角色",
    "引擎通用不分领域，差异全在元件和参数",
    "物理元件、水文元件、模型元件统一基类"
]

colors = [C['blue'], C['orange'], C['green'], C['yellow']]

for i, (text, color) in enumerate(zip(principles, colors)):
    row = i // 2
    col = i % 2
    x = 1 + col * 6
    y = 1.5 + row * 2.5

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.text = str(i+1)
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    text_box = slide.shapes.add_textbox(Inches(x+0.8), Inches(y), Inches(4.5), Inches(0.6))
    tf = text_box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = C['text']
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

print(f"第{page_count}页: 设计哲学概述")

# ========== 第4页：设计哲学 - 原则1详解 ==========
page_count += 1
slide = prs.slides.add_slide(blank)
add_header(slide, "原则1：确定性用规则，不确定性用大模型", C['blue'])
add_page_number(slide, page_count, 50)

# 左侧：规则引擎
left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.5))
left_box.fill.solid()
left_box.fill.fore_color.rgb = C['orange']
left_box.line.color.rgb = C['orange']
left_box.line.width = Pt(3)

tf = left_box.text_frame
tf.text = "规则引擎处理确定性"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

items = [
    "查询闸门当前水位",
    "执行标准巡检流程",
    "ODD阈值判别",
    "预警等级确定",
    "标准操作流程",
    "物理约束检查",
    "安全规则验证"
]
for item in items:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(16)
    p.font.color.rgb = C['white']
    p.space_before = Pt(10)
    p.level = 0

# 右侧：大模型
right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5))
right_box.fill.solid()
right_box.fill.fore_color.rgb = C['blue']
right_box.line.color.rgb = C['blue']
right_box.line.width = Pt(3)

tf = right_box.text_frame
tf.text = "大模型处理不确定性"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C['white']
p.alignment = PP_ALIGN.CENTER

items = [
    "分析工况并推荐方案",
    "设计方案比选",
    "跨域关联分析",
    "可解释性生成",
    "未预见的异常工况",
    "应急情况类比推理",
    "创新性方案构想"
]
for item in items:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(16)
    p.font.color.rgb = C['white']
    p.space_before = Pt(10)
    p.level = 0

print(f"第{page_count}页: 原则1详解")

# ========== 第5页：设计哲学 - 原则2详解 ==========
page_count += 1
slide = prs.slides.add_slide(blank)
add_header(slide, "原则2：选Skill即获得角色能力", C['blue'])
add_page_number(slide, page_count, 50)

# 概念说明
concept_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.8))
tf = concept_box.text_frame
tf.text = "用户不需要声明角色，不需要切换模式。选择Skill就自动获得对应的规则约束、回复深度和术语风格。"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = C['primary']
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# 示例
examples = [
    ("设计工程师", "断面优化Skill", "自动加载设计规则", C['blue']),
    ("运维人员", "ODD预警Skill", "自动加载运维规则", C['orange']),
    ("科研人员", "仿真对比Skill", "自动加载科研规则", C['green']),
    ("学生教师", "教学实验Skill", "自动加载教学规则", C['purple'])
]

for i, (role, skill, rule, color) in enumerate(examples):
    y = 2.5 + i * 1.1

    # 角色
    role_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(2.5), Inches(0.8))
    role_box.fill.solid()
    role_box.fill.fore_color.rgb = C['light_bg']
    role_box.line.color.rgb = color
    role_box.line.width = Pt(2)
    tf = role_box.text_frame
    tf.text = role
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = C['text']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 箭头
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.5), Inches(y+0.2), Inches(1.2), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()

    # Skill
    skill_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(y), Inches(2.8), Inches(0.8))
    skill_box.fill.solid()
    skill_box.fill.fore_color.rgb = color
    skill_box.line.fill.background()
    tf = skill_box.text_frame
    tf.text = skill
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 箭头2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.9), Inches(y+0.2), Inches(1.2), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = color
    arrow2.line.fill.background()

    # 规则
    rule_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(y), Inches(3.2), Inches(0.8))
    rule_box.fill.solid()
    rule_box.fill.fore_color.rgb = C['light_bg']
    rule_box.line.color.rgb = color
    rule_box.line.width = Pt(2)
    tf = rule_box.text_frame
    tf.text = rule
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = C['text']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

print(f"第{page_count}页: 原则2详解")

# 保存第一部分
print("\n保存第一部分（前5页）...")
output = "HydroClaw_50页宽屏版_part1.pptx"
prs.save(output)
print(f"已保存: {output}")
print(f"当前页数: {page_count}")
