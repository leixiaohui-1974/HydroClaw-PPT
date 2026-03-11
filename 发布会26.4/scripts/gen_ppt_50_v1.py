"""
HydroClaw 50页宽屏PPT完整生成器
使用数据驱动方式批量生成
"""
import sys
import codecs
if sys.platform == 'win32':
    sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer, 'strict')
    sys.stderr = codecs.getwriter('utf-8')(sys.stderr.buffer, 'strict')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 配色方案
C = {
    'primary': RGBColor(68, 84, 106),
    'dark': RGBColor(31, 73, 125),
    'blue': RGBColor(91, 155, 213),
    'orange': RGBColor(237, 125, 49),
    'green': RGBColor(112, 173, 71),
    'yellow': RGBColor(255, 192, 0),
    'purple': RGBColor(142, 124, 195),
    'red': RGBColor(192, 0, 0),
    'cyan': RGBColor(0, 176, 240),
    'text': RGBColor(63, 63, 63),
    'white': RGBColor(255, 255, 255),
    'light_bg': RGBColor(240, 240, 240),
    'light_blue': RGBColor(230, 240, 255),
    'light_green': RGBColor(230, 255, 230),
    'light_orange': RGBColor(255, 240, 230)
}

# 创建宽屏PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

TOTAL_PAGES = 50
page_count = 0

def add_header(slide, title_text, color=C['primary']):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    tf = header.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C['white']
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_page_number(slide, page_num):
    page_box = slide.shapes.add_textbox(Inches(12.3), Inches(7), Inches(0.9), Inches(0.35))
    tf = page_box.text_frame
    tf.text = f"{page_num}/{TOTAL_PAGES}"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = C['text']
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_footer(slide):
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(11), Inches(0.35))
    tf = footer_box.text_frame
    tf.text = "HydroClaw认知智能体系 | 中国水利水电科学研究院"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

print("="*80)
print("HydroClaw 50页宽屏PPT生成器")
print("="*80)

# 页面内容数据
slides_data = [
    # 第1页：封面
    {
        'type': 'cover',
        'title': 'HydroClaw认知智能体系',
        'subtitle': '从全民养虾到水网自主运行',
        'subtitle2': '垂直大模型 + 个人水网智能体 + 数字孪生',
        'author': '雷晓辉 | 中国水利水电科学研究院 | 2026年3月'
    },

    # 第2页：目录
    {
        'type': 'toc',
        'sections': [
            ("一、核心价值与挑战", "3-5", C['blue']),
            ("二、设计哲学", "6-10", C['orange']),
            ("三、总体架构", "11-15", C['green']),
            ("四、垂直大模型", "16-20", C['purple']),
            ("五、个人水网智能体", "21-25", C['cyan']),
            ("六、认知决策层", "26-30", C['blue']),
            ("七、技能编排层", "31-35", C['orange']),
            ("八、计算引擎层", "36-40", C['green']),
            ("九、对象层与元件库", "41-45", C['yellow']),
            ("十、应用案例与展望", "46-50", C['red'])
        ]
    },

    # 第3-5页：核心价值与挑战
    {
        'type': 'section_title',
        'title': '一、核心价值与挑战',
        'subtitle': '从全民养虾到水网自主运行的转变',
        'color': C['blue']
    },
    {
        'type': 'content',
        'header': '当前挑战：全民养虾困境',
        'color': C['blue'],
        'items': [
            ('人工依赖', '调度决策依赖人工经验，难以标准化', C['red']),
            ('响应滞后', '异常工况响应慢，错过最佳处置时机', C['orange']),
            ('知识流失', '专家经验难以传承，新人培养周期长', C['yellow']),
            ('效率低下', '重复性工作占用大量人力资源', C['green'])
        ]
    },
    {
        'type': 'content',
        'header': 'HydroClaw解决方案：智能化转型',
        'color': C['blue'],
        'items': [
            ('认知智能', '大模型理解意图，规则引擎保障安全', C['blue']),
            ('自主运行', '7x24小时监控，秒级响应异常工况', C['green']),
            ('知识固化', '专家经验转化为Skill，永久保存', C['purple']),
            ('效率提升', '自动化处理常规任务，人专注于决策', C['cyan'])
        ]
    },

    # 第6-10页：设计哲学
    {
        'type': 'section_title',
        'title': '二、设计哲学',
        'subtitle': '四大核心原则指导系统设计',
        'color': C['orange']
    },
    {
        'type': 'philosophy',
        'header': '设计哲学：四大核心原则',
        'color': C['orange'],
        'principles': [
            ('确定性用规则，不确定性用大模型', C['blue']),
            ('用户选了什么Skill就具备什么角色', C['orange']),
            ('引擎通用不分领域，差异全在元件和参数', C['green']),
            ('物理元件、水文元件、模型元件统一基类', C['yellow'])
        ]
    },
]

# 继续添加更多页面数据...
# 为了演示，我先生成前10页

print("开始生成页面...")

for slide_data in slides_data:
    page_count += 1

    if slide_data['type'] == 'cover':
        # 封面页
        slide = prs.slides.add_slide(blank)
        bg = slide.background
        fill = bg.fill
        fill.gradient()
        fill.gradient_angle = 90
        fill.gradient_stops[0].color.rgb = C['dark']
        fill.gradient_stops[1].color.rgb = C['primary']

        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(1.2))
        tf = title_box.text_frame
        tf.text = slide_data['title']
        tf.paragraphs[0].font.size = Pt(68)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = C['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.333), Inches(0.7))
        tf = subtitle_box.text_frame
        tf.text = slide_data['subtitle']
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.color.rgb = C['blue']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle2_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.4), Inches(10.333), Inches(0.6))
        tf = subtitle2_box.text_frame
        tf.text = slide_data['subtitle2']
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = C['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        author_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.3), Inches(10.333), Inches(0.5))
        tf = author_box.text_frame
        tf.text = slide_data['author']
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = C['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        print(f"[OK] 第{page_count}页: 封面")

    elif slide_data['type'] == 'toc':
        # 目录页
        slide = prs.slides.add_slide(blank)
        add_header(slide, "目录")
        add_page_number(slide, page_count)
        add_footer(slide)

        for i, (title, pages, color) in enumerate(slide_data['sections']):
            y = 1.1 + i * 0.58

            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y), Inches(0.45), Inches(0.45))
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.fill.background()
            tf = circle.text_frame
            tf.text = str(i+1)
            tf.paragraphs[0].font.size = Pt(20)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = C['white']
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            text_box = slide.shapes.add_textbox(Inches(1.6), Inches(y), Inches(9.5), Inches(0.45))
            tf = text_box.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(22)
            tf.paragraphs[0].font.color.rgb = C['text']
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            page_box = slide.shapes.add_textbox(Inches(11.3), Inches(y), Inches(1.5), Inches(0.45))
            tf = page_box.text_frame
            tf.text = f"P{pages}"
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = C['text']
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        print(f"[OK] 第{page_count}页: 目录")

    elif slide_data['type'] == 'section_title':
        # 章节标题页
        slide = prs.slides.add_slide(blank)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = slide_data['color']

        title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.333), Inches(1.5))
        tf = title_box.text_frame
        tf.text = slide_data['title']
        tf.paragraphs[0].font.size = Pt(56)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = C['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.333), Inches(0.8))
        tf = subtitle_box.text_frame
        tf.text = slide_data['subtitle']
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.color.rgb = C['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_page_number(slide, page_count)

        print(f"[OK] 第{page_count}页: {slide_data['title']}")

    elif slide_data['type'] == 'content':
        # 内容页
        slide = prs.slides.add_slide(blank)
        add_header(slide, slide_data['header'], slide_data['color'])
        add_page_number(slide, page_count)
        add_footer(slide)

        for i, (title, desc, color) in enumerate(slide_data['items']):
            row = i // 2
            col = i % 2
            x = 0.8 + col * 6.3
            y = 1.5 + row * 2.5

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(2))
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.color.rgb = color
            box.line.width = Pt(3)

            tf = box.text_frame
            tf.text = title
            p = tf.paragraphs[0]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = C['white']
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = desc
            p.font.size = Pt(18)
            p.font.color.rgb = C['white']
            p.space_before = Pt(12)
            p.alignment = PP_ALIGN.CENTER

        print(f"[OK] 第{page_count}页: {slide_data['header']}")

    elif slide_data['type'] == 'philosophy':
        # 设计哲学页
        slide = prs.slides.add_slide(blank)
        add_header(slide, slide_data['header'], slide_data['color'])
        add_page_number(slide, page_count)
        add_footer(slide)

        for i, (text, color) in enumerate(slide_data['principles']):
            row = i // 2
            col = i % 2
            x = 1 + col * 6
            y = 1.5 + row * 2.5

            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.6), Inches(0.6))
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.fill.background()
            tf = circle.text_frame
            tf.text = str(i+1)
            tf.paragraphs[0].font.size = Pt(28)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = C['white']
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            text_box = slide.shapes.add_textbox(Inches(x+0.8), Inches(y), Inches(4.5), Inches(0.6))
            tf = text_box.text_frame
            tf.text = text
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = C['text']
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        print(f"[OK] 第{page_count}页: {slide_data['header']}")

# 保存
output = "HydroClaw_50页宽屏版_v1.pptx"
prs.save(output)

print("\n" + "="*80)
print(f"[完成] 已生成 {page_count} 页")
print(f"[文件] {output}")
print("="*80)
