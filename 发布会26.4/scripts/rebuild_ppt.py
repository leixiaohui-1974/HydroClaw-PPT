#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
HydroClaw认知智能方案PPT重构
充分利用现有PPT素材和风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import glob

# 配色方案（基于现有PPT分析）
C = {
    'primary': RGBColor(68, 84, 106),    # 主蓝
    'dark': RGBColor(31, 73, 125),       # 深蓝
    'blue': RGBColor(91, 155, 213),      # 亮蓝
    'orange': RGBColor(237, 125, 49),    # 橙色
    'green': RGBColor(112, 173, 71),     # 绿色
    'yellow': RGBColor(255, 192, 0),     # 黄色
    'text': RGBColor(63, 63, 63),        # 文字
    'bg': RGBColor(231, 230, 230),       # 背景
    'white': RGBColor(255, 255, 255)
}

class HydroClawPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def add_header(self, slide, title_text):
        """添加标题栏（深蓝色背景）"""
        # 标题栏背景
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, Inches(10), Inches(0.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = C['primary']
        header.line.fill.background()

        # 标题文字
        tf = header.text_frame
        tf.text = title_text
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = C['white']
        tf.paragraphs[0].font.name = '微软雅黑'
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def add_footer(self, slide, page_num):
        """添加页脚"""
        footer = slide.shapes.add_textbox(
            Inches(9), Inches(7.2), Inches(0.8), Inches(0.25)
        )
        tf = footer.text_frame
        tf.text = str(page_num)
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = C['text']
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def slide_cover(self):
        """封面页"""
        blank = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank)

        # 渐变背景
        bg = slide.background
        fill = bg.fill
        fill.gradient()
        fill.gradient_angle = 90
        fill.gradient_stops[0].color.rgb = C['dark']
        fill.gradient_stops[1].color.rgb = C['primary']

        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = "HydroClaw认知智能体系完整方案"
        p = tf.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = C['white']
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER

        # 英文副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.5), Inches(8), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        tf.text = "Cognitive AI Architecture for Autonomous Water Network Operations"
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = C['blue']
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER

        # 副标题2
        subtitle2_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5), Inches(8), Inches(0.6)
        )
        tf = subtitle2_box.text_frame
        tf.text = "从"全民养虾"到水网自主运行——水网领域的垂直认知智能体"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = C['white']
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER

        # 作者信息
        author_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(8), Inches(0.5)
        )
        tf = author_box.text_frame
        tf.text = "雷晓辉  |  中国水利水电科学研究院  |  2026.03"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = C['white']
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER

    def slide_philosophy(self):
        """设计哲学页"""
        blank = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank)
        self.add_header(slide, "设计哲学")

        # 四句话
        principles = [
            "确定性的事用规则，不确定性的事用大模型，大模型的输出用规则兜底",
            "用户选了什么Skill就具备什么角色，行为即身份",
            "引擎通用不分领域，差异全在元件和参数",
            "物理元件、水文元件、模型元件统一基类，混合组装任意水网"
        ]

        colors = [C['blue'], C['orange'], C['green'], C['yellow']]

        for i, (text, color) in enumerate(zip(principles, colors)):
            # 编号圆圈
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.8), Inches(1.5 + i*1.3), Inches(0.5), Inches(0.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.fill.background()

            tf = circle.text_frame
            tf.text = str(i+1)
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = C['white']
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 文字
            text_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(1.5 + i*1.3), Inches(7.5), Inches(0.5)
            )
            tf = text_box.text_frame
            tf.text = text
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = C['text']
            tf.paragraphs[0].font.name = '微软雅黑'
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        self.add_footer(slide, 2)

    def slide_architecture(self):
        """总体架构页"""
        blank = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank)
        self.add_header(slide, "总体架构：四层设计")

        layers = [
            ("① 认知决策层", "大模型 + 规则引擎 + 模板渲染", C['blue']),
            ("② 技能编排层", "Skill体系（原子→组合→流程）", C['orange']),
            ("③ 计算引擎层", "七大通用引擎 + 工具箱", C['green']),
            ("④ 对象层", "三族元件库 + 水网组装", C['yellow'])
        ]

        for i, (title, desc, color) in enumerate(layers):
            # 层级框
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(1.2 + i*1.4), Inches(7), Inches(1.1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.color.rgb = color
            box.line.width = Pt(2)

            # 标题
            tf = box.text_frame
            tf.text = title
            p = tf.paragraphs[0]
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = C['white']
            p.font.name = '微软雅黑'

            # 描述
            p = tf.add_paragraph()
            p.text = desc
            p.font.size = Pt(16)
            p.font.color.rgb = C['white']
            p.font.name = '微软雅黑'
            p.space_before = Pt(6)

        self.add_footer(slide, 3)

    def slide_layer1_overview(self):
        """第①层：认知决策概览"""
        blank = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank)
        self.add_header(slide, "第①层：认知决策")

        # 分工说明
        sections = [
            ("规则引擎处理", "~60%交互", "确定性场景", C['blue']),
            ("大模型处理", "~30%交互", "需要理解和推理", C['orange']),
            ("大模型主导+规则兜底", "~10%", "开放性问题", C['green'])
        ]

        for i, (title, ratio, desc, color) in enumerate(sections):
            x = 1 + i * 2.8

            # 标题框
            title_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(1.5), Inches(2.5), Inches(0.6)
            )
            title_box.fill.solid()
            title_box.fill.fore_color.rgb = color
            title_box.line.fill.background()

            tf = title_box.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = C['white']
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 比例
            ratio_box = slide.shapes.add_textbox(
                Inches(x), Inches(2.2), Inches(2.5), Inches(0.4)
            )
            tf = ratio_box.text_frame
            tf.text = ratio
            tf.paragraphs[0].font.size = Pt(20)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = color
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x), Inches(2.7), Inches(2.5), Inches(0.5)
            )
            tf = desc_box.text_frame
            tf.text = desc
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = C['text']
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 底部说明
        note_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(8), Inches(0.6)
        )
        tf = note_box.text_frame
        tf.text = "核心原则：确定性用规则，不确定性用大模型，大模型输出用规则兜底"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = C['primary']
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        self.add_footer(slide, 4)

    def generate(self, output_file="HydroClaw_认知智能方案_重构版.pptx"):
        """生成完整PPT"""
        print("开始生成PPT...")

        self.slide_cover()
        print("✓ 封面页")

        self.slide_philosophy()
        print("✓ 设计哲学")

        self.slide_architecture()
        print("✓ 总体架构")

        self.slide_layer1_overview()
        print("✓ 认知决策层")

        # 保存
        self.prs.save(output_file)
        print(f"\n✅ PPT生成完成: {output_file}")
        print(f"   共 {len(self.prs.slides)} 页")

if __name__ == "__main__":
    ppt = HydroClawPPT()
    ppt.generate()

