#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HydroClaw认知智能方案PPT生成器
基于现有PPT风格重构
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案（基于分析的现有PPT）
COLORS = {
    'primary_blue': RGBColor(68, 84, 106),      # #44546A 主色调
    'dark_blue': RGBColor(31, 73, 125),         # #1F497D 深蓝
    'accent_blue': RGBColor(91, 155, 213),      # #5B9BD5 强调蓝
    'accent_orange': RGBColor(237, 125, 49),    # #ED7D31 橙色
    'accent_green': RGBColor(112, 173, 71),     # #70AD47 绿色
    'accent_yellow': RGBColor(255, 192, 0),     # #FFC000 黄色
    'text_dark': RGBColor(63, 63, 63),          # #3F3F3F 深灰文字
    'bg_light': RGBColor(231, 230, 230),        # #E7E6E6 浅灰背景
    'white': RGBColor(255, 255, 255)
}

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title, subtitle="", author="", date=""):
        """添加封面页"""
        blank_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(blank_layout)

        # 背景色
        background = slide.background
