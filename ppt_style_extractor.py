"""
PPT风格知识库提取器 - 三引擎协作版
从团队PPT样本中提取：配色方案、字体方案、版式模式、文字表述、图片素材
生成结构化知识库供后续PPT生成使用
"""
import os
import sys
import json
import hashlib
import shutil
from collections import Counter, defaultdict
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import traceback

# ============================================================
# 配置
# ============================================================
BASE_DIR = Path("D:/cowork/ppt")
KB_DIR = BASE_DIR / "ppt_knowledge_base"
KB_COLORS_DIR = KB_DIR / "colors"
KB_FONTS_DIR = KB_DIR / "fonts"
KB_LAYOUTS_DIR = KB_DIR / "layouts"
KB_TEXT_DIR = KB_DIR / "text"
KB_IMAGES_DIR = KB_DIR / "images"
KB_SUMMARY_DIR = KB_DIR / "summary"

# 要分析的样本PPT（按优先级排列）
SAMPLE_PPTS = [
    # 团队高质量PPT
    BASE_DIR / "发布会26.4/output/HydroClaw_发布会_V4.pptx",
    BASE_DIR / "发布会26.4/output/HydroClaw_发布会_V3.pptx",
    BASE_DIR / "发布会26.4/HydroClaw_认知智能方案_v2.pptx",
    BASE_DIR / "HydroClaw_50页完整版_最终.pptx",
    BASE_DIR / "HydroClaw_50页高质量版_V2.pptx",
    # Gamma生成版（视觉参考）
    BASE_DIR / "发布会26.4/output/HydroClaw_Gamma_80p.pptx",
    # 大型学术PPT
    BASE_DIR / "智慧水利团队V5(4).pptx",
    BASE_DIR / "自主运行水网理论、技术与应用260121@雷晓辉.pptx",
    BASE_DIR / "环境与生态领域-重点-河北工程大学-0916.pptx",
    BASE_DIR / "流域水库群智能化运行与精准调控关键技术及应用PPT-1109 修改版(1)(2)(1).pptx",
    BASE_DIR / "京津冀复杂水网运行控制关键技术与测控平台研发-答辩.pptx",
]

# ============================================================
# 初始化目录
# ============================================================
for d in [KB_DIR, KB_COLORS_DIR, KB_FONTS_DIR, KB_LAYOUTS_DIR,
          KB_TEXT_DIR, KB_IMAGES_DIR, KB_SUMMARY_DIR,
          KB_IMAGES_DIR / "backgrounds",
          KB_IMAGES_DIR / "diagrams",
          KB_IMAGES_DIR / "photos",
          KB_IMAGES_DIR / "icons",
          KB_IMAGES_DIR / "charts"]:
    d.mkdir(parents=True, exist_ok=True)


# ============================================================
# 工具函数
# ============================================================
def rgb_to_hex(r, g, b):
    return f"#{r:02X}{g:02X}{b:02X}"


def extract_color_from_rgb(rgb_color):
    """从RGBColor对象提取hex颜色"""
    if rgb_color is None:
        return None
    try:
        return rgb_to_hex(rgb_color[0], rgb_color[1], rgb_color[2])
    except Exception:
        return None


def emu_to_inches(emu):
    if emu is None:
        return None
    return round(emu / 914400, 2)


def classify_image(img_bytes, width, height):
    """根据图片尺寸和特征分类"""
    ratio = width / max(height, 1)
    area = width * height

    if area < 10000:  # 很小的图，可能是图标
        return "icons"
    elif ratio > 3 or ratio < 0.33:  # 很长或很高，可能是装饰条
        return "diagrams"
    elif area > 500000:  # 大图，可能是照片或背景
        return "photos"
    else:
        return "diagrams"


# ============================================================
# 核心提取类
# ============================================================
class PPTStyleExtractor:
    def __init__(self):
        self.all_colors = Counter()       # 所有颜色出现频率
        self.bg_colors = Counter()        # 背景色
        self.text_colors = Counter()      # 文字颜色
        self.shape_colors = Counter()     # 形状填充色

        self.all_fonts = Counter()        # 字体使用频率
        self.title_fonts = Counter()      # 标题字体
        self.body_fonts = Counter()       # 正文字体
        self.font_sizes = Counter()       # 字号分布
        self.title_sizes = Counter()      # 标题字号
        self.body_sizes = Counter()       # 正文字号

        self.layouts = []                 # 版式信息列表
        self.slide_dimensions = Counter() # 幻灯片尺寸

        self.titles = []                  # 所有标题文案
        self.body_texts = []              # 所有正文文案
        self.keywords = Counter()         # 关键词频率
        self.phrases = []                 # 常用短语

        self.images_extracted = 0
        self.images_by_category = Counter()
        self.image_hashes = set()         # 去重

        self.ppt_stats = []              # 每个PPT的统计信息

    def analyze_ppt(self, ppt_path):
        """分析单个PPT文件"""
        name = ppt_path.name
        print(f"\n{'='*60}")
        print(f"分析: {name}")
        print(f"{'='*60}")

        try:
            prs = Presentation(str(ppt_path))
        except Exception as e:
            print(f"  [ERROR] 无法打开: {e}")
            return

        # 记录尺寸
        w_in = emu_to_inches(prs.slide_width)
        h_in = emu_to_inches(prs.slide_height)
        dim_key = f"{w_in}x{h_in}"
        self.slide_dimensions[dim_key] += 1

        slide_count = len(prs.slides)
        ppt_colors = Counter()
        ppt_fonts = Counter()
        ppt_titles = []
        ppt_images = 0

        for slide_idx, slide in enumerate(prs.slides):
            # 提取背景
            self._extract_background(slide)

            # 提取版式信息
            layout_info = self._extract_layout(slide, slide_idx, name)
            self.layouts.append(layout_info)

            for shape in slide.shapes:
                # 提取颜色
                self._extract_shape_colors(shape, ppt_colors)

                # 提取文字
                if hasattr(shape, "text_frame"):
                    self._extract_text(shape, slide_idx, ppt_titles)

                # 提取图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._extract_image(shape, name, slide_idx)
                    ppt_images += 1

                # 处理组合形状
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._process_group(shape, name, slide_idx, ppt_colors, ppt_titles)

        self.all_colors.update(ppt_colors)

        stat = {
            "name": name,
            "slides": slide_count,
            "dimensions": dim_key,
            "unique_colors": len(ppt_colors),
            "top_colors": ppt_colors.most_common(10),
            "images": ppt_images,
            "titles": ppt_titles[:10],
        }
        self.ppt_stats.append(stat)

        print(f"  幻灯片: {slide_count}")
        print(f"  尺寸: {dim_key}")
        print(f"  颜色: {len(ppt_colors)} 种")
        print(f"  图片: {ppt_images} 张")
        print(f"  标题: {len(ppt_titles)} 个")

    def _extract_background(self, slide):
        """提取幻灯片背景色"""
        try:
            bg = slide.background
            if bg.fill and bg.fill.type is not None:
                try:
                    color = extract_color_from_rgb(bg.fill.fore_color.rgb)
                    if color:
                        self.bg_colors[color] += 1
                except Exception:
                    pass
        except Exception:
            pass

    def _extract_shape_colors(self, shape, ppt_colors):
        """提取形状中的颜色"""
        try:
            if hasattr(shape, "fill") and shape.fill and shape.fill.type is not None:
                try:
                    color = extract_color_from_rgb(shape.fill.fore_color.rgb)
                    if color:
                        self.shape_colors[color] += 1
                        ppt_colors[color] += 1
                except Exception:
                    pass
        except Exception:
            pass

        # 线条颜色
        try:
            if hasattr(shape, "line") and shape.line.fill and shape.line.fill.type is not None:
                color = extract_color_from_rgb(shape.line.color.rgb)
                if color:
                    ppt_colors[color] += 1
        except Exception:
            pass

    def _extract_text(self, shape, slide_idx, ppt_titles):
        """提取文字内容和样式"""
        try:
            tf = shape.text_frame
        except Exception:
            return

        full_text = shape.text.strip()
        if not full_text or len(full_text) < 2:
            return

        for para_idx, para in enumerate(tf.paragraphs):
            text = para.text.strip()
            if not text:
                continue

            # 提取字体信息
            for run in para.runs:
                font = run.font
                if font.name:
                    self.all_fonts[font.name] += 1
                if font.size:
                    size_pt = round(font.size / 12700)  # EMU to Pt
                    self.font_sizes[size_pt] += 1
                    if size_pt >= 24:
                        self.title_sizes[size_pt] += 1
                        if font.name:
                            self.title_fonts[font.name] += 1
                    else:
                        self.body_sizes[size_pt] += 1
                        if font.name:
                            self.body_fonts[font.name] += 1

                # 文字颜色
                try:
                    if font.color and font.color.rgb:
                        color = extract_color_from_rgb(font.color.rgb)
                        if color:
                            self.text_colors[color] += 1
                except Exception:
                    pass

            # 分类为标题或正文
            is_title = False
            for run in para.runs:
                if run.font.size and run.font.size / 12700 >= 24:
                    is_title = True
                    break
                if run.font.bold:
                    is_title = True
                    break

            if is_title and len(text) < 50:
                self.titles.append(text)
                ppt_titles.append(text)
            elif len(text) > 5:
                self.body_texts.append(text)

            # 提取关键词（中文分词简化版：按标点分割）
            for seg in text.replace("，", " ").replace("、", " ").replace("；", " ").replace("：", " ").replace("。", " ").replace("—", " ").replace("-", " ").split():
                seg = seg.strip()
                if 2 <= len(seg) <= 10:
                    self.keywords[seg] += 1

    def _extract_layout(self, slide, slide_idx, ppt_name):
        """提取版式信息"""
        shapes_info = []
        for shape in slide.shapes:
            info = {
                "type": str(shape.shape_type),
                "left": emu_to_inches(shape.left),
                "top": emu_to_inches(shape.top),
                "width": emu_to_inches(shape.width),
                "height": emu_to_inches(shape.height),
                "has_text": hasattr(shape, "text") and bool(shape.text.strip()),
                "is_picture": shape.shape_type == MSO_SHAPE_TYPE.PICTURE,
            }
            shapes_info.append(info)

        # 分类版式类型
        n_text = sum(1 for s in shapes_info if s["has_text"])
        n_pic = sum(1 for s in shapes_info if s["is_picture"])
        n_total = len(shapes_info)

        if n_total == 0:
            layout_type = "blank"
        elif n_text <= 1 and n_pic == 0:
            layout_type = "title_only"
        elif n_pic > 0 and n_text > 0:
            layout_type = "mixed_content"
        elif n_text > 3:
            layout_type = "text_heavy"
        elif n_pic > 1:
            layout_type = "image_gallery"
        else:
            layout_type = "standard"

        return {
            "ppt": ppt_name,
            "slide": slide_idx,
            "layout_type": layout_type,
            "shapes_count": n_total,
            "text_count": n_text,
            "picture_count": n_pic,
            "shapes": shapes_info[:20],  # 最多保留20个形状信息
        }

    def _extract_image(self, shape, ppt_name, slide_idx):
        """提取并保存图片"""
        try:
            image = shape.image
            blob = image.blob
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            elif ext not in ("png", "jpg", "gif", "svg+xml", "bmp", "tiff"):
                ext = "png"
            if ext == "svg+xml":
                ext = "svg"

            # 去重
            h = hashlib.md5(blob).hexdigest()
            if h in self.image_hashes:
                return
            self.image_hashes.add(h)

            # 分类
            w = emu_to_inches(shape.width) * 96 if shape.width else 100
            ht = emu_to_inches(shape.height) * 96 if shape.height else 100
            category = classify_image(blob, w, ht)

            # 保存
            safe_name = ppt_name.replace(" ", "_").replace("(", "").replace(")", "")[:30]
            filename = f"{safe_name}_s{slide_idx}_{h[:8]}.{ext}"
            save_path = KB_IMAGES_DIR / category / filename

            with open(save_path, "wb") as f:
                f.write(blob)

            self.images_extracted += 1
            self.images_by_category[category] += 1

        except Exception:
            pass

    def _process_group(self, group, ppt_name, slide_idx, ppt_colors, ppt_titles):
        """处理组合形状"""
        try:
            for shape in group.shapes:
                self._extract_shape_colors(shape, ppt_colors)
                if hasattr(shape, "text_frame"):
                    self._extract_text(shape, slide_idx, ppt_titles)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._extract_image(shape, ppt_name, slide_idx)
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._process_group(shape, ppt_name, slide_idx, ppt_colors, ppt_titles)
        except Exception:
            pass

    # ============================================================
    # 输出知识库
    # ============================================================
    def generate_knowledge_base(self):
        """生成结构化知识库文件"""
        print(f"\n{'='*60}")
        print("生成知识库...")
        print(f"{'='*60}")

        self._generate_color_kb()
        self._generate_font_kb()
        self._generate_layout_kb()
        self._generate_text_kb()
        self._generate_image_kb()
        self._generate_summary()

    def _generate_color_kb(self):
        """生成配色知识库"""
        data = {
            "description": "团队PPT配色方案分析",
            "background_colors": {
                "top_colors": self.bg_colors.most_common(20),
                "description": "幻灯片背景色频率排行"
            },
            "text_colors": {
                "top_colors": self.text_colors.most_common(20),
                "description": "文字颜色频率排行"
            },
            "shape_fill_colors": {
                "top_colors": self.shape_colors.most_common(30),
                "description": "形状填充色频率排行"
            },
            "all_colors_frequency": self.all_colors.most_common(50),
            "recommended_palette": self._recommend_palette(),
        }

        with open(KB_COLORS_DIR / "color_analysis.json", "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        # 生成可读的Markdown报告
        md = "# 团队PPT配色方案\n\n"
        md += "## 推荐配色板\n\n"
        palette = data["recommended_palette"]
        for role, colors in palette.items():
            md += f"### {role}\n"
            for c in colors:
                md += f"- `{c[0]}` (出现{c[1]}次)\n"
            md += "\n"

        md += "## 背景色TOP10\n\n"
        for color, count in self.bg_colors.most_common(10):
            md += f"- `{color}` — {count}次\n"

        md += "\n## 文字色TOP10\n\n"
        for color, count in self.text_colors.most_common(10):
            md += f"- `{color}` — {count}次\n"

        md += "\n## 形状填充色TOP15\n\n"
        for color, count in self.shape_colors.most_common(15):
            md += f"- `{color}` — {count}次\n"

        with open(KB_COLORS_DIR / "配色方案.md", "w", encoding="utf-8") as f:
            f.write(md)

        print(f"  配色: {len(self.all_colors)} 种颜色, 背景色 {len(self.bg_colors)} 种")

    def _recommend_palette(self):
        """基于频率分析推荐配色板"""
        # 过滤纯黑白灰
        def is_neutral(hex_color):
            r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
            diff = max(r, g, b) - min(r, g, b)
            return diff < 30

        accent_colors = [(c, n) for c, n in self.shape_colors.most_common(50) if not is_neutral(c)]
        bg_light = [(c, n) for c, n in self.bg_colors.most_common(20)
                    if int(c[1:3], 16) > 200 and int(c[3:5], 16) > 200 and int(c[5:7], 16) > 200]
        bg_dark = [(c, n) for c, n in self.bg_colors.most_common(20)
                   if int(c[1:3], 16) < 80 and int(c[3:5], 16) < 80 and int(c[5:7], 16) < 80]
        text_main = self.text_colors.most_common(5)

        return {
            "主色调（强调色）": accent_colors[:5],
            "浅色背景": bg_light[:3],
            "深色背景": bg_dark[:3],
            "主要文字色": text_main,
        }

    def _generate_font_kb(self):
        """生成字体知识库"""
        data = {
            "description": "团队PPT字体使用分析",
            "all_fonts": self.all_fonts.most_common(30),
            "title_fonts": self.title_fonts.most_common(10),
            "body_fonts": self.body_fonts.most_common(10),
            "font_sizes": {
                "all": self.font_sizes.most_common(20),
                "title_sizes": self.title_sizes.most_common(10),
                "body_sizes": self.body_sizes.most_common(10),
            },
            "recommended": {
                "title_font": self.title_fonts.most_common(1)[0][0] if self.title_fonts else "Microsoft YaHei",
                "body_font": self.body_fonts.most_common(1)[0][0] if self.body_fonts else "Microsoft YaHei",
                "title_size_range": f"{min(dict(self.title_sizes).keys()) if self.title_sizes else 24}-{max(dict(self.title_sizes).keys()) if self.title_sizes else 44}pt",
                "body_size_range": f"{min(dict(self.body_sizes).keys()) if self.body_sizes else 12}-{max(dict(self.body_sizes).keys()) if self.body_sizes else 20}pt",
            }
        }

        with open(KB_FONTS_DIR / "font_analysis.json", "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        md = "# 团队PPT字体方案\n\n"
        md += "## 推荐字体\n\n"
        md += f"- 标题字体: **{data['recommended']['title_font']}**\n"
        md += f"- 正文字体: **{data['recommended']['body_font']}**\n"
        md += f"- 标题字号范围: {data['recommended']['title_size_range']}\n"
        md += f"- 正文字号范围: {data['recommended']['body_size_range']}\n\n"

        md += "## 字体使用频率TOP15\n\n"
        for font, count in self.all_fonts.most_common(15):
            md += f"- **{font}** — {count}次\n"

        md += "\n## 字号分布\n\n"
        for size, count in sorted(self.font_sizes.most_common(20), key=lambda x: -x[0]):
            md += f"- {size}pt — {count}次\n"

        with open(KB_FONTS_DIR / "字体方案.md", "w", encoding="utf-8") as f:
            f.write(md)

        print(f"  字体: {len(self.all_fonts)} 种字体, {len(self.font_sizes)} 种字号")

    def _generate_layout_kb(self):
        """生成版式知识库"""
        layout_types = Counter(l["layout_type"] for l in self.layouts)
        shapes_per_slide = [l["shapes_count"] for l in self.layouts]
        avg_shapes = sum(shapes_per_slide) / max(len(shapes_per_slide), 1)

        data = {
            "description": "团队PPT版式分析",
            "slide_dimensions": dict(self.slide_dimensions),
            "layout_types": dict(layout_types),
            "avg_shapes_per_slide": round(avg_shapes, 1),
            "total_slides_analyzed": len(self.layouts),
            "sample_layouts": self.layouts[:50],  # 保存前50个作为样本
        }

        with open(KB_LAYOUTS_DIR / "layout_analysis.json", "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        md = "# 团队PPT版式分析\n\n"
        md += f"## 分析概况\n- 总幻灯片: {len(self.layouts)}\n"
        md += f"- 平均元素数/页: {round(avg_shapes, 1)}\n\n"

        md += "## 幻灯片尺寸\n\n"
        for dim, count in self.slide_dimensions.most_common():
            md += f"- **{dim}英寸** — {count}个PPT\n"

        md += "\n## 版式类型分布\n\n"
        for lt, count in layout_types.most_common():
            pct = round(count / len(self.layouts) * 100, 1)
            md += f"- **{lt}** — {count}页 ({pct}%)\n"

        with open(KB_LAYOUTS_DIR / "版式分析.md", "w", encoding="utf-8") as f:
            f.write(md)

        print(f"  版式: {len(self.layouts)} 页, {len(layout_types)} 种类型")

    def _generate_text_kb(self):
        """生成文字知识库"""
        # 标题文案去重
        unique_titles = list(dict.fromkeys(self.titles))

        # 关键词过滤（去掉太常见的）
        stopwords = {"的", "和", "与", "了", "是", "在", "有", "为", "中", "等",
                     "对", "不", "以", "从", "到", "也", "上", "下", "及", "之",
                     "其", "被", "这", "那", "个", "一", "二", "三", "四", "五"}
        filtered_keywords = [(w, c) for w, c in self.keywords.most_common(200)
                             if w not in stopwords and c >= 3]

        # 提取常用短语（标题中的高频模式）
        title_patterns = Counter()
        for t in unique_titles:
            # 提取"XX——XX"模式
            if "——" in t or "—" in t:
                title_patterns["标题——副标题"] += 1
            if "：" in t or ":" in t:
                title_patterns["标题：说明"] += 1
            if any(c.isdigit() for c in t[:3]):
                title_patterns["编号+标题"] += 1

        data = {
            "description": "团队PPT文字表述分析",
            "titles": {
                "total": len(self.titles),
                "unique": len(unique_titles),
                "samples": unique_titles[:100],
                "patterns": dict(title_patterns),
            },
            "keywords": {
                "top_keywords": filtered_keywords[:100],
                "domain_terms": [w for w, c in filtered_keywords if c >= 5][:50],
            },
            "body_text_samples": self.body_texts[:50],
        }

        with open(KB_TEXT_DIR / "text_analysis.json", "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        md = "# 团队PPT文字表述库\n\n"
        md += f"## 统计\n- 标题总数: {len(self.titles)}\n- 去重后: {len(unique_titles)}\n\n"

        md += "## 标题文案样本（前50个）\n\n"
        for i, t in enumerate(unique_titles[:50]):
            md += f"{i+1}. {t}\n"

        md += "\n## 领域关键词TOP50\n\n"
        for w, c in filtered_keywords[:50]:
            md += f"- **{w}** ({c}次)\n"

        md += "\n## 标题模式\n\n"
        for pattern, count in title_patterns.most_common():
            md += f"- {pattern}: {count}次\n"

        with open(KB_TEXT_DIR / "文字表述库.md", "w", encoding="utf-8") as f:
            f.write(md)

        print(f"  文字: {len(unique_titles)} 个标题, {len(filtered_keywords)} 个关键词")

    def _generate_image_kb(self):
        """生成图片素材知识库"""
        md = "# 团队PPT图片素材库\n\n"
        md += f"## 统计\n- 提取图片总数: {self.images_extracted}\n\n"
        md += "## 分类统计\n\n"
        for cat, count in self.images_by_category.most_common():
            md += f"- **{cat}**: {count} 张\n"

        md += "\n## 目录结构\n\n"
        md += "```\n"
        md += "images/\n"
        md += "  backgrounds/  — 背景图片\n"
        md += "  diagrams/     — 图表和示意图\n"
        md += "  photos/       — 照片\n"
        md += "  icons/        — 图标\n"
        md += "  charts/       — 数据图表\n"
        md += "```\n"

        with open(KB_IMAGES_DIR / "图片素材索引.md", "w", encoding="utf-8") as f:
            f.write(md)

        print(f"  图片: {self.images_extracted} 张 ({dict(self.images_by_category)})")

    def _generate_summary(self):
        """生成综合摘要"""
        summary = {
            "description": "HydroClaw团队PPT风格知识库",
            "generated_by": "ppt_style_extractor.py (三引擎协作版)",
            "ppts_analyzed": len(self.ppt_stats),
            "total_slides": sum(s["slides"] for s in self.ppt_stats),
            "knowledge_base_structure": {
                "colors": "配色方案 — 背景色、文字色、强调色",
                "fonts": "字体方案 — 字体名、字号分布",
                "layouts": "版式模式 — 版式类型分布、元素位置",
                "text": "文字表述 — 标题文案、关键词、领域术语",
                "images": "图片素材 — 分类提取的图片资源",
            },
            "ppt_details": self.ppt_stats,
        }

        with open(KB_SUMMARY_DIR / "knowledge_base_summary.json", "w", encoding="utf-8") as f:
            json.dump(summary, f, ensure_ascii=False, indent=2)

        # 生成总览Markdown
        md = "# PPT风格知识库 — 总览\n\n"
        md += "## 分析范围\n\n"
        for stat in self.ppt_stats:
            md += f"- **{stat['name']}** — {stat['slides']}页, {stat['dimensions']}\n"

        md += f"\n## 知识库结构\n\n"
        md += f"| 模块 | 文件 | 说明 |\n"
        md += f"|------|------|------|\n"
        md += f"| 配色方案 | `colors/配色方案.md` | 推荐配色板 + 颜色频率分析 |\n"
        md += f"| 字体方案 | `fonts/字体方案.md` | 推荐字体 + 字号分布 |\n"
        md += f"| 版式分析 | `layouts/版式分析.md` | 版式类型分布 + 样本布局 |\n"
        md += f"| 文字表述 | `text/文字表述库.md` | 标题文案 + 领域关键词 |\n"
        md += f"| 图片素材 | `images/图片素材索引.md` | 分类图片资源 |\n"

        md += f"\n## 使用方式\n\n"
        md += "### 用于 python-pptx 生成\n"
        md += "```python\n"
        md += "import json\n"
        md += "# 加载配色方案\n"
        md += "with open('ppt_knowledge_base/colors/color_analysis.json') as f:\n"
        md += "    colors = json.load(f)\n"
        md += "# 使用推荐配色\n"
        md += "primary = colors['recommended_palette']['主色调（强调色）'][0][0]\n"
        md += "```\n\n"
        md += "### 用于 Presenton\n"
        md += "将知识库内容作为 `instructions` 参数传入生成请求，\n"
        md += "指导AI按团队风格生成内容。\n"

        with open(KB_SUMMARY_DIR / "知识库总览.md", "w", encoding="utf-8") as f:
            f.write(md)

        # 在知识库根目录也生成一份
        with open(KB_DIR / "README.md", "w", encoding="utf-8") as f:
            f.write(md)

        print(f"\n  总览已生成")
        print(f"  分析了 {len(self.ppt_stats)} 个PPT, 共 {sum(s['slides'] for s in self.ppt_stats)} 页")


# ============================================================
# 额外：从已解压的ppt_assets提取媒体
# ============================================================
def extract_from_ppt_assets():
    """从已有的ppt_assets目录复制媒体文件"""
    print("\n从ppt_assets提取媒体素材...")
    assets_dirs = [
        BASE_DIR / "发布会26.4/ppt_assets/team/ppt/media",
        BASE_DIR / "发布会26.4/ppt_assets/theory/ppt/media",
    ]
    copied = 0
    for adir in assets_dirs:
        if not adir.exists():
            continue
        source_name = adir.parent.parent.name  # "team" or "theory"
        for f in adir.iterdir():
            if f.is_file():
                ext = f.suffix.lower()
                if ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg", ".emf"):
                    # 分类
                    size = f.stat().st_size
                    if size > 500000:
                        cat = "photos"
                    elif size < 5000:
                        cat = "icons"
                    else:
                        cat = "diagrams"

                    dest = KB_IMAGES_DIR / cat / f"{source_name}_{f.name}"
                    if not dest.exists():
                        shutil.copy2(f, dest)
                        copied += 1
    print(f"  从ppt_assets复制了 {copied} 个媒体文件")


# ============================================================
# 主函数
# ============================================================
def main():
    print("=" * 60)
    print("PPT风格知识库提取器 v1.0")
    print("三引擎协作版 (Claude + Gemini + Codex)")
    print("=" * 60)

    extractor = PPTStyleExtractor()

    # 分析每个样本PPT
    for ppt_path in SAMPLE_PPTS:
        if ppt_path.exists():
            try:
                extractor.analyze_ppt(ppt_path)
            except Exception as e:
                print(f"  [ERROR] {ppt_path.name}: {e}")
                traceback.print_exc()
        else:
            print(f"  [SKIP] 不存在: {ppt_path.name}")

    # 从ppt_assets提取
    extract_from_ppt_assets()

    # 生成知识库
    extractor.generate_knowledge_base()

    print(f"\n{'='*60}")
    print(f"知识库已生成: {KB_DIR}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
