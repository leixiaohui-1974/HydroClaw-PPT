from __future__ import annotations

import asyncio
import json
import shutil
import sys
import threading
from collections import defaultdict
from collections.abc import Iterable
from pathlib import Path

import torch
from pptx import Presentation
from tqdm.asyncio import tqdm_asyncio

from deeppresenter.main import InputRequest
from deeppresenter.utils.webview import PlaywrightConverter, convert_html_to_pptx
from pptagent.model_utils import get_image_embedding, get_image_model
from pptagent.utils import ppt_to_images

SIMILARITY_THRESHOLD = 0.8
MAX_CONCURRENCY = 4
_IMAGE_MODEL = None
GLOBAL_LOCK = threading.Lock()


def _get_image_model() -> tuple[object, object]:
    global _IMAGE_MODEL
    if _IMAGE_MODEL is None:
        _IMAGE_MODEL = get_image_model()
    return _IMAGE_MODEL


def _image_embeddings(image_dir: Path) -> dict[str, list[float]]:
    with GLOBAL_LOCK:
        extractor, model = _get_image_model()
        return get_image_embedding(str(image_dir), extractor, model, batchsize=8)


def _compare_images(
    left: Iterable[Path],
    right: Iterable[Path],
    left_embeddings: dict[str, list[float]],
    right_embeddings: dict[str, list[float]],
) -> str | None:
    for left_img, right_img in zip(left, right):
        left_embedding = torch.tensor(left_embeddings[left_img.name])
        right_embedding = torch.tensor(right_embeddings[right_img.name])
        sim = torch.cosine_similarity(left_embedding, right_embedding, dim=0).item()
        if sim < SIMILARITY_THRESHOLD:
            return f"Image similarity missmatch at {left_img.name} vs {right_img.name}: {sim}"
    return None


def _check_pptx_validity(pptx_path: Path) -> tuple[bool, str | None]:
    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        return False, f"pptx open failed: {exc}"

    for idx, slide in enumerate(prs.slides):
        # Track (text, position) -> shape_id to detect duplicate elements across shapes
        # Position is rounded to 0.1 inch to catch overlapping shapes
        duplicate_texts = set()
        seen = set()
        for shape_id, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text.strip()
                    if len(text) < 5:
                        continue
                    if not text:
                        continue
                    if text in seen:
                        duplicate_texts.add(text)
                    else:
                        seen.add(text)
                    if len(duplicate_texts) >= 3:
                        return (
                            False,
                            f"Duplicate text: {text} at slide {idx} shape {shape_id}",
                        )

    return True, None


async def _process_slides(
    workspace: Path, semaphore: asyncio.Semaphore
) -> tuple[Path, str] | None:
    intermediate_path = workspace / "intermediate_output.json"
    if not intermediate_path.exists():
        return None

    intermediate = json.loads(intermediate_path.read_text(encoding="utf-8"))
    slide_html_dir = Path(intermediate.get("slide_html_dir", ""))
    if not slide_html_dir:
        return None
    if not slide_html_dir.is_absolute():
        slide_html_dir = workspace / slide_html_dir
    if not slide_html_dir.exists():
        return workspace, "slide_html_dir missing"

    html_files = sorted(slide_html_dir.glob("*.html"))
    if not html_files:
        return None

    final_path = Path(intermediate.get("final", ""))
    pptx_path = None
    if "pptx" in intermediate:
        pptx_path = Path(intermediate["pptx"])
    elif final_path and final_path.suffix.lower() == ".pptx":
        pptx_path = final_path
    elif final_path and final_path.suffix.lower() == ".pdf":
        pptx_path = final_path.with_suffix(".pptx")
    if not pptx_path.is_absolute():
        pptx_path = workspace / pptx_path
    pptx_images_dir = workspace / ".pptx-images"
    pdf_path = (
        final_path
        if final_path and final_path.suffix.lower() == ".pdf"
        else pptx_path.with_suffix(".pdf")
    )
    if not pdf_path.is_absolute():
        pdf_path = workspace / pdf_path
    pdf_images_dir = pdf_path.parent / f".slide_images-pdf-{pdf_path.stem}"
    req = InputRequest.model_validate_json(
        (workspace / ".input_request.json").read_text()
    )
    aspect_ratio = req.powerpoint_type

    async with semaphore:
        try:
            if pptx_path.exists():
                shutil.rmtree(str(pptx_path), ignore_errors=True)

            try:
                await convert_html_to_pptx(
                    slide_html_dir, pptx_path, aspect_ratio=aspect_ratio
                )
            except:  # noqa: E722
                return None

            if pptx_images_dir.exists():
                shutil.rmtree(str(pptx_images_dir), ignore_errors=True)
            await ppt_to_images(str(pptx_path), str(pptx_images_dir))

            if pdf_path.exists():
                shutil.rmtree(str(pdf_path), ignore_errors=True)
            if pdf_images_dir.exists():
                shutil.rmtree(str(pdf_images_dir), ignore_errors=True)
            async with PlaywrightConverter() as pc:
                await pc.convert_to_pdf(
                    [str(p) for p in html_files],
                    pdf_path,
                    aspect_ratio=aspect_ratio,
                )
        except Exception as exc:
            return workspace, f"conversion failed: {exc}"

        is_valid, error_reason = _check_pptx_validity(pptx_path)
        if not is_valid:
            return workspace, error_reason

        pptx_images = sorted(list(pptx_images_dir.glob("*.jpg")))
        pdf_images = sorted(list(pdf_images_dir.glob("*.jpg")))
        if len(pptx_images) != len(pdf_images):
            return workspace, (
                f"image count mismatch: pptx={len(pptx_images)} pdf={len(pdf_images)}"
            )

        pptx_embeddings = _image_embeddings(pptx_images_dir)
        pdf_embeddings = _image_embeddings(pdf_images_dir)
        error = _compare_images(
            pptx_images, pdf_images, pptx_embeddings, pdf_embeddings
        )
        if error:
            return workspace, error

    return None


async def main():
    max_folders = int(sys.argv[1]) if len(sys.argv) > 1 else None
    workspaces = sorted([p for p in Path("/opt/workspace").iterdir() if p.is_dir()])
    if Path("test_html2pptx").exists():
        shutil.rmtree("test_html2pptx")
    if max_folders is not None:
        workspaces = workspaces[:max_folders]

    semaphore = asyncio.Semaphore(MAX_CONCURRENCY)
    tasks = [_process_slides(workspace, semaphore) for workspace in workspaces]
    counter = defaultdict(int)
    for task in tqdm_asyncio.as_completed(tasks):
        result = await task
        if result is None:
            continue
        workspace, error = result
        counter[error] += 1
        (workspace / ".html2pptx.error.txt").write_text(error)
        shutil.copytree(workspace, Path("test_html2pptx") / workspace.name)
    print(f"Detected {sum(counter.values())} errors out of {len(workspaces)} folders")


if __name__ == "__main__":
    import asyncio

    asyncio.run(main())
