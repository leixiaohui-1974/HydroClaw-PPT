"""Test document generation with python-pptx."""

import pytest

PPTX_SCRIPT = """\
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Python-PPTX 综合测试"
slide.placeholders[1].text = "自动生成演示文稿\\n沙箱环境测试"

# Content slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "主要功能特性"
tf = slide.placeholders[1].text_frame
tf.text = "功能 1: 基础文本支持"
for text, level in [("功能 2: 多级列表", 1), ("功能 3: 中文字体渲染", 1), ("功能 4: 自定义样式", 0)]:
    p = tf.add_paragraph()
    p.text, p.level = text, level

# Shapes slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_box.text_frame.text = "图形与形状演示"
title_box.text_frame.paragraphs[0].font.size = Pt(32)

rect = slide.shapes.add_shape(1, Inches(1), Inches(2), Inches(3), Inches(2))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0, 128, 255)
rect.text_frame.text = "蓝色矩形"

circle = slide.shapes.add_shape(9, Inches(5), Inches(2), Inches(2), Inches(2))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(255, 128, 0)
circle.text_frame.text = "圆形"

prs.save('综合演示.pptx')
print(f'SUCCESS: {len(prs.slides)} slides')
"""


@pytest.mark.asyncio
async def test_python_pptx_generation(agent_env, workspace, tool_call_helper):
    """Test python-pptx with multiple slides, shapes, and Chinese content."""
    write_call = tool_call_helper(
        "write_file", {"path": str(workspace / "test_pptx.py"), "content": PPTX_SCRIPT}
    )
    await agent_env.tool_execute(write_call)

    exec_call = tool_call_helper(
        "execute_command", {"command": f"python {workspace / 'test_pptx.py'}"}
    )
    result = await agent_env.tool_execute(exec_call)
    assert not result.is_error and "SUCCESS" in result.text

    list_call = tool_call_helper("list_directory", {"path": str(workspace)})
    list_result = await agent_env.tool_execute(list_call)
    assert "综合演示.pptx" in list_result.text
